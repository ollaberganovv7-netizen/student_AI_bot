from __future__ import annotations
import io
import os
import copy
import logging
from typing import List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE

logger = logging.getLogger(__name__)

def _duplicate_slide(prs, source_slide):
    """Deep copy a slide from the same presentation, preserving all formatting."""
    slide_layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)
    
    # Copy background
    if source_slide.background.fill.type is not None:
        try:
            bg_elem = source_slide.background._element
            new_bg = copy.deepcopy(bg_elem)
            new_slide.background._element.getparent().replace(
                new_slide.background._element, new_bg
            )
        except:
            pass

    # Remove default placeholders from new slide
    for ph in list(new_slide.placeholders):
        try:
            sp = ph._element
            sp.getparent().remove(sp)
        except:
            pass

    # Copy all shapes from template
    for shape in source_slide.shapes:
        try:
            el = copy.deepcopy(shape._element)
            new_slide.shapes._spTree.append(el)
        except Exception as e:
            logger.warning(f"Could not copy shape {shape.name}: {e}")

    return new_slide

def _fill_text_frame(tf, text: str):
    if not tf.paragraphs:
        return
    
    p0 = tf.paragraphs[0]
    alignment = p0.alignment
    
    # Capture font from first run if available
    font_props = {}
    if p0.runs:
        run = p0.runs[0]
        try:
            if run.font.size is not None: font_props['size'] = run.font.size
            if run.font.bold is not None: font_props['bold'] = run.font.bold
            if run.font.italic is not None: font_props['italic'] = run.font.italic
            if run.font.name is not None: font_props['name'] = run.font.name
            if run.font.color and run.font.color.rgb is not None:
                font_props['color_rgb'] = run.font.color.rgb
        except:
            pass

    if isinstance(text, list):
        lines = text
    elif "\n" in text:
        lines = [l.strip() for l in text.split("\n") if l.strip()]
    else:
        lines = [text]
    
    tf.clear()
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = str(line)
        if alignment is not None:
            para.alignment = alignment
        for run in para.runs:
            try:
                if 'size' in font_props: run.font.size = font_props['size']
                if 'bold' in font_props: run.font.bold = font_props['bold']
                if 'italic' in font_props: run.font.italic = font_props['italic']
                if 'name' in font_props: run.font.name = font_props['name']
                if 'color_rgb' in font_props: run.font.color.rgb = font_props['color_rgb']
            except:
                pass
    try:
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except:
        pass

def _replace_all_text_shapes(slide, content: dict):
    main_text = content.get("main_text", "")
    title_text = content.get("title_text", "")
    
    if not main_text and not title_text:
        return
        
    text_shapes = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        w = shape.width / 914400
        h = shape.height / 914400
        area = w * h
        if area < 0.5:
            continue
        text_shapes.append((area, shape))
        
    text_shapes.sort(key=lambda x: x[0], reverse=True)
    
    if not text_shapes:
        return
        
    # The largest shape is body (main_text)
    if main_text:
        _fill_text_frame(text_shapes[0][1].text_frame, main_text)
        
    # The second largest is title (title_text)
    if title_text and len(text_shapes) > 1:
        _fill_text_frame(text_shapes[1][1].text_frame, title_text)
    elif title_text and len(text_shapes) == 1:
        # If there is only one shape, maybe just append it
        pass

def generate_test_pptx(content: dict, template_path: str) -> bytes:
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template not found: {template_path}")

    prs = Presentation(template_path)
    slides = list(prs.slides)
    total_existing = len(slides)
    
    # Fill Title slide (slide 0)
    if total_existing > 0:
        _replace_all_text_shapes(slides[0], {
            "title_text": content.get("topic", "Mavzu"),
            "main_text": f"Bajardi: {content.get('author', 'Foydalanuvchi')}"
        })
        
    ordered_content = []
    # Index 1: Reja
    plan = content.get("plan", [])
    plan_text = "\n".join(f"{i+1}. {p}" for i, p in enumerate(plan))
    ordered_content.append({"title_text": "REJA", "main_text": plan_text})
    
    # Index 2: Kirish
    ordered_content.append({"title_text": "KIRISH", "main_text": content.get("kirish", "")})
    
    # Build cycle
    for sec in content.get("sections", []):
        sec_title = sec.get("title", "")
        # Slide 4: Subpoints
        subpoints_str = "\n".join(f"• {sp}" for sp in sec.get("subpoints", []))
        ordered_content.append({"title_text": sec_title, "main_text": subpoints_str})
        
        for sc in sec.get("slides_content", []):
            if isinstance(sc, str):
                ordered_content.append({
                    "title_text": sec_title, 
                    "main_text": sc
                })
            else:
                ordered_content.append({
                    "title_text": sc.get("title", sec_title), 
                    "main_text": sc.get("body", sc.get("content", ""))
                })
            
    # Now we populate the slides
    # Fill existing slides 1 and 2
    if total_existing > 1 and len(ordered_content) > 0:
        _replace_all_text_shapes(slides[1], ordered_content[0])
    if total_existing > 2 and len(ordered_content) > 1:
        _replace_all_text_shapes(slides[2], ordered_content[1])
        
    # For index 3 onwards, we need to duplicate template slides
    # We will use slide 3 for subpoints layout, and slide 4 for content layout (if exists, else 3)
    subpoints_ref_idx = 3 if total_existing > 3 else total_existing - 1
    content_ref_idx = 4 if total_existing > 4 else subpoints_ref_idx
    
    if subpoints_ref_idx < 0:
        subpoints_ref_idx = 0
    if content_ref_idx < 0:
        content_ref_idx = 0

    # For everything from index 2 in ordered_content
    for i in range(2, len(ordered_content)):
        is_subpoints = ((i - 2) % 3 == 0) # 0, 3, 6 ... are subpoints slides
        
        if is_subpoints:
            ref_slide = slides[subpoints_ref_idx]
        else:
            ref_slide = slides[content_ref_idx]
            
        new_slide = _duplicate_slide(prs, ref_slide)
        _replace_all_text_shapes(new_slide, ordered_content[i])
        
    # Finally, remove the original reference slides (from index 3 to end of original slides)
    # We remove from back to front to avoid index shifting issues
    for i in range(total_existing - 1, 2, -1):
        if i < len(prs.slides):
            sldId = prs.slides._sldIdLst[i]
            try: prs.part.drop_rel(sldId.rId)
            except: pass
            prs.slides._sldIdLst.remove(sldId)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
