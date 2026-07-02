"""
Premium Presentation Design System
====================================
Generates modern, professional, premium designs for presentations.
Each slide gets unique design based on topic and slide type.
No templates copied - every presentation is uniquely designed by AI logic.

Author: Student AI Bot
"""

from __future__ import annotations
import random
import re
from typing import Optional, List, Dict
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.oxml import parse_xml


# ════════════════════════════════════════════════════════════════════════════════
# HELPER UTILITIES
# ════════════════════════════════════════════════════════════════════════════════

def _hex(h: str) -> RGBColor:
    """Convert hex string to RGBColor."""
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _lighter(hex_color: str, factor: float = 0.15) -> str:
    """Make a hex color lighter by factor (0-1)."""
    h = hex_color.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    r = min(255, int(r + (255 - r) * factor))
    g = min(255, int(g + (255 - g) * factor))
    b = min(255, int(b + (255 - b) * factor))
    return f"{r:02X}{g:02X}{b:02X}"


def _darker(hex_color: str, factor: float = 0.2) -> str:
    """Make a hex color darker by factor (0-1)."""
    h = hex_color.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    r = max(0, int(r * (1 - factor)))
    g = max(0, int(g * (1 - factor)))
    b = max(0, int(b * (1 - factor)))
    return f"{r:02X}{g:02X}{b:02X}"


# ════════════════════════════════════════════════════════════════════════════════
# COLOR PALETTES (15+ professional palettes with topic keyword matching)
# ════════════════════════════════════════════════════════════════════════════════

_PALETTES = [
    {
        "name": "apple_dark_mode",
        "bg1": "000000", "bg2": "1C1C1E",
        "accent": "0A84FF", "accent2": "5E5CE6",
        "text1": "FFFFFF", "text2": "EBEBF5",
        "on_accent": "FFFFFF", "is_dark": True,
        "keywords": ["texnologiya", "kompyuter", "dastur", "raqamli", "it", "sun'iy",
                     "intellekt", "robot", "axborot", "tizim", "elektron", "internet",
                     "dasturlash", "algoritm", "ma'lumot", "baza", "tarmoq", "server",
                     "kiberhujum", "xavfsizlik", "programma"]
    },
    {
        "name": "canva_minimalist",
        "bg1": "FFFFFF", "bg2": "F5F5F7",
        "accent": "7B61FF", "accent2": "00C4CC",
        "text1": "1D1D1F", "text2": "424245",
        "on_accent": "FFFFFF", "is_dark": False,
        "keywords": ["san'at", "madaniyat", "adabiyot", "musiqa", "dizayn",
                     "ijodkorlik", "falsafa", "psixologiya", "she'riyat", "arxitektura"]
    },
    {
        "name": "tedx_cinematic",
        "bg1": "111111", "bg2": "000000",
        "accent": "E62B1E", "accent2": "FF4D4D",
        "text1": "FFFFFF", "text2": "CCCCCC",
        "on_accent": "FFFFFF", "is_dark": True,
        "keywords": ["jamiyat", "motivatsiya", "nutq", "liderlik", "biznes", "tadbirkor",
                     "innovatsiya", "g'oya", "global", "inson", "ijtimoiy", "huquq"]
    },
    {
        "name": "corporate_prestige",
        "bg1": "0A1128", "bg2": "1C2541",
        "accent": "D4AF37", "accent2": "F3E5AB",
        "text1": "FFFFFF", "text2": "E0E0E0",
        "on_accent": "000000", "is_dark": True,
        "keywords": ["moliya", "iqtisod", "bank", "investitsiya", "biznes", "tijorat",
                     "menejment", "boshqaruv", "davlat", "siyosat", "huquq", "strategiya"]
    },
    {
        "name": "apple_light_mode",
        "bg1": "FBFBFD", "bg2": "FFFFFF",
        "accent": "FF2D55", "accent2": "FF3B30",
        "text1": "000000", "text2": "86868B",
        "on_accent": "FFFFFF", "is_dark": False,
        "keywords": ["tibbiyot", "sog'liq", "kasallik", "dori", "shifoxona",
                     "jarrohlik", "anatomiya", "fiziologiya", "biologiya", "kimyo"]
    },
    {
        "name": "eco_minimalist",
        "bg1": "F7F9F6", "bg2": "E8EFE6",
        "accent": "34C759", "accent2": "28A745",
        "text1": "1A2518", "text2": "3D4C3A",
        "on_accent": "FFFFFF", "is_dark": False,
        "keywords": ["tabiat", "ekologiya", "atrof-muhit", "botanika", "o'simlik",
                     "hayvon", "qishloq", "dehqonchilik", "fermer", "yashil", "o'rmon"]
    },
    {
        "name": "arctic_ice",
        "bg1": "F4F7FC", "bg2": "E8EEF8",
        "accent": "1565C0", "accent2": "0D47A1",
        "text1": "1A237E", "text2": "3949AB",
        "on_accent": "FFFFFF", "is_dark": False,
        "keywords": ["huquq", "qonun", "davlat", "boshqaruv", "siyosat",
                     "demokratiya", "konstitutsiya", "parlament"]
    },
    {
        "name": "pearl_light",
        "bg1": "FAFAFA", "bg2": "F0F0F0",
        "accent": "FF5722", "accent2": "E64A19",
        "text1": "212121", "text2": "616161",
        "on_accent": "FFFFFF", "is_dark": False,
        "keywords": ["arxitektura", "qurilish", "loyiha", "bino", "shaharsozlik",
                     "infrastruktura", "transport"]
    },
    {
        "name": "golden_finance",
        "bg1": "0F0E0A", "bg2": "1E1C14",
        "accent": "FFD600", "accent2": "FFC107",
        "text1": "FFFDE7", "text2": "FFF9C4",
        "on_accent": "0F0E0A", "is_dark": True,
        "keywords": ["iqtisodiyot", "moliya", "buxgalteriya", "soliq", "budget",
                     "investitsiya", "valyuta", "kredit", "fond", "bozor"]
    },
    {
        "name": "teal_science",
        "bg1": "0A1A1A", "bg2": "102929",
        "accent": "00BFA5", "accent2": "009688",
        "text1": "E0F2F1", "text2": "80CBC4",
        "on_accent": "FFFFFF", "is_dark": True,
        "keywords": ["kimyo", "fizika", "matematika", "fan", "tadqiqot",
                     "laboratoriya", "tajriba", "formula", "nazariya"]
    },
    {
        "name": "coral_education",
        "bg1": "1A0E14", "bg2": "2E1824",
        "accent": "FF4081", "accent2": "F50057",
        "text1": "FCE4EC", "text2": "F48FB1",
        "on_accent": "FFFFFF", "is_dark": True,
        "keywords": ["ta'lim", "o'qitish", "pedagogika", "maktab", "universitet",
                     "talaba", "o'quvchi", "dars", "metodika"]
    },
    {
        "name": "steel_industrial",
        "bg1": "1A1C20", "bg2": "2C2F36",
        "accent": "78909C", "accent2": "546E7A",
        "text1": "ECEFF1", "text2": "B0BEC5",
        "on_accent": "FFFFFF", "is_dark": True,
        "keywords": ["muhandislik", "texnika", "mashina", "metallurgiya",
                     "avtomatika", "mexanika", "elektronika"]
    },
    {
        "name": "olive_history",
        "bg1": "141410", "bg2": "22221A",
        "accent": "CDDC39", "accent2": "AFB42B",
        "text1": "F9FBE7", "text2": "DCED6B",
        "on_accent": "141410", "is_dark": True,
        "keywords": ["tarix", "davlat", "sivilizatsiya", "arxeologiya", "qadimiy",
                     "mustaqillik", "millat", "xalq"]
    },
    {
        "name": "sky_blue_fresh",
        "bg1": "F5FBFF", "bg2": "E3F2FD",
        "accent": "2196F3", "accent2": "1976D2",
        "text1": "0D1B2A", "text2": "37474F",
        "on_accent": "FFFFFF", "is_dark": False,
        "keywords": ["aviatsiya", "kosmos", "samolyot", "parvoz", "meteorologiya",
                     "atmosfera", "ob-havo"]
    },
    {
        "name": "dark_gradient_universal",
        "bg1": "111827", "bg2": "1F2937",
        "accent": "818CF8", "accent2": "6366F1",
        "text1": "F9FAFB", "text2": "D1D5DB",
        "on_accent": "FFFFFF", "is_dark": True,
        "keywords": []  # universal fallback for dark
    },
    {
        "name": "warm_light_universal",
        "bg1": "FFF8F0", "bg2": "FFF0E0",
        "accent": "F97316", "accent2": "EA580C",
        "text1": "1C1917", "text2": "57534E",
        "on_accent": "FFFFFF", "is_dark": False,
        "keywords": []  # universal fallback for light
    },
]


def _select_palette(topic: str) -> dict:
    """Select the best color palette based on topic keywords."""
    if not topic:
        return random.choice(_PALETTES[:-2])  # exclude fallbacks

    topic_lower = topic.lower()
    best_match = None
    best_score = 0

    for pal in _PALETTES:
        score = 0
        for kw in pal["keywords"]:
            if kw in topic_lower:
                score += len(kw)  # longer matches score higher
        if score > best_score:
            best_score = score
            best_match = pal

    if best_match and best_score > 0:
        return best_match

    # Random from main palettes (exclude last 2 universal ones)
    return random.choice(_PALETTES[:-2])


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE TYPE DETECTION
# ════════════════════════════════════════════════════════════════════════════════

_TITLE_WORDS = {"muqova", "title"}
_SECTION_WORDS = {"o'quv savoli", "bo'lim", "section", "qism"}
_PLAN_WORDS = {"reja", "plan", "mundarija", "o'quv savollari"}
_INTRO_WORDS = {"kirish", "introduction", "maqsad"}
_CONCLUSION_WORDS = {"xulosa", "conclusion", "yakun"}
_REFERENCE_WORDS = {"adabiyotlar", "adabiyot", "manbalar", "references", "foydalanilgan"}
_QUOTE_WORDS = {"iqtibos", "tsitata", "quote"}
_FINAL_WORDS = {"rahmat", "thank", "e'tibor", "yakuniy"}
_GOAL_WORDS = {"maqsad", "maqsadlar", "tarbiyaviy"}


def _detect_slide_type(slide, idx: int, total: int) -> str:
    """Detect slide type from content and position."""
    if idx == 0:
        return "title"
    if idx == total - 1:
        return "final"

    text = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            text += " " + shape.text_frame.text.lower()

    text = text.strip()

    for w in _QUOTE_WORDS:
        if w in text:
            return "quote"
    for w in _CONCLUSION_WORDS:
        if w in text:
            return "conclusion"
    for w in _REFERENCE_WORDS:
        if w in text:
            return "references"
    for w in _GOAL_WORDS:
        if w in text:
            return "goals"
    for w in _PLAN_WORDS:
        if w in text:
            return "plan"
    for w in _INTRO_WORDS:
        if w in text:
            return "intro"
    for w in _SECTION_WORDS:
        if w in text:
            return "section"
    for w in _FINAL_WORDS:
        if w in text:
            return "final"

    return "content"


# ════════════════════════════════════════════════════════════════════════════════
# GRADIENT BACKGROUNDS (XML injection)
# ════════════════════════════════════════════════════════════════════════════════

_NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _set_gradient_bg(slide, c1: str, c2: str, angle: int = 5400000):
    """Apply 2-stop gradient background. Falls back to solid on error."""
    try:
        slide.follow_master_background = False
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _hex(c1)

        bg_elem = slide._element.find(f'{{{_NS_P}}}bg')
        if bg_elem is not None:
            bgPr = bg_elem.find(f'{{{_NS_P}}}bgPr')
            if bgPr is not None:
                for child in list(bgPr):
                    local = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                    if 'Fill' in local or 'fill' in local:
                        bgPr.remove(child)

                grad = parse_xml(
                    f'<a:gradFill xmlns:a="{_NS_A}" rotWithShape="1">'
                    f'<a:gsLst>'
                    f'<a:gs pos="0"><a:srgbClr val="{c1}"/></a:gs>'
                    f'<a:gs pos="100000"><a:srgbClr val="{c2}"/></a:gs>'
                    f'</a:gsLst>'
                    f'<a:lin ang="{angle}" scaled="0"/>'
                    f'</a:gradFill>'
                )
                bgPr.insert(0, grad)
    except Exception:
        try:
            slide.follow_master_background = False
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = _hex(c1)
        except Exception:
            pass


def _set_gradient_3stop(slide, c1: str, c2: str, c3: str, angle: int = 5400000):
    """Apply 3-stop gradient background for extra premium feel."""
    try:
        slide.follow_master_background = False
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _hex(c1)

        bg_elem = slide._element.find(f'{{{_NS_P}}}bg')
        if bg_elem is not None:
            bgPr = bg_elem.find(f'{{{_NS_P}}}bgPr')
            if bgPr is not None:
                for child in list(bgPr):
                    local = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                    if 'Fill' in local or 'fill' in local:
                        bgPr.remove(child)

                grad = parse_xml(
                    f'<a:gradFill xmlns:a="{_NS_A}" rotWithShape="1">'
                    f'<a:gsLst>'
                    f'<a:gs pos="0"><a:srgbClr val="{c1}"/></a:gs>'
                    f'<a:gs pos="50000"><a:srgbClr val="{c2}"/></a:gs>'
                    f'<a:gs pos="100000"><a:srgbClr val="{c3}"/></a:gs>'
                    f'</a:gsLst>'
                    f'<a:lin ang="{angle}" scaled="0"/>'
                    f'</a:gradFill>'
                )
                bgPr.insert(0, grad)
    except Exception:
        _set_gradient_bg(slide, c1, c3, angle)


# ════════════════════════════════════════════════════════════════════════════════
# DECORATIVE SHAPES (geometric elements with transparency and shadows)
# ════════════════════════════════════════════════════════════════════════════════

def _send_to_back(shape):
    """Sends a shape to the back of the shape tree so it doesn't cover text."""
    try:
        element = shape.element
        spTree = element.getparent()
        spTree.remove(element)
        spTree.insert(2, element)
    except Exception:
        pass

def _add_shape_no_border(slide, shape_type, left, top, width, height, fill_hex,
                          alpha_pct: int = 100):
    """Add shape with fill, no border, optional transparency."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex(fill_hex)
    shape.line.fill.background()

    if alpha_pct < 100:
        try:
            from pptx.oxml.xmlchemy import OxmlElement
            alpha_val = int(alpha_pct * 1000)
            solidFill = shape.fill._xPr.find(f'{{{_NS_A}}}solidFill')
            if solidFill is not None:
                srgb = solidFill.find(f'{{{_NS_A}}}srgbClr')
                if srgb is not None:
                    alpha_elem = OxmlElement('a:alpha')
                    alpha_elem.set('val', str(alpha_val))
                    srgb.append(alpha_elem)
        except Exception:
            pass
            
    _send_to_back(shape)
    return shape


def _add_shadow(shape, blur_pt: int = 5, dist_pt: int = 3, angle: int = 45):
    """Add soft outer shadow to a shape."""
    try:
        shadow = shape.shadow
        shadow.inherit = False
        shadow.visible = True
        shadow.blur_radius = Pt(blur_pt)
        shadow.distance = Pt(dist_pt)
        shadow.angle = angle
    except Exception:
        pass


def _decorate_title_slide(slide, pal: dict, sw, sh):
    """Ultra-Premium Title Slide: Glassmorphism, Cinematic Lights, Minimalist UI."""
    accent = pal["accent"]
    accent2 = pal["accent2"]

    # 1. Base dark/light cinematic gradient background orb (Very large)
    _add_shape_no_border(slide, MSO_SHAPE.OVAL,
                         -Inches(4), -Inches(4),
                         sw + Inches(8), sh + Inches(8), accent, alpha_pct=12)
                         
    # 2. Floating Light Effects (3D abstract glowing orbs)
    _add_shape_no_border(slide, MSO_SHAPE.OVAL,
                         sw - Inches(5), -Inches(2),
                         Inches(8), Inches(8), accent2, alpha_pct=25)
    _add_shape_no_border(slide, MSO_SHAPE.OVAL,
                         -Inches(2), sh - Inches(4),
                         Inches(7), Inches(7), accent, alpha_pct=20)
                         
    # 3. Glassmorphism overlay (simulated by a semi-transparent full screen box in front of the orbs)
    glass = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, sh)
    glass.fill.solid()
    glass.fill.fore_color.rgb = _hex(pal["bg1"])
    glass.line.fill.background()
    try:
        from pptx.oxml.xmlchemy import OxmlElement
        alpha_elem = OxmlElement('a:alpha')
        alpha_elem.set('val', '75000') # 75% opacity for frosted glass effect
        srgb = glass.fill._xPr.find(f'{{{_NS_A}}}solidFill').find(f'{{{_NS_A}}}srgbClr')
        srgb.append(alpha_elem)
    except Exception:
        pass
    _send_to_back(glass)
    
    # 4. We must send the glowing orbs BEHIND the glass overlay!
    # By sending glass to back, and then sending orbs to back AGAIN, the orbs are behind the glass.
    # Actually, the shapes were added before glass. So glass is on top. We want: Background -> Orbs -> Glass -> Text.
    # We will just leave them as is (glass is added AFTER orbs, so it's in front of orbs naturally).
    # We must send ALL of them behind the text boxes.
    for shape in list(slide.shapes):
        if not shape.has_text_frame:
            _send_to_back(shape)

    # 5. Add a sleek vertical accent bar on the left edge over the glass
    _add_shape_no_border(slide, MSO_SHAPE.RECTANGLE,
                         Inches(0.5), Inches(2), Inches(0.06), sh - Inches(4), accent, alpha_pct=100)
    
    # 6. Apply Soft Shadow to Title Text
    for shape in slide.shapes:
        if shape.has_text_frame and _is_title_shape(shape, slide):
            try:
                from pptx.oxml.xmlchemy import OxmlElement
                spPr = shape.element.spPr
                effectLst = OxmlElement('a:effectLst')
                outerShdw = OxmlElement('a:outerShdw')
                outerShdw.set('blurRad', '100000')  # Soft blur
                outerShdw.set('dist', '50000')
                outerShdw.set('dir', '2700000')    # 45 degrees
                outerShdw.set('algn', 'tl')
                srgbClr = OxmlElement('a:srgbClr')
                srgbClr.set('val', '000000')
                alpha = OxmlElement('a:alpha')
                alpha.set('val', '30000') # 30% transparency
                srgbClr.append(alpha)
                outerShdw.append(srgbClr)
                effectLst.append(outerShdw)
                spPr.append(effectLst)
            except Exception:
                pass


def _decorate_section_slide(slide, pal: dict, sw, sh):
    """Cinematic Episode Intro for Section Slides."""
    accent = pal["accent"]
    accent2 = pal["accent2"]
    
    pic = _find_picture(slide)
    if pic:
        # 1. Full Screen Picture
        pic.left = 0
        pic.top = 0
        pic.width = sw
        pic.height = sh
        _send_to_back(pic)
        
        # 2. Cinematic Gradient Overlay
        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, sh)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = _hex("050505") # pure cinematic dark
        overlay.line.fill.background()
        try:
            from pptx.oxml.xmlchemy import OxmlElement
            alpha = OxmlElement('a:alpha')
            alpha.set('val', '80000') # 80% opacity
            overlay.fill._xPr.find(f'{{{_NS_A}}}solidFill').find(f'{{{_NS_A}}}srgbClr').append(alpha)
        except:
            pass
            
    # 3. Dynamic Geometric Elements (Cinematic Letterbox feel)
    _add_shape_no_border(slide, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), sw - Inches(1.0), Inches(0.02), accent, alpha_pct=100)
    _add_shape_no_border(slide, MSO_SHAPE.RECTANGLE, Inches(0.5), sh - Inches(0.5), sw - Inches(1.0), Inches(0.02), accent, alpha_pct=100)
    
    # Center floating diamond
    _add_shape_no_border(slide, MSO_SHAPE.DIAMOND, sw/2 - Inches(0.15), Inches(0.8), Inches(0.3), Inches(0.3), accent2, alpha_pct=100)
    
    title, body = _find_text_boxes(slide)
    
    # 4. Cinematic Typography Layout (Center Aligned)
    if title:
        title.left = Inches(1)
        title.width = sw - Inches(2)
        title.top = sh/2 - Inches(1.5)
        
    if body:
        body.left = Inches(1)
        body.width = sw - Inches(2)
        body.top = sh/2 + Inches(0.2)


def _find_picture(slide):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return shape
    return None

def _find_text_boxes(slide):
    title = None
    body = None
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if _is_title_shape(shape, slide):
            title = shape
        else:
            if not body or shape.width > body.width:
                body = shape
    return title, body

def _detect_and_draw_infographics(slide, card, cx, cy, cw, ch, text, pal):
    """Draw a mini visual chart if the text contains a percentage or number."""
    accent = pal["accent"]
    accent2 = pal["accent2"]
    
    match = re.search(r'(\d{1,3})%', text)
    if match:
        pct = int(match.group(1))
        if pct > 100: pct = 100
        
        bar_w = cw - Inches(1.0)
        bar_h = Inches(0.12)
        bar_x = cx + Inches(0.8)
        bar_y = cy + ch - Inches(0.4)
        
        track = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bar_x, bar_y, bar_w, bar_h)
        track.fill.solid()
        track.fill.fore_color.rgb = _hex(pal["bg1"])
        track.line.fill.background()
        created = [track.shape_id]
        
        fill_w = bar_w * (pct / 100.0)
        if fill_w > 0:
            fill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bar_x, bar_y, fill_w, bar_h)
            fill.fill.solid()
            fill.fill.fore_color.rgb = _hex(accent2)
            fill.line.fill.background()
            _apply_glow_effect(fill, accent2)
            created.append(fill.shape_id)
            
        return True, created
    return False, []

def _convert_to_cards(slide, body, text_left, text_width, start_y, available_h, pal, highlight_last=False):
    """Convert body text paragraphs into premium interactive-looking cards."""
    points = []
    for p in body.text_frame.paragraphs:
        if p.text.strip():
            points.append(p.text.strip())
            
    body.left = slide.part.part.package.presentation_part.presentation.slide_width + Inches(10) # Hide original
    
    num_pts = len(points)
    if num_pts == 0: return
    
    gap = Inches(0.3)
    card_h = min(Inches(1.5), (available_h - (num_pts - 1) * gap) / num_pts)
    accent = pal["accent"]
    
    anim_groups = []
    
    for i, pt in enumerate(points):
        group_ids = []
        cy = start_y + i * (card_h + gap)
        is_highlight = highlight_last and (i == num_pts - 1)
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, text_left, cy, text_width, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = _hex(accent if is_highlight else pal["bg2"])
        
        if not is_highlight:
            try:
                from pptx.oxml.xmlchemy import OxmlElement
                alpha = OxmlElement('a:alpha')
                alpha.set('val', '85000') 
                card.fill._xPr.find(f'{{{_NS_A}}}solidFill').find(f'{{{_NS_A}}}srgbClr').append(alpha)
            except: pass
            
        card.line.color.rgb = _hex(accent)
        card.line.width = Pt(1.5 if not is_highlight else 0)
        _apply_card_shadow(card)
        group_ids.append(card.shape_id)
        
        has_chart = False
        if not is_highlight:
            has_chart, chart_ids = _detect_and_draw_infographics(slide, card, text_left, cy, text_width, card_h, pt, pal)
            group_ids.extend(chart_ids)
        
        icon_size = Inches(0.35)
        shapes = [MSO_SHAPE.STAR_5_POINT, MSO_SHAPE.DIAMOND, MSO_SHAPE.HEXAGON]
        icon = slide.shapes.add_shape(shapes[i % len(shapes)], text_left + Inches(0.3), cy + Inches(0.2), icon_size, icon_size)
        icon.fill.solid()
        icon.fill.fore_color.rgb = _hex(pal["bg1"] if is_highlight else accent)
        icon.line.fill.background()
        group_ids.append(icon.shape_id)
        
        tb = slide.shapes.add_textbox(text_left + Inches(0.8), cy + Inches(0.05), text_width - Inches(1.0), card_h - (Inches(0.4) if has_chart else Inches(0.2)))
        group_ids.append(tb.shape_id)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP if has_chart else MSO_ANCHOR.MIDDLE
        p = tf.add_paragraph()
        
        display_text = pt
        if len(display_text) > 180:
            display_text = display_text[:177] + "..."
            
        p.text = display_text
        p.font.name = "Arial"
        p.font.size = Pt(16) if text_width > Inches(6) else Pt(14)
        p.font.color.rgb = _hex("FFFFFF" if is_highlight else pal["text1"])
        
        _send_to_back(card)
        anim_groups.append(group_ids)
        
    _apply_object_animations(slide, [{"ids": g, "preset": "fade", "on_click": True} for g in anim_groups])

def _decorate_split_layout(slide, pal: dict, sw, sh, is_left_image=True):
    """Premium Split-Screen Layout (50/50). Image on one side, cards on the other."""
    pic = _find_picture(slide)
    title, body = _find_text_boxes(slide)
    
    if pic:
        half_w = sw / 2
        pic.top = 0
        pic.height = sh
        pic.width = half_w
        pic.left = 0 if is_left_image else half_w
        _send_to_back(pic)

    if title and body:
        half_w = sw / 2
        margin = Inches(0.5)
        text_left = half_w + margin if is_left_image else margin
        text_width = half_w - (margin * 2)
        
        title.left = text_left
        title.width = text_width
        title.top = Inches(1.0)
        
        start_y = title.top + title.height + Inches(0.5)
        available_h = sh - start_y - Inches(0.5)
        
        _convert_to_cards(slide, body, text_left, text_width, start_y, available_h, pal)

def _decorate_cinematic_layout(slide, pal: dict, sw, sh):
    """Premium Cinematic Layout. Image covers whole slide with a dark overlay."""
    pic = _find_picture(slide)
    title, body = _find_text_boxes(slide)
    
    if pic:
        pic.left = 0
        pic.top = 0
        pic.width = sw
        pic.height = sh
        _send_to_back(pic)
        
        # Dark overlay
        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, sh)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = _hex("000000")
        overlay.line.fill.background()
        try:
            from pptx.oxml.xmlchemy import OxmlElement
            alpha_elem = OxmlElement('a:alpha')
            alpha_elem.set('val', '60000') # 60% opacity
            srgb = overlay.fill._xPr.find(f'{{{_NS_A}}}solidFill').find(f'{{{_NS_A}}}srgbClr')
            srgb.append(alpha_elem)
        except Exception:
            pass
            
        _send_to_back(overlay)
        _send_to_back(pic) # pic goes behind overlay
        
    if title and body:
        margin = Inches(1)
        title.left = margin
        title.width = sw - (margin * 2)
        title.top = Inches(2)
        
        body.left = margin
        body.width = sw - (margin * 2)
        body.top = title.top + title.height + Inches(0.5)

def _decorate_minimal_typographic(slide, pal: dict, sw, sh):
    """Premium Minimalist text layout (Cards) when no image is present."""
    accent = pal["accent"]
    _add_shape_no_border(slide, MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.5), sh, accent, alpha_pct=100)
    
    title, body = _find_text_boxes(slide)
    if title and body:
        title.left = Inches(1.5)
        title.width = sw - Inches(2.5)
        title.top = Inches(1.0)
        
        start_y = title.top + title.height + Inches(0.5)
        available_h = sh - start_y - Inches(0.5)
        
        _convert_to_cards(slide, body, Inches(1.5), sw - Inches(2.5), start_y, available_h, pal)


def _decorate_quote_slide(slide, pal: dict, sw, sh):
    """Premium State Style Quote Slide (President Portrait with Gold/Blue accents)."""
    # Force state colors (Gold and Deep Navy) for the quote slide background
    accent_gold = "D4AF37"
    navy_bg = "0F172A"
    
    # 1. State/Modern Background (Deep Blue)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, sh)
    bg.fill.solid()
    bg.fill.fore_color.rgb = _hex(navy_bg)
    bg.line.fill.background()
    _send_to_back(bg)
    
    # Transparent geometric accents
    _add_shape_no_border(slide, MSO_SHAPE.OVAL, -Inches(2), -Inches(2), Inches(6), Inches(6), accent_gold, alpha_pct=10)
    _add_shape_no_border(slide, MSO_SHAPE.OVAL, sw - Inches(4), sh - Inches(4), Inches(8), Inches(8), "FFFFFF", alpha_pct=5)
    
    pic = _find_picture(slide)
    title, body = _find_text_boxes(slide)
    
    if pic:
        # 2. Premium Portrait Framing (Left side)
        # Attempt to make it a standard portrait ratio
        pic.height = sh - Inches(2)
        pic.width = int(pic.height * 0.75) 
        pic.left = Inches(1)
        pic.top = Inches(1)
        
        # Add Gold Frame behind portrait
        frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, pic.left - Inches(0.1), pic.top - Inches(0.1), pic.width + Inches(0.2), pic.height + Inches(0.2))
        frame.fill.background() # transparent fill
        frame.line.color.rgb = _hex(accent_gold)
        frame.line.width = Pt(4)
        
        # Ensure proper z-ordering: Background -> Frame -> Portrait
        _send_to_back(frame)
        _send_to_back(bg) 
    
    # 3. Katta qo'shtirnoq (Huge Quotation Mark)
    if body:
        quote_mark = slide.shapes.add_textbox(sw / 2, Inches(0.5), Inches(2), Inches(2))
        tf = quote_mark.text_frame
        p = tf.paragraphs[0]
        p.text = '“'
        p.font.name = "Georgia"
        p.font.size = Pt(120)
        p.font.color.rgb = _hex(accent_gold)
        try:
            # Make it semi-transparent
            from pptx.oxml.xmlchemy import OxmlElement
            alpha = OxmlElement('a:alpha')
            alpha.set('val', '40000') # 40% opacity
            p.runs[0].font.color._xFill.find(f'{{{_NS_A}}}srgbClr').append(alpha)
        except:
            pass

    if title and body:
        title.left = sw / 2 + Inches(0.5)
        title.width = sw / 2 - Inches(1)
        title.top = Inches(1.5)
        
        body.left = sw / 2 + Inches(0.5)
        body.width = sw / 2 - Inches(1)
        body.top = title.top + title.height + Inches(0.5)
        
        tf.word_wrap = True
        
        anim_configs = []
        if pic:
            anim_configs.append({"ids": [pic.shape_id, frame.shape_id], "preset": "zoom", "on_click": False})
        
        quote_ids = [quote_mark.shape_id, body.shape_id]
        if title: quote_ids.append(title.shape_id)
        anim_configs.append({"ids": quote_ids, "preset": "fade", "on_click": True})
        
        _apply_object_animations(slide, anim_configs)

def _decorate_plan_slide(slide, pal: dict, sw, sh):
    """Premium Timeline Step Navigation for Plan/O'quv savollari slide."""
    accent = pal["accent"]
    accent2 = pal["accent2"]
    
    # Premium Gradient Background using large soft shapes
    _add_shape_no_border(slide, MSO_SHAPE.OVAL, sw - Inches(4), -Inches(4), Inches(8), Inches(8), accent, alpha_pct=15)
    _add_shape_no_border(slide, MSO_SHAPE.OVAL, -Inches(3), sh - Inches(3), Inches(6), Inches(6), accent2, alpha_pct=10)
    
    title, body = _find_text_boxes(slide)
    if not body: return
    
    questions = []
    for p in body.text_frame.paragraphs:
        if p.text.strip():
            questions.append(p.text.strip())
            
    body.left = sw + Inches(10) # Hide original
    
    num_qs = len(questions)
    if num_qs == 0: return
    
    start_y = Inches(2.2)
    available_h = sh - Inches(2.5)
    gap = Inches(0.4)
    block_h = min(Inches(1.2), (available_h - (num_qs - 1) * gap) / num_qs)
    
    line_x = Inches(1.5)
    
    # Draw colorful progress line connecting the steps
    timeline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, line_x, start_y + block_h/2, Inches(0.08), max(Inches(0.1), (num_qs - 1) * (block_h + gap)))
    timeline.fill.solid()
    timeline.fill.fore_color.rgb = _hex(accent2)
    timeline.line.fill.background()
    _send_to_back(timeline)
    
    for i, q in enumerate(questions):
        cy = start_y + i * (block_h + gap)
        is_active = (i == 0) # The first item acts as the "Active" glowing point on the timeline
        
        # Step Marker (Circle on the line)
        marker_size = Inches(0.5)
        marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, line_x - marker_size/2 + Inches(0.04), cy + (block_h - marker_size)/2, marker_size, marker_size)
        marker.fill.solid()
        marker.fill.fore_color.rgb = _hex(accent if is_active else pal["bg2"])
        marker.line.color.rgb = _hex(accent) 
        marker.line.width = Pt(2)
        if is_active:
            _apply_glow_effect(marker, accent)
        
        # Premium Block for Question
        block_w = sw - Inches(3)
        block = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, line_x + Inches(0.6), cy, block_w, block_h)
        block.fill.solid()
        block.fill.fore_color.rgb = _hex(pal["bg1"])
        
        # Glassmorphism effect
        alpha_val = '90000' if is_active else '75000'
        try:
            from pptx.oxml.xmlchemy import OxmlElement
            alpha = OxmlElement('a:alpha')
            alpha.set('val', alpha_val)
            block.fill._xPr.find(f'{{{_NS_A}}}solidFill').find(f'{{{_NS_A}}}srgbClr').append(alpha)
        except:
            pass
            
        block.line.color.rgb = _hex(accent2 if is_active else accent)
        block.line.width = Pt(2 if is_active else 1)
        
        if is_active:
            _apply_glow_effect(block, accent2) # Active item glows!
        else:
            _apply_card_shadow(block)
        
        # AI Icon (Simulated geometric icons)
        icon_size = block_h * 0.4
        shapes = [MSO_SHAPE.STAR_5_POINT, MSO_SHAPE.DIAMOND, MSO_SHAPE.HEXAGON, MSO_SHAPE.OVAL, MSO_SHAPE.PENTAGON]
        shape_type = shapes[i % len(shapes)]
        icon = slide.shapes.add_shape(shape_type, line_x + Inches(0.9), cy + (block_h - icon_size)/2, icon_size, icon_size)
        icon.fill.solid()
        icon.fill.fore_color.rgb = _hex(accent2 if is_active else accent)
        icon.line.fill.background()
        
        # Text
        tb = slide.shapes.add_textbox(line_x + Inches(1.5), cy, block_w - Inches(1.6), block_h)
        tf = tb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.add_paragraph()
        p.text = q
        p.font.name = "Arial"
        p.font.size = Pt(22)
        p.font.color.rgb = _hex(pal["text1"])
        if is_active:
            p.font.bold = True
            
        _send_to_back(block)
        
    # Order background properly
    for s in list(slide.shapes):
        if not s.has_text_frame and s != timeline:
            _send_to_back(s)
    _send_to_back(timeline)

def _apply_glow_effect(shape, color_hex):
    try:
        from pptx.oxml.xmlchemy import OxmlElement
        spPr = shape.element.spPr
        effectLst = OxmlElement('a:effectLst')
        glow = OxmlElement('a:glow')
        glow.set('rad', '150000')  # 15pt blur radius for glow
        srgbClr = OxmlElement('a:srgbClr')
        srgbClr.set('val', color_hex)
        alpha = OxmlElement('a:alpha')
        alpha.set('val', '60000') # 60% opacity glow
        srgbClr.append(alpha)
        glow.append(srgbClr)
        effectLst.append(glow)
        spPr.append(effectLst)
    except Exception:
        pass

def _apply_card_shadow(shape):
    try:
        from pptx.oxml.xmlchemy import OxmlElement
        spPr = shape.element.spPr
        effectLst = OxmlElement('a:effectLst')
        outerShdw = OxmlElement('a:outerShdw')
        outerShdw.set('blurRad', '200000')  # 20pt blur
        outerShdw.set('dist', '50000')      # 5pt distance
        outerShdw.set('dir', '5400000')     # 90 degrees (straight down)
        srgbClr = OxmlElement('a:srgbClr')
        srgbClr.set('val', '000000')
        alpha = OxmlElement('a:alpha')
        alpha.set('val', '15000') # 15% opacity
        srgbClr.append(alpha)
        outerShdw.append(srgbClr)
        effectLst.append(outerShdw)
        spPr.append(effectLst)
    except Exception:
        pass

def _decorate_goals_slide(slide, pal: dict, sw, sh):
    """Premium Card-based decorations for Goals/Maqsadlar slide."""
    accent = pal["accent"]
    accent2 = pal["accent2"]
    
    # Large soft glowing background circle in the center right
    _add_shape_no_border(slide, MSO_SHAPE.OVAL, sw/2, -Inches(2), Inches(8), Inches(8), accent, alpha_pct=10)
    _add_shape_no_border(slide, MSO_SHAPE.OVAL, -Inches(2), sh - Inches(4), Inches(6), Inches(6), accent2, alpha_pct=8)
    
    title, body = _find_text_boxes(slide)
    if not body: return
    
    # Extract goals from the body text frame
    goals = []
    for p in body.text_frame.paragraphs:
        if p.text.strip():
            goals.append(p.text.strip())
            
    # Hide the original body
    body.left = sw + Inches(10)
    
    num_cards = len(goals)
    if num_cards == 0: return
    
    # Vertical Card Layout Engine
    card_width = sw - Inches(2)
    start_x = Inches(1)
    
    # Dynamic height calculation
    available_height = sh - Inches(2.5)
    gap = Inches(0.4)
    # Ensure cards aren't too tall
    card_height = min(Inches(1.5), (available_height - (num_cards - 1) * gap) / num_cards)
    start_y = Inches(2.5)
    
    for i, goal in enumerate(goals):
        cy = start_y + i * (card_height + gap)
        
        # 1. Premium Frosted Card Background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, start_x, cy, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = _hex(pal["bg1"])
        
        # Simulate glassmorphism transparency if dark mode, otherwise solid white-ish
        try:
            from pptx.oxml.xmlchemy import OxmlElement
            alpha = OxmlElement('a:alpha')
            alpha.set('val', '85000') # 85% opacity
            card.fill._xPr.find(f'{{{_NS_A}}}solidFill').find(f'{{{_NS_A}}}srgbClr').append(alpha)
        except Exception:
            pass
            
        card.line.color.rgb = _hex(accent)
        card.line.width = Pt(1)
        
        # 2. Hover / Float Shadow Effect
        _apply_card_shadow(card)
        
        # 3. Dynamic Icon (Using different shapes based on index)
        icon_size = card_height * 0.4
        shapes = [MSO_SHAPE.STAR_5_POINT, MSO_SHAPE.DIAMOND, MSO_SHAPE.HEXAGON, MSO_SHAPE.OVAL]
        shape_type = shapes[i % len(shapes)]
        
        icon = slide.shapes.add_shape(shape_type, start_x + Inches(0.5), cy + (card_height - icon_size)/2, icon_size, icon_size)
        icon.fill.solid()
        icon.fill.fore_color.rgb = _hex(accent2)
        icon.line.fill.background()
        
        # 4. Goal Text
        tb = slide.shapes.add_textbox(start_x + Inches(1.5), cy, card_width - Inches(2), card_height)
        tf = tb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.add_paragraph()
        p.text = goal
        p.font.name = "Arial"
        p.font.size = Pt(24)
        p.font.color.rgb = _hex(pal["text1"])
        p.font.bold = True
        
        # Ensures text z-ordering is above the card
        _send_to_back(card)
        # Send background elements completely to back
        for s in slide.shapes:
            if not s.has_text_frame and s != card and s != icon:
                _send_to_back(s)


def _decorate_conclusion_slide(slide, pal: dict, sw, sh):
    """Premium minimalist decorations for conclusion slides with Cards & Infographics."""
    accent = pal["accent"]
    accent2 = pal["accent2"]
    
    # Elegant top and bottom bars
    _add_shape_no_border(slide, MSO_SHAPE.RECTANGLE,
                         Inches(1), Inches(0.5),
                         sw - Inches(2), Inches(0.05), accent, alpha_pct=100)
    _add_shape_no_border(slide, MSO_SHAPE.RECTANGLE,
                         Inches(1), sh - Inches(0.5),
                         sw - Inches(2), Inches(0.05), accent2, alpha_pct=100)
                         
    title, body = _find_text_boxes(slide)
    
    if title and body:
        title.left = Inches(1.5)
        title.width = sw - Inches(3)
        title.top = Inches(1.0)
        
        start_y = title.top + title.height + Inches(0.5)
        available_h = sh - start_y - Inches(0.8)
        
        _convert_to_cards(slide, body, Inches(1.5), sw - Inches(3), start_y, available_h, pal, highlight_last=True)


def _decorate_references_slide(slide, pal: dict, sw, sh):
    """Premium minimalist decorations for References slides with Cards list."""
    accent = pal["accent"]
    accent2 = pal["accent2"]
    
    # Clean accent bar on the right edge
    _add_shape_no_border(slide, MSO_SHAPE.RECTANGLE,
                         sw - Inches(0.2), 0,
                         Inches(0.2), sh, accent2, alpha_pct=100)
                         
    # Single subtle glow bottom left
    _add_shape_no_border(slide, MSO_SHAPE.OVAL,
                         -Inches(3), sh - Inches(3),
                         Inches(6), Inches(6), accent, alpha_pct=5)
                         
    title, body = _find_text_boxes(slide)
    if title and body:
        title.left = Inches(1.5)
        title.width = sw - Inches(3)
        title.top = Inches(1.0)
        
        start_y = title.top + title.height + Inches(0.5)
        available_h = sh - start_y - Inches(0.8)
        
        _convert_to_cards(slide, body, Inches(1.5), sw - Inches(3), start_y, available_h, pal)


def _decorate_final_slide(slide, pal: dict, sw, sh):
    """Premium 'Thank You' slide decorations (Cinematic Ending)."""
    accent = pal["accent"]
    accent2 = pal["accent2"]
    
    pic = _find_picture(slide)
    if pic:
        pic.left = 0
        pic.top = 0
        pic.width = sw
        pic.height = sh
        _send_to_back(pic)
        
        # Cinematic Gradient Overlay
        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, sh)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = _hex("000000")
        overlay.line.fill.background()
        try:
            from pptx.oxml.xmlchemy import OxmlElement
            alpha = OxmlElement('a:alpha')
            alpha.set('val', '75000') # 75% opacity dark overlay
            overlay.fill._xPr.find(f'{{{_NS_A}}}solidFill').find(f'{{{_NS_A}}}srgbClr').append(alpha)
        except: pass
    
    # Transparent Geometric Shapes / Glows
    _add_shape_no_border(slide, MSO_SHAPE.OVAL, (sw - Inches(8)) // 2, (sh - Inches(8)) // 2, Inches(8), Inches(8), accent, alpha_pct=15)
    _add_shape_no_border(slide, MSO_SHAPE.OVAL, (sw - Inches(5)) // 2, (sh - Inches(5)) // 2, Inches(5), Inches(5), accent2, alpha_pct=10)
    
    # Premium glowing lines (Light rays)
    _add_shape_no_border(slide, MSO_SHAPE.RECTANGLE, sw/2 - Inches(2), Inches(1), Inches(4), Inches(0.05), accent, alpha_pct=100)
    _add_shape_no_border(slide, MSO_SHAPE.RECTANGLE, sw/2 - Inches(2), sh - Inches(1), Inches(4), Inches(0.05), accent2, alpha_pct=100)

    title, body = _find_text_boxes(slide)
    if title:
        title.left = Inches(1)
        title.width = sw - Inches(2)
        title.top = sh / 2 - Inches(1.5)
        
    if body:
        # Glassmorphism Card for Subtitle
        card_w = sw - Inches(4)
        card_h = Inches(1.5)
        card_x = Inches(2)
        card_y = sh / 2 + Inches(0.5)
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_x, card_y, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = _hex(pal["bg1"])
        try:
            from pptx.oxml.xmlchemy import OxmlElement
            alpha = OxmlElement('a:alpha')
            alpha.set('val', '80000') 
            card.fill._xPr.find(f'{{{_NS_A}}}solidFill').find(f'{{{_NS_A}}}srgbClr').append(alpha)
        except: pass
        card.line.color.rgb = _hex(accent)
        card.line.width = Pt(1)
        _apply_glow_effect(card, accent)
        
        body.left = card_x + Inches(0.5)
        body.width = card_w - Inches(1)
        body.top = card_y + (card_h - Inches(0.5)) / 2
        
        _send_to_back(body)
        _send_to_back(card)
        if title: _send_to_back(title)
        
    for s in list(slide.shapes):
        if not s.has_text_frame and s not in (title, body):
            _send_to_back(s)
    if pic:
        _send_to_back(pic)




# ════════════════════════════════════════════════════════════════════════════════
# TEXT FORMATTING (Arial font, keyword highlighting, proper sizing)
# ════════════════════════════════════════════════════════════════════════════════

def _is_title_shape(shape, slide) -> bool:
    """Check if a shape is likely the title (placeholder or by position/size)."""
    if shape.is_placeholder:
        from pptx.enum.shapes import PP_PLACEHOLDER
        pf = shape.placeholder_format
        if pf.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
            return True
    font_size = 0
    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size and run.font.size.pt > font_size:
                    font_size = run.font.size.pt
    except Exception:
        pass
    return font_size >= 22 or (shape.top < Inches(2) and shape.width > Inches(5))


def _apply_object_animations(slide, anim_configs):
    """Injects a single XML timing tree for all object animations on the slide.
    anim_configs is a list of dicts: {"ids": [spid1, spid2], "preset": "fade"/"zoom", "on_click": True/False}
    """
    if not anim_configs: return
    try:
        from pptx.oxml import parse_xml
        import xml.etree.ElementTree as ET
        
        timing = ET.Element(f'{{{_NS_P}}}timing')
        tnLst = ET.SubElement(timing, f'{{{_NS_P}}}tnLst')
        par1 = ET.SubElement(tnLst, f'{{{_NS_P}}}par')
        cTn1 = ET.SubElement(par1, f'{{{_NS_P}}}cTn', {'id': '1', 'dur': 'indefinite', 'restart': 'never', 'nodeType': 'tmRoot'})
        childTnLst1 = ET.SubElement(cTn1, f'{{{_NS_P}}}childTnLst')
        seq = ET.SubElement(childTnLst1, f'{{{_NS_P}}}seq', {'concurrent': '1', 'nextAc': 'seek'})
        cTn2 = ET.SubElement(seq, f'{{{_NS_P}}}cTn', {'id': '2', 'dur': 'indefinite', 'nodeType': 'mainSeq'})
        main_childTnLst = ET.SubElement(cTn2, f'{{{_NS_P}}}childTnLst')

        current_id = 3
        for config in anim_configs:
            group = config.get("ids", [])
            preset = config.get("preset", "fade")
            on_click = config.get("on_click", False)
            
            preset_id = '10' # fade
            if preset == "zoom": preset_id = '64' # Zoom is sometimes 64 or 51 or 12. Actually let's use presetClass="entr". For Zoom in Office it's 55 or 64. Let's use 10 for safety if we don't know, but 12 is Zoom in some specs. Actually, presetID="10" is Fade, presetID="1" is Appear, presetID="2" is Fly In. Let's stick to "10" for all to avoid corruption, as PowerPoint can be very strict about preset classes. Wait, we'll use 10.
            
            # Main paragraph for this sequence step
            par_main = ET.SubElement(main_childTnLst, f'{{{_NS_P}}}par')
            cTn_main = ET.SubElement(par_main, f'{{{_NS_P}}}cTn', {'id': str(current_id), 'fill': 'hold'})
            stCondLst_main = ET.SubElement(cTn_main, f'{{{_NS_P}}}stCondLst')
            ET.SubElement(stCondLst_main, f'{{{_NS_P}}}cond', {'delay': 'indefinite' if on_click else '0'})
            childTnLst_main = ET.SubElement(cTn_main, f'{{{_NS_P}}}childTnLst')
            current_id += 1
            
            for spid in group:
                par_shape = ET.SubElement(childTnLst_main, f'{{{_NS_P}}}par')
                cTn_shape = ET.SubElement(par_shape, f'{{{_NS_P}}}cTn', {'id': str(current_id), 'fill': 'hold'})
                stCondLst_shape = ET.SubElement(cTn_shape, f'{{{_NS_P}}}stCondLst')
                ET.SubElement(stCondLst_shape, f'{{{_NS_P}}}cond', {'delay': '0'}) # With previous within group
                childTnLst_shape = ET.SubElement(cTn_shape, f'{{{_NS_P}}}childTnLst')
                
                par_effect = ET.SubElement(childTnLst_shape, f'{{{_NS_P}}}par')
                cTn_effect = ET.SubElement(par_effect, f'{{{_NS_P}}}cTn', {
                    'id': str(current_id+1), 'presetID': '10', 'presetClass': 'entr', 'presetSubtype': '0', 'fill': 'hold', 'nodeType': 'withEffect' if not on_click else 'clickEffect'
                })
                stCondLst_effect = ET.SubElement(cTn_effect, f'{{{_NS_P}}}stCondLst')
                ET.SubElement(stCondLst_effect, f'{{{_NS_P}}}cond', {'delay': '0'})
                childTnLst_effect = ET.SubElement(cTn_effect, f'{{{_NS_P}}}childTnLst')
                
                set_node = ET.SubElement(childTnLst_effect, f'{{{_NS_P}}}set')
                cBhvr = ET.SubElement(set_node, f'{{{_NS_P}}}cBhvr')
                cTn_dur = ET.SubElement(cBhvr, f'{{{_NS_P}}}cTn', {'id': str(current_id+2), 'dur': '800' if preset == 'zoom' else '500', 'fill': 'hold'})
                stCondLst_dur = ET.SubElement(cTn_dur, f'{{{_NS_P}}}stCondLst')
                ET.SubElement(stCondLst_dur, f'{{{_NS_P}}}cond', {'delay': '0'})
                tgtEl = ET.SubElement(cBhvr, f'{{{_NS_P}}}tgtEl')
                ET.SubElement(tgtEl, f'{{{_NS_P}}}spTgt', {'spid': str(spid)})
                attrNameLst = ET.SubElement(cBhvr, f'{{{_NS_P}}}attrNameLst')
                attrName = ET.SubElement(attrNameLst, f'{{{_NS_P}}}attrName')
                attrName.text = 'style.visibility'
                
                to_node = ET.SubElement(set_node, f'{{{_NS_P}}}to')
                ET.SubElement(to_node, f'{{{_NS_P}}}strVal', {'val': 'visible'})
                
                current_id += 3
            
        xml_str = ET.tostring(timing, encoding='unicode')
        slide._element.append(parse_xml(xml_str))
    except Exception as e:
        pass


def _format_all_text(slide, pal: dict, slide_type: str):
    """Reformat all text on the slide: Arial font, proper colors, sizes."""
    text1 = _hex(pal["text1"])
    text2 = _hex(pal["text2"])
    accent_rgb = _hex(pal["accent"])
    is_dark = pal["is_dark"]

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        is_title = _is_title_shape(shape, slide)

        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                # Set Arial font
                run.font.name = "Arial"

                # Determine appropriate color
                if is_title:
                    if slide_type in ("title", "section", "final", "goals", "plan"):
                        run.font.color.rgb = accent_rgb
                    else:
                        run.font.color.rgb = text1
                    run.font.bold = True
                    
                    # Center align title on most premium layouts
                    if slide_type in ("goals", "plan", "section", "conclusion", "references", "content", "intro"):
                        para.alignment = PP_ALIGN.CENTER
                        
                    # Huge uppercase for section and conclusion titles
                    if slide_type in ("section", "conclusion", "references"):
                        run.text = run.text.upper()
                        try:
                            run.font.size = Pt(48)
                        except Exception:
                            pass
                    else:
                        # Make other titles quite large too
                        try:
                            run.font.size = Pt(40)
                        except Exception:
                            pass
                else:
                    # Force all body text to use the palette's text1 color
                    # to ensure readability against our custom background gradients
                    try:
                        run.font.color.rgb = text1
                    except Exception:
                        pass
                            
                    # Remove bullets and make text italic and larger for quote slides
                    if slide_type == "quote":
                        run.font.italic = True
                        run.font.color.rgb = _hex("FFFFFF") # Force white text on the deep navy background
                        try:
                            # Increase font size for quote text if possible
                            if run.font.size:
                                run.font.size = Pt(int(run.font.size.pt * 1.2))
                            else:
                                run.font.size = Pt(28)
                            # Remove bullet point format if exists
                            if para.level > 0:
                                para.level = 0
                        except Exception:
                            pass
                            
                    # Force white text and center alignment for Cinematic Section intro slides and Final slides
                    if slide_type in ("section", "final"):
                        run.font.color.rgb = _hex("FFFFFF")
                        para.alignment = PP_ALIGN.CENTER
                        if run.font.size and slide_type == "section":
                            run.font.size = Pt(int(run.font.size.pt * 1.1))
                            
                    # Section slide body text formatting (center, bold, huge)
                    if slide_type == "section":
                        run.font.bold = True
                        para.alignment = PP_ALIGN.CENTER
                        if para.level > 0:
                            para.level = 0
                        try:
                            if run.font.size:
                                run.font.size = Pt(int(run.font.size.pt * 1.5))
                            else:
                                run.font.size = Pt(32)
                        except Exception:
                            pass
                    
                    # Justify body text on most slides
                    if slide_type in ("goals", "plan", "content", "intro", "conclusion", "references"):
                        para.alignment = PP_ALIGN.JUSTIFY
                        
                    # Conclusion slide body text formatting (remove bullets, bold)
                    if slide_type == "conclusion":
                        run.font.bold = True
                        if para.level > 0:
                            para.level = 0
                        try:
                            if run.font.size:
                                run.font.size = Pt(int(run.font.size.pt * 1.2))
                            else:
                                run.font.size = Pt(24)
                        except Exception:
                            pass


# ════════════════════════════════════════════════════════════════════════════════
# PROFESSIONAL TRANSITIONS (Morph, Fade, Push, Zoom)
# ════════════════════════════════════════════════════════════════════════════════

def _add_transition(slide, transition_type: str = "fade", speed: str = "med"):
    """Add slide transition via XML injection."""
    try:
        transitions = {
            "fade": f'<p:transition xmlns:p="{_NS_P}" spd="{speed}"><p:fade/></p:transition>',
            "push": f'<p:transition xmlns:p="{_NS_P}" spd="{speed}"><p:push dir="u"/></p:transition>',
            "cover": f'<p:transition xmlns:p="{_NS_P}" spd="{speed}"><p:cover dir="l"/></p:transition>',
            "wipe": f'<p:transition xmlns:p="{_NS_P}" spd="{speed}"><p:wipe dir="d"/></p:transition>',
            "split": f'<p:transition xmlns:p="{_NS_P}" spd="{speed}"><p:split orient="horz" dir="out"/></p:transition>',
        }
        xml_str = transitions.get(transition_type, transitions["fade"])
        slide._element.insert(1, parse_xml(xml_str))
    except Exception:
        pass


# ════════════════════════════════════════════════════════════════════════════════
# BACKGROUND APPLICATION per slide type
# ════════════════════════════════════════════════════════════════════════════════

_GRADIENT_ANGLES = [5400000, 2700000, 0, 10800000]  # down, diagonal, right, up


def _apply_bg_for_type(slide, pal: dict, slide_type: str, idx: int):
    """Apply the right background based on slide type."""
    bg1, bg2 = pal["bg1"], pal["bg2"]
    accent, accent2 = pal["accent"], pal["accent2"]

    if slide_type == "title":
        _set_gradient_3stop(slide, bg1, _darker(bg2, 0.1), _lighter(bg1, 0.05), angle=5400000)
    elif slide_type == "section":
        _set_gradient_bg(slide, _darker(bg1, 0.1), bg2, angle=2700000)
    elif slide_type == "quote":
        _set_gradient_bg(slide, bg1, _lighter(bg2, 0.05), angle=5400000)
    elif slide_type == "conclusion":
        _set_gradient_bg(slide, bg2, bg1, angle=10800000)
    elif slide_type == "final":
        _set_gradient_3stop(slide, _darker(bg1, 0.15), bg1, _lighter(bg2, 0.08), angle=5400000)
    else:
        angle = _GRADIENT_ANGLES[idx % len(_GRADIENT_ANGLES)]
        if idx % 3 == 0:
            _set_gradient_bg(slide, bg1, bg2, angle)
        elif idx % 3 == 1:
            _set_gradient_bg(slide, bg2, bg1, angle)
        else:
            _set_gradient_bg(slide, bg1, _lighter(bg2, 0.03), angle)


# ════════════════════════════════════════════════════════════════════════════════
# MAIN PUBLIC API
# ════════════════════════════════════════════════════════════════════════════════

def apply_premium_design(prs, topic: str = ""):
    """
    Apply full premium design to an existing presentation.
    Called AFTER slides are created and filled with content.
    
    This function:
    1. Selects a color palette based on topic keywords
    2. Applies gradient backgrounds per slide type
    3. Adds decorative geometric elements (non-overlapping with text)
    4. Reformats all text to Arial with proper contrast
    """
    palette = _select_palette(topic)
    slides = list(prs.slides)
    total = len(slides)
    if total == 0:
        return

    sw = prs.slide_width
    sh = prs.slide_height

    prev_decorator_idx = -1

    for idx, slide in enumerate(slides):
        slide_type = _detect_slide_type(slide, idx, total)

        # 1. Apply gradient background
        _apply_bg_for_type(slide, palette, slide_type, idx)

        # 2. Add decorative elements (varies per slide type)
        if slide_type == "title":
            _decorate_title_slide(slide, palette, sw, sh)
        elif slide_type == "section":
            _decorate_section_slide(slide, palette, sw, sh)
        elif slide_type == "quote":
            _decorate_quote_slide(slide, palette, sw, sh)
        elif slide_type == "goals":
            _decorate_goals_slide(slide, palette, sw, sh)
        elif slide_type == "plan":
            _decorate_plan_slide(slide, palette, sw, sh)
        elif slide_type == "conclusion":
            _decorate_conclusion_slide(slide, palette, sw, sh)
        elif slide_type == "references":
            _decorate_references_slide(slide, palette, sw, sh)
        elif slide_type == "final":
            _add_transition(slide, "zoom", "slow") # Smooth Zoom Out for cinematic ending
            _decorate_final_slide(slide, palette, sw, sh)
        else:
            has_pic = _find_picture(slide) is not None
            if has_pic:
                if idx % 3 == 0:
                    _decorate_cinematic_layout(slide, palette, sw, sh)
                elif idx % 3 == 1:
                    _decorate_split_layout(slide, palette, sw, sh, is_left_image=True)
                else:
                    _decorate_split_layout(slide, palette, sw, sh, is_left_image=False)
            else:
                _decorate_minimal_typographic(slide, palette, sw, sh)

        # 3. Reformat text
        _format_all_text(slide, palette, slide_type)


def apply_premium_transitions(prs):
    """Apply professional slide transitions."""
    slides = list(prs.slides)
    total = len(slides)
    if total == 0:
        return

    for idx, slide in enumerate(slides):
        slide_type = _detect_slide_type(slide, idx, total)

        if slide_type == "title":
            _add_transition(slide, "morph", "slow") # Cinematic Morph/Fade In for the title
        elif slide_type == "section":
            _add_transition(slide, "fade", "slow") # Cinematic fade in for episode intros
        elif slide_type == "quote":
            _add_transition(slide, "zoom", "slow") # Smooth Zoom for portrait
        elif slide_type == "goals":
            _add_transition(slide, "push", "med") # Simulates floating up
        elif slide_type == "plan":
            _add_transition(slide, "wipe", "slow") # Wipe mimics a timeline drawing itself
        elif slide_type == "conclusion":
            _add_transition(slide, "wipe", "med")
        elif slide_type == "references":
            _add_transition(slide, "cover", "med") # Smooth cover from side for references
        elif slide_type == "final":
            _add_transition(slide, "zoom", "slow") # Smooth cinematic zoom out
        else:
            # Randomize transitions for standard content slides
            options = ["morph", "fade", "zoom", "push", "wipe"]
            t = options[idx % len(options)]
            _add_transition(slide, t, "slow" if t in ["morph", "zoom"] else "med")
