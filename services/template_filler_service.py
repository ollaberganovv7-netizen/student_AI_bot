from __future__ import annotations
"""
Template Filler Service
-----------------------
Deep PPTX template analyzer + smart content filler.

Unlike the basic pptx_service, this service:
1. Deeply analyzes EVERY slide in the template (layout type, text areas, sizes, image placeholders)
2. Classifies each slide by type (title, info, cards, list, comparison, gallery, conclusion)
3. Calculates exact text capacity (max chars / words) per text block
4. Fills template IN-PLACE without breaking design — no slide cloning/deletion
5. Auto-trims text to fit physical text areas
6. Inserts images into image placeholders preserving size/position
"""

import io
import os
import re
import logging
from typing import List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)

# ══════════════════════════════════════════════════════════════════════════════
# CONSTANTS
# ══════════════════════════════════════════════════════════════════════════════

TITLE_PH_TYPES = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}
BODY_PH_TYPES = {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT}
SUBTITLE_PH_TYPES = {PP_PLACEHOLDER.SUBTITLE}
PICTURE_PH_TYPES = {PP_PLACEHOLDER.PICTURE, PP_PLACEHOLDER.BITMAP, PP_PLACEHOLDER.MEDIA_CLIP}

# Average characters per square inch for typical PowerPoint text
CHARS_PER_SQ_INCH_TITLE = 6    # large font titles
CHARS_PER_SQ_INCH_BODY = 14    # body text (~16-18pt)
CHARS_PER_SQ_INCH_SMALL = 18   # smaller text boxes (~12-14pt)

# Slide type classification keywords
TITLE_SLIDE_KEYWORDS = {"title", "титульный", "cover", "sarlavha", "muqova"}
CONCLUSION_KEYWORDS = {"xulosa", "conclusion", "заключение", "итог", "summary", "yakuniy"}
PLAN_KEYWORDS = {"reja", "plan", "план", "mundarija", "содержание", "agenda", "outline"}

# ── TEMPLATE TAGS ─────────────────────────────────────────────────────────────
# Users can mark template shapes with tags like {{title}}, {{body}}, {{remove}}
# to give the AI/filler exact instructions on what content goes where.
#
# Supported tags (medium tag system):
#   {{title}}       — main slide title
#   {{subtitle}}    — subtitle (e.g. for cover slide)
#   {{author}}      — author name
#   {{topic}}       — presentation topic (alias of title for cover)
#   {{date}}        — current date
#   {{body}}        — main text / bullet list (single block)
#   {{body_1}}, {{body_2}}, ... — multiple text columns/blocks
#   {{point_1}}, {{point_2}}, ... — individual list items / step titles
#   {{image}} or {{image_1}}, {{image_2}} — picture placeholder
#   {{remove}}      — DELETE this shape (watermarks, attribution, decoration)

TAG_PATTERN = re.compile(r"\{\{\s*([a-zA-Z_][a-zA-Z0-9_]*)\s*\}\}")


def _iter_all_shapes(slide_or_group):
    """Iterate all shapes including sub-shapes inside GroupShape."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for shape in slide_or_group.shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_all_shapes(shape)


def parse_tags(text: str) -> List[str]:
    """Extract all {{tag}} names from text. Returns lowercase list."""
    if not text:
        return []
    return [m.group(1).lower() for m in TAG_PATTERN.finditer(text)]


def get_primary_tag(text: str) -> Optional[str]:
    """Get the first/primary tag from text, lowercase, or None."""
    tags = parse_tags(text)
    return tags[0] if tags else None


# ══════════════════════════════════════════════════════════════════════════════
# DATA CLASSES
# ══════════════════════════════════════════════════════════════════════════════

class TextBlockInfo:
    """Info about a single text area in a slide."""
    def __init__(self, shape, role: str, max_chars: int, max_words: int,
                 width_inches: float, height_inches: float, font_size_pt: float = 0):
        self.shape = shape
        self.shape_name = shape.name  # store name for cross-Presentation lookup
        self.role = role          # "title", "subtitle", "body", "textbox"
        self.max_chars = max_chars
        self.max_words = max_words
        self.width_inches = width_inches
        self.height_inches = height_inches
        self.font_size_pt = font_size_pt
        self.original_text = ""
        self.tag: Optional[str] = None        # primary {{tag}} found in original text (e.g. "title", "body_1")
        self.all_tags: List[str] = []         # all tags found
        try:
            self.original_text = shape.text_frame.text if shape.has_text_frame else ""
            self.all_tags = parse_tags(self.original_text)
            self.tag = self.all_tags[0] if self.all_tags else None
        except:
            pass


class ImagePlaceholderInfo:
    """Info about an image placeholder in a slide."""
    def __init__(self, shape, width_inches: float, height_inches: float, tag: Optional[str] = None):
        self.shape = shape
        self.shape_name = shape.name  # store name for cross-Presentation lookup
        self.width_inches = width_inches
        self.height_inches = height_inches
        self.tag = tag  # e.g. "image_1" if shape was tagged


class SlideAnalysis:
    """Full analysis of a single slide."""
    def __init__(self, index: int):
        self.index = index
        self.slide_type = "info"         # title, plan, info, cards, list, comparison, gallery, conclusion
        self.text_blocks: List[TextBlockInfo] = []
        self.image_placeholders: List[ImagePlaceholderInfo] = []
        self.shapes_to_remove: List[str] = []   # shape names tagged with {{remove}}
        self.has_title = False
        self.has_body = False
        self.has_subtitle = False
        self.has_images = False
        self.body_block_count = 0
        self.total_max_chars = 0
        self.total_max_words = 0
        self.has_tags = False              # True if any shape has {{tag}} markup

    def to_dict(self) -> dict:
        return {
            "index": self.index,
            "slide_type": self.slide_type,
            "has_title": self.has_title,
            "has_body": self.has_body,
            "has_subtitle": self.has_subtitle,
            "has_images": self.has_images,
            "has_tags": self.has_tags,
            "body_block_count": self.body_block_count,
            "total_max_chars": self.total_max_chars,
            "total_max_words": self.total_max_words,
            "image_count": len(self.image_placeholders),
            "text_blocks": [
                {
                    "role": b.role,
                    "tag": b.tag,
                    "all_tags": b.all_tags,
                    "max_chars": b.max_chars,
                    "max_words": b.max_words,
                    "width": round(b.width_inches, 1),
                    "height": round(b.height_inches, 1),
                    "font_pt": round(b.font_size_pt, 1),
                    "sample_text": b.original_text[:60],
                }
                for b in self.text_blocks
            ],
            "image_tags": [p.tag for p in self.image_placeholders if p.tag],
        }


class TemplateAnalysis:
    """Full analysis of the entire PPTX template."""
    def __init__(self, path: str):
        self.path = path
        self.total_slides = 0
        self.slides: List[SlideAnalysis] = []

    def to_prompt_context(self) -> str:
        """Generate a concise text description for the AI prompt."""
        lines = [f"SHABLON TAHLILI: {self.total_slides} ta slayd"]
        for sa in self.slides:
            blocks_desc = []
            for b in sa.text_blocks:
                blocks_desc.append(f"{b.role}(max {b.max_words} so'z)")
            img_desc = f", {len(sa.image_placeholders)} ta rasm joyi" if sa.image_placeholders else ""
            line = (
                f"  Slayd {sa.index}: tur={sa.slide_type}, "
                f"bloklar=[{', '.join(blocks_desc)}]{img_desc}"
            )
            lines.append(line)
        return "\n".join(lines)

    def to_dict(self) -> dict:
        return {
            "total_slides": self.total_slides,
            "slides": [s.to_dict() for s in self.slides],
        }


# ══════════════════════════════════════════════════════════════════════════════
# TEMPLATE ANALYSIS
# ══════════════════════════════════════════════════════════════════════════════

def _resolve_path(template_path: str) -> str:
    """Resolve template path to absolute."""
    if template_path and not os.path.isabs(template_path):
        base = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        full = os.path.join(base, template_path)
        if os.path.exists(full):
            return full
    return template_path


def _get_max_font_size(shape) -> float:
    """Get the largest font size in a shape (in pt)."""
    max_size = 0.0
    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    pt_val = run.font.size.pt
                    if pt_val > max_size:
                        max_size = pt_val
    except:
        pass
    return max_size


def _estimate_capacity(width_inches: float, height_inches: float,
                       font_size_pt: float, role: str) -> tuple:
    """Estimate max chars and words that fit in a text area."""
    area = width_inches * height_inches

    if role == "title":
        chars_per_inch = CHARS_PER_SQ_INCH_TITLE
    elif font_size_pt and font_size_pt >= 20:
        chars_per_inch = CHARS_PER_SQ_INCH_TITLE
    elif font_size_pt and font_size_pt <= 14:
        chars_per_inch = CHARS_PER_SQ_INCH_SMALL
    else:
        chars_per_inch = CHARS_PER_SQ_INCH_BODY

    # Adjust for actual font size if known
    if font_size_pt and font_size_pt > 0:
        # Larger font → fewer chars, smaller font → more chars
        size_factor = 16.0 / max(font_size_pt, 8)
        chars_per_inch = int(chars_per_inch * size_factor)

    max_chars = max(10, int(area * chars_per_inch))
    max_words = max(3, max_chars // 6)  # average 6 chars per word

    return max_chars, max_words


def _classify_slide_type(analysis: SlideAnalysis, slide_index: int, total_slides: int) -> str:
    """Determine slide type based on its structure and position."""
    # First slide is almost always title
    if slide_index == 0:
        return "title"

    # Last slide is often conclusion
    if slide_index == total_slides - 1:
        return "conclusion"

    # Check original text for keywords
    all_text = " ".join(b.original_text.lower() for b in analysis.text_blocks)

    if any(kw in all_text for kw in PLAN_KEYWORDS):
        return "plan"
    if any(kw in all_text for kw in CONCLUSION_KEYWORDS):
        return "conclusion"
    if any(kw in all_text for kw in TITLE_SLIDE_KEYWORDS) and slide_index == 0:
        return "title"

    # Gallery: has images but minimal text
    if len(analysis.image_placeholders) >= 2 and analysis.body_block_count <= 1:
        return "gallery"

    # Cards: multiple body blocks of similar size (columns)
    if analysis.body_block_count >= 2:
        return "cards"

    # Comparison: 2 body blocks
    if analysis.body_block_count == 2:
        return "comparison"

    # Image + text
    if analysis.has_images and analysis.has_body:
        return "image_info"

    # List/Info: standard
    return "info"


def analyze_template_deep(template_path: str) -> TemplateAnalysis:
    """
    Deeply analyze a PPTX template: every slide, every shape, every text area.
    Returns TemplateAnalysis with full capacity data.
    """
    template_path = _resolve_path(template_path)
    result = TemplateAnalysis(template_path)

    if not template_path or not os.path.exists(template_path):
        logger.warning(f"Template not found: {template_path}")
        return result

    prs = Presentation(template_path)
    result.total_slides = len(prs.slides)

    for slide_idx, slide in enumerate(prs.slides):
        sa = SlideAnalysis(slide_idx)

        # Analyze all shapes (including sub-shapes in groups)
        for shape in _iter_all_shapes(slide):
            # ── Pre-check: shape text for tags ────────────────────────────
            shape_text = ""
            primary_tag = None
            try:
                if shape.has_text_frame:
                    shape_text = shape.text_frame.text or ""
                    primary_tag = get_primary_tag(shape_text)
            except:
                pass

            # ── {{remove}} tag → mark for deletion ──────────────────────
            if primary_tag == "remove":
                sa.shapes_to_remove.append(shape.name)
                sa.has_tags = True
                continue

            # ── Image placeholders ───────────────────────────────────────
            if shape.is_placeholder:
                pf = shape.placeholder_format
                if pf.type in PICTURE_PH_TYPES:
                    w = shape.width / 914400
                    h = shape.height / 914400
                    img_tag = primary_tag if primary_tag and primary_tag.startswith("image") else None
                    sa.image_placeholders.append(ImagePlaceholderInfo(shape, w, h, tag=img_tag))
                    sa.has_images = True
                    if img_tag:
                        sa.has_tags = True
                    continue

            # Existing images (pics already in template) — image tags can't apply here
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                w = shape.width / 914400
                h = shape.height / 914400
                sa.image_placeholders.append(ImagePlaceholderInfo(shape, w, h))
                sa.has_images = True
                continue

            # ── {{image}} tag inside text shapes (e.g. rectangles) ────────
            if primary_tag and primary_tag.startswith("image") and shape.has_text_frame:
                w = shape.width / 914400
                h = shape.height / 914400
                sa.image_placeholders.append(ImagePlaceholderInfo(shape, w, h, tag=primary_tag))
                sa.has_images = True
                sa.has_tags = True
                continue

            # ── Text frames ──────────────────────────────────────────────
            if not shape.has_text_frame:
                continue

            w_inches = shape.width / 914400
            h_inches = shape.height / 914400
            font_pt = _get_max_font_size(shape)

            # ── Skip decorative shapes (empty text, no tag) ──────────────
            if not primary_tag and not shape_text.strip():
                # Thin bars (accent lines/strips)
                if h_inches < 0.5 or w_inches < 0.5:
                    continue
                # Empty auto shapes (ovals, rounded rects, arrows) are decorative
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    try:
                        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
                        if shape.auto_shape_type in (
                            MSO_AUTO_SHAPE_TYPE.OVAL,
                            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                            MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
                        ):
                            continue
                    except:
                        pass
                # Full-width accent header bars (wide but short empty rectangles)
                slide_w = 13.333
                slide_h = 7.5
                if w_inches > slide_w * 0.8 and h_inches < 2.0:
                    continue
                # Full-slide background rectangles (cover >90% of slide)
                if w_inches > slide_w * 0.85 and h_inches > slide_h * 0.85:
                    continue

            # ── Tag-based role assignment (overrides heuristics) ─────────
            tagged_role = None
            if primary_tag:
                sa.has_tags = True
                if primary_tag in ("title", "topic"):
                    tagged_role = "title"
                    sa.has_title = True
                elif primary_tag == "subtitle":
                    tagged_role = "subtitle"
                    sa.has_subtitle = True
                elif primary_tag == "author":
                    tagged_role = "author"
                elif primary_tag == "date":
                    tagged_role = "date"
                elif primary_tag == "body" or primary_tag.startswith("body_"):
                    tagged_role = "body"
                    sa.has_body = True
                    sa.body_block_count += 1
                elif primary_tag.startswith("point_"):
                    tagged_role = "point"
                    sa.has_body = True

            # ── Fallback role detection (when no tag) ────────────────────
            role = tagged_role or "textbox"
            if not tagged_role:
                if shape.is_placeholder:
                    pf = shape.placeholder_format
                    if pf.type in TITLE_PH_TYPES:
                        role = "title"
                        sa.has_title = True
                    elif pf.type in BODY_PH_TYPES:
                        role = "body"
                        sa.has_body = True
                        sa.body_block_count += 1
                    elif pf.type in SUBTITLE_PH_TYPES:
                        role = "subtitle"
                        sa.has_subtitle = True
                else:
                    # Heuristic for non-placeholder text boxes
                    if font_pt >= 24 and h_inches < 1.5:
                        role = "title"
                        sa.has_title = True
                    elif w_inches * h_inches >= 4.0:
                        role = "body"
                        sa.has_body = True
                        sa.body_block_count += 1
                    elif w_inches < 1.2 or h_inches < 0.3:
                        role = "decorative"  # skip tiny elements
                        continue

            max_chars, max_words = _estimate_capacity(w_inches, h_inches, font_pt, role)
            block = TextBlockInfo(shape, role, max_chars, max_words, w_inches, h_inches, font_pt)
            sa.text_blocks.append(block)
            sa.total_max_chars += max_chars
            sa.total_max_words += max_words

        # Classify slide type
        sa.slide_type = _classify_slide_type(sa, slide_idx, result.total_slides)
        result.slides.append(sa)

    return result


# ══════════════════════════════════════════════════════════════════════════════
# TEXT FORMATTING & TRIMMING
# ══════════════════════════════════════════════════════════════════════════════

def _trim_text(text: str, max_chars: int, max_words: int) -> str:
    """Trim text to fit within character/word limits, cutting at sentence boundary."""
    if not text:
        return text

    # First check word count
    words = text.split()
    if len(words) > max_words:
        text = " ".join(words[:max_words])
        # Try to end at sentence boundary
        last_period = text.rfind(".")
        last_semicolon = text.rfind(";")
        cut_point = max(last_period, last_semicolon)
        if cut_point > len(text) * 0.6:
            text = text[:cut_point + 1]
        else:
            text = text.rstrip(" ,;:") + "..."

    # Then check char count
    if len(text) > max_chars:
        text = text[:max_chars]
        last_period = text.rfind(".")
        last_space = text.rfind(" ")
        cut_point = max(last_period, last_space)
        if cut_point > len(text) * 0.5:
            text = text[:cut_point + 1].rstrip()
        else:
            text = text.rstrip(" ,;:") + "..."

    return text.strip()


def _capture_run_font(run) -> dict:
    """Capture all font properties from a Run."""
    f = run.font
    props = {}
    try:
        if f.size is not None:
            props["size"] = f.size
    except:
        pass
    try:
        if f.bold is not None:
            props["bold"] = f.bold
    except:
        pass
    try:
        if f.italic is not None:
            props["italic"] = f.italic
    except:
        pass
    try:
        if f.name is not None:
            props["name"] = f.name
    except:
        pass
    try:
        if f.color and f.color.type is not None:
            ctype = str(f.color.type)
            if "RGB" in ctype or f.color.type == 1:
                props["color_rgb"] = f.color.rgb
            elif "SCHEME" in ctype or f.color.type == 2:
                props["color_theme"] = getattr(f.color, "theme_color", None)
    except:
        pass
    return props


def _apply_run_font(run, props: dict):
    """Apply captured font properties to a Run."""
    f = run.font
    try:
        if "size" in props:
            f.size = props["size"]
        if "bold" in props:
            f.bold = props["bold"]
        if "italic" in props:
            f.italic = props["italic"]
        if "name" in props:
            f.name = props["name"]
        if "color_rgb" in props:
            f.color.rgb = props["color_rgb"]
        elif "color_theme" in props and props["color_theme"]:
            f.color.theme_color = props["color_theme"]
    except:
        pass


def _fill_shape_text(shape, text: str, max_chars: int = 0, max_words: int = 0,
                     is_title: bool = False):
    """
    Fill a shape's text frame with new text, preserving original formatting.
    Auto-trims if text exceeds capacity.
    Works correctly with PowerPoint placeholders (does NOT use tf.clear()).
    """
    if not shape.has_text_frame:
        return

    tf = shape.text_frame

    # Capture formatting from first paragraph/run
    font_props = {}
    alignment = None
    if tf.paragraphs:
        p0 = tf.paragraphs[0]
        alignment = p0.alignment
        if p0.runs:
            font_props = _capture_run_font(p0.runs[0])

    # Trim text if needed
    if max_chars > 0 or max_words > 0:
        text = _trim_text(text, max_chars or 99999, max_words or 99999)

    # Handle bullet points (list of items)
    if isinstance(text, list):
        items = text
    elif "\n" in text and not is_title:
        items = [line.strip() for line in text.split("\n") if line.strip()]
    else:
        items = [text]

    # ── Fill using direct XML manipulation to preserve placeholder status ──
    from pptx.oxml.ns import qn
    from lxml import etree

    txBody = tf._txBody
    if txBody is None:
        return

    # Remove ALL existing <a:p> paragraphs
    for old_p in txBody.findall(qn('a:p')):
        txBody.remove(old_p)

    # Insert new paragraphs with text
    for i, item in enumerate(items):
        item_text = str(item).strip()

        # Trim each item individually
        if max_chars > 0 and len(items) > 1:
            per_item_chars = max(20, max_chars // len(items))
            per_item_words = max(5, max_words // len(items))
            item_text = _trim_text(item_text, per_item_chars, per_item_words)

        # Add bullet for body text if not already present
        if not is_title and item_text:
            first_char = item_text[0] if item_text else ''
            has_bullet = first_char in ('\u2022', '\u2013', '-', '\u2713', '\u2714', '\u25ba', '\u25b8', '\u25cf')
            has_number = len(item_text) > 1 and item_text[0].isdigit() and item_text[1] in '.)'
            if not has_bullet and not has_number:
                item_text = '\u2022  ' + item_text

        # Build paragraph XML: <a:p><a:pPr/><a:r><a:rPr/><a:t>text</a:t></a:r></a:p>
        p_elem = etree.SubElement(txBody, qn('a:p'))

        # Add paragraph properties (alignment + spacing)
        pPr = etree.SubElement(p_elem, qn('a:pPr'))
        if alignment is not None:
            algn_map = {
                PP_ALIGN.LEFT: 'l', PP_ALIGN.CENTER: 'ctr',
                PP_ALIGN.RIGHT: 'r', PP_ALIGN.JUSTIFY: 'just',
            }
            algn_val = algn_map.get(alignment, 'l')
            pPr.set('algn', algn_val)

        # Add paragraph spacing for body text
        if not is_title:
            spc_aft = etree.SubElement(pPr, qn('a:spcAft'))
            spc_pts = etree.SubElement(spc_aft, qn('a:spcPts'))
            spc_pts.set('val', '800')  # 8pt after each paragraph

        # Add run with text
        r_elem = etree.SubElement(p_elem, qn('a:r'))

        # Add run properties (font)
        rPr = etree.SubElement(r_elem, qn('a:rPr'))
        rPr.set('lang', 'uz-UZ')
        rPr.set('dirty', '0')
        if font_props.get('bold'):
            rPr.set('b', '1')
        if font_props.get('italic'):
            rPr.set('i', '1')
        if font_props.get('size'):
            rPr.set('sz', str(int(font_props['size'].pt * 100)))
        if font_props.get('name'):
            latin = etree.SubElement(rPr, qn('a:latin'))
            latin.set('typeface', font_props['name'])

        # Add text element
        t_elem = etree.SubElement(r_elem, qn('a:t'))
        t_elem.text = item_text

    # Set word wrap
    try:
        tf.word_wrap = True
    except:
        pass


# ══════════════════════════════════════════════════════════════════════════════
# IMAGE INSERTION
# ══════════════════════════════════════════════════════════════════════════════

def _insert_image_into_placeholder(shape, img_bytes: bytes) -> bool:
    """Try to insert image into a picture placeholder. Returns True on success."""
    try:
        stream = io.BytesIO(img_bytes)
        shape.insert_picture(stream)
        return True
    except:
        return False


def _replace_existing_image(slide, old_shape, img_bytes: bytes):
    """Replace an existing image shape with new image at same position/size."""
    left = old_shape.left
    top = old_shape.top
    width = old_shape.width
    height = old_shape.height

    try:
        # Remove old shape
        sp = old_shape._element
        sp.getparent().remove(sp)
    except:
        pass

    try:
        stream = io.BytesIO(img_bytes)
        slide.shapes.add_picture(stream, left, top, width, height)
    except Exception as e:
        logger.error(f"Image replacement error: {e}")


# ══════════════════════════════════════════════════════════════════════════════
# MAIN: FILL TEMPLATE
# ══════════════════════════════════════════════════════════════════════════════

def _resolve_tag_content(
    tag: str,
    content: dict,
    topic: str = "",
    author: str = "",
    date_str: str = "",
    is_first_slide: bool = False,
) -> Optional[str]:
    """
    Map a {{tag}} to its content from the AI-generated dict.

    Supported tags:
        title, topic   → content["title"] (or topic on slide 0)
        subtitle       → content["subtitle"]
        author         → author name (or content["author"])
        date           → today's date
        body           → content["body"] or joined content["points"]
        body_N         → content["body_N"] or content["points"][N-1]
        point_N        → content["points"][N-1] or content["point_N"]

    Returns the text to fill, or None if no matching content was found
    (in which case the original tag will remain — handler should clear it).
    """
    tag = tag.lower().strip()

    # ── Direct mappings ──────────────────────────────────────────────
    if tag in ("title", "topic"):
        text = content.get("title", "")
        if not text and is_first_slide:
            text = topic
        return text or None

    if tag == "subtitle":
        text = content.get("subtitle", "")
        if is_first_slide and author:
            # Title slide: always show author
            if text:
                text = f"{text}\nBajardi: {author}"
            else:
                text = f"Bajardi: {author}"
        return text or None

    if tag == "author":
        return content.get("author") or author or None

    if tag == "date":
        return content.get("date") or date_str or None

    if tag == "body":
        body = content.get("body") or content.get("content")
        if body:
            if isinstance(body, list):
                return "\n".join(str(p) for p in body)
            return str(body)
        # fallback: join points
        points = content.get("points")
        if isinstance(points, list) and points:
            return "\n".join(str(p) for p in points)
        return None

    # ── Indexed body_N tags ─────────────────────────────────────────
    if tag.startswith("body_"):
        # Try direct key first: content["body_1"]
        direct = content.get(tag)
        if direct:
            return str(direct) if not isinstance(direct, list) else "\n".join(str(p) for p in direct)
        # Fallback: take Nth element from points
        try:
            idx = int(tag.split("_", 1)[1]) - 1
            points = content.get("points")
            if isinstance(points, list) and 0 <= idx < len(points):
                return str(points[idx])
        except (ValueError, IndexError):
            pass
        return None

    # ── Indexed subtitle_N tags ───────────────────────────────────────
    if tag.startswith("subtitle_"):
        direct = content.get(tag)
        if direct:
            return str(direct)
        # Fallback: take Nth element from points
        try:
            idx = int(tag.split("_", 1)[1]) - 1
            points = content.get("points")
            if isinstance(points, list) and 0 <= idx < len(points):
                return str(points[idx])
        except (ValueError, IndexError):
            pass
        return None

    # ── Indexed point_N tags ─────────────────────────────────────────
    if tag.startswith("point_"):
        # Try direct key
        direct = content.get(tag)
        if direct:
            return str(direct)
        # Try points[N-1]
        try:
            idx = int(tag.split("_", 1)[1]) - 1
            points = content.get("points")
            if isinstance(points, list) and 0 <= idx < len(points):
                return str(points[idx])
        except (ValueError, IndexError):
            pass
        return None

    # ── Image tags handled elsewhere; return None for any other ──────
    if tag.startswith("image"):
        return None

    # Unknown tag — try direct lookup as last resort
    raw = content.get(tag)
    if raw:
        return str(raw) if not isinstance(raw, list) else "\n".join(str(p) for p in raw)
    return None


def fill_template(
    template_path: str,
    slides_content: List[dict],
    topic: str,
    author: str = "",
    slide_images: Optional[dict] = None,
    analysis: Optional[TemplateAnalysis] = None,
    target_slides: Optional[int] = None,
    subject: str = "",
    university: str = "",
    doc_type_label: str = "",
) -> bytes:
    """
    Fill an existing PPTX template with AI-generated content.

    This does NOT clone or delete slides. It fills each slide IN-PLACE,
    preserving the original design.

    Args:
        template_path: path to .pptx template
        slides_content: list of dicts, one per slide:
            {
                "title": "Slide title",
                "subtitle": "Optional subtitle",
                "body": "Main text or list of bullet points",
                "points": ["Point 1", "Point 2"],    # alternative to body
                "image_keyword": "search keyword",     # for auto image search
            }
        topic: presentation topic (for title slide)
        author: author name (for title slide)
        slide_images: {slide_index: image_bytes} for custom images
        analysis: pre-computed TemplateAnalysis (if None, will be computed)
        target_slides: desired number of slides; excess slides removed from end

    Returns:
        bytes of the filled .pptx file
    """
    template_path = _resolve_path(template_path)
    if not template_path or not os.path.exists(template_path):
        raise FileNotFoundError(f"Template not found: {template_path}")

    # Analyze if not provided
    if analysis is None:
        analysis = analyze_template_deep(template_path)

    prs = Presentation(template_path)
    slide_images = slide_images or {}

    slides = list(prs.slides)
    num_template_slides = len(slides)
    num_content = len(slides_content)

    logger.info(f"Filling template: {num_template_slides} slides, {num_content} content items")

    # Academic title slide overlay for first slide
    if (subject or university) and len(slides) > 0:
        from services.pptx_service import build_academic_title_slide
        build_academic_title_slide(
            slides[0], topic=topic, author=author,
            subject=subject, university=university,
            doc_type_label=doc_type_label, prs=prs
        )

    for slide_idx, slide in enumerate(slides):
        if slide_idx >= len(analysis.slides):
            break

        sa = analysis.slides[slide_idx]

        # Build a lookup: shape_name → shape in current prs
        shape_by_name = {shape.name: shape for shape in _iter_all_shapes(slide)}

        # ── Step 1: Remove shapes tagged with {{remove}} (watermarks etc) ──
        for shape_name in sa.shapes_to_remove:
            shape_to_kill = shape_by_name.get(shape_name)
            if shape_to_kill is not None:
                try:
                    sp = shape_to_kill._element
                    sp.getparent().remove(sp)
                    logger.debug(f"Slide {slide_idx}: removed shape '{shape_name}' (tag=remove)")
                except Exception as e:
                    logger.warning(f"Slide {slide_idx}: failed to remove '{shape_name}': {e}")

        # Get content for this slide
        if slide_idx < num_content:
            content = slides_content[slide_idx]
            if not isinstance(content, dict):
                content = {"body": str(content)} if content else {}
        else:
            content = {}

        # Pre-compute helpers for tag-based filling
        from datetime import datetime
        date_str = datetime.now().strftime("%d.%m.%Y")
        body_content_default = content.get("content") or content.get("points") or content.get("body", "")
        points_list = content.get("points") if isinstance(content.get("points"), list) else None

        # ── Step 2: Fill text blocks ──────────────────────────────────────
        for block in sa.text_blocks:
            # Find the actual shape in the CURRENT presentation by name
            actual_shape = shape_by_name.get(block.shape_name)
            if actual_shape is None or not actual_shape.has_text_frame:
                continue

            # ── Tag-based filling (precise) ──────────────────────────────
            if block.tag:
                # If shape has MULTIPLE tags, skip here — second pass handles inline
                shape_full_text = actual_shape.text_frame.text or ""
                all_tags_in_shape = TAG_PATTERN.findall(shape_full_text)
                if len(all_tags_in_shape) > 1:
                    continue  # handled by inline second pass

                text_to_fill = _resolve_tag_content(
                    tag=block.tag,
                    content=content,
                    topic=topic,
                    author=author,
                    date_str=date_str,
                    is_first_slide=(slide_idx == 0),
                )
                if text_to_fill is not None:
                    is_title_like = block.tag in ("title", "topic")
                    _fill_shape_text(
                        actual_shape, text_to_fill,
                        max_chars=block.max_chars,
                        max_words=block.max_words,
                        is_title=is_title_like
                    )
                else:
                    # No content for this tag — clear it so raw {{tag}} doesn't show
                    _fill_shape_text(actual_shape, "", max_chars=0, max_words=0)
                continue

            # ── Legacy role-based filling (when no tag) ──────────────────
            if block.role == "title":
                title_text = content.get("title", "")
                if not title_text and slide_idx == 0:
                    title_text = topic
                if title_text:
                    _fill_shape_text(
                        actual_shape, title_text,
                        max_chars=block.max_chars,
                        max_words=block.max_words,
                        is_title=True
                    )

            elif block.role == "subtitle":
                sub_text = content.get("subtitle", "")
                if not sub_text and slide_idx == 0:
                    sub_text = f"Bajardi: {author}" if author else ""
                if sub_text:
                    _fill_shape_text(
                        actual_shape, sub_text,
                        max_chars=block.max_chars,
                        max_words=block.max_words
                    )

            elif block.role in ("body", "textbox"):
                body_content = body_content_default

                if isinstance(body_content, list):
                    joined = "\n".join(str(p) for p in body_content)
                    _fill_shape_text(
                        actual_shape, joined,
                        max_chars=block.max_chars,
                        max_words=block.max_words
                    )
                elif body_content:
                    _fill_shape_text(
                        actual_shape, str(body_content),
                        max_chars=block.max_chars,
                        max_words=block.max_words
                    )

        # ── Step 3: Insert images ────────────────────────────────────────
        if sa.image_placeholders:
            # Build a list of (image_bytes, target_placeholder) pairs
            # Priority:
            #   1. content["images"] dict keyed by tag ({"image_1": bytes, ...})
            #   2. slide_images[slide_idx] → first image placeholder
            images_by_tag = content.get("_images_by_tag") or {}

            def _do_insert(placeholder_info, img_bytes):
                """Insert image into any type of placeholder shape."""
                actual_shape = shape_by_name.get(placeholder_info.shape_name)
                if not actual_shape:
                    return
                if actual_shape.is_placeholder:
                    if not _insert_image_into_placeholder(actual_shape, img_bytes):
                        _replace_existing_image(slide, actual_shape, img_bytes)
                else:
                    _replace_existing_image(slide, actual_shape, img_bytes)

            # Insert by tag first
            inserted = set()
            for placeholder in sa.image_placeholders:
                if placeholder.tag and placeholder.tag in images_by_tag:
                    _do_insert(placeholder, images_by_tag[placeholder.tag])
                    inserted.add(placeholder.shape_name)

            # Fallback: insert slide_images into ALL uninserted placeholders
            # Handler uses 1-based keys (slide 1, 2, 3...), so check both conventions
            img_bytes_fallback = slide_images.get(slide_idx) or slide_images.get(slide_idx + 1)
            if img_bytes_fallback:
                for placeholder in sa.image_placeholders:
                    if placeholder.shape_name in inserted:
                        continue
                    _do_insert(placeholder, img_bytes_fallback)
                    inserted.add(placeholder.shape_name)

    # ── Second pass: replace any remaining {{tag}} patterns inline ────────
    from datetime import datetime
    date_str = datetime.now().strftime("%d.%m.%Y")
    for slide_idx, slide in enumerate(slides):
        if slide_idx >= num_content:
            break
        content = slides_content[slide_idx]
        if not isinstance(content, dict):
            content = {"body": str(content)} if content else {}
        for shape in _iter_all_shapes(slide):
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                full_text = para.text or ""
                remaining_tags = TAG_PATTERN.findall(full_text)
                if not remaining_tags:
                    continue
                # Build replaced text from full paragraph
                new_text = full_text
                for tag_name in remaining_tags:
                    tag_lower = tag_name.lower()
                    replacement = _resolve_tag_content(
                        tag=tag_lower, content=content, topic=topic,
                        author=author, date_str=date_str,
                        is_first_slide=(slide_idx == 0),
                    )
                    if replacement is None:
                        replacement = ""
                    pat = re.compile(r"\{\{\s*" + re.escape(tag_name) + r"\s*\}\}", re.IGNORECASE)
                    new_text = pat.sub(replacement, new_text)
                # Apply: put all text into first run, clear the rest
                if new_text != full_text and para.runs:
                    para.runs[0].text = new_text
                    for run in para.runs[1:]:
                        run.text = ""

    # ── Duplicate slides if target_slides > template slides ─────────────
    if target_slides and target_slides > len(prs.slides):
        needed = target_slides - len(prs.slides)
        logger.info(f"Duplicating {needed} slides (template={len(prs.slides)}, target={target_slides})")
        # Find last content slide to duplicate (skip closing/thank-you slides at end)
        last_content_idx = len(prs.slides) - 1
        for check_idx in range(len(prs.slides) - 1, 0, -1):
            if check_idx < len(analysis.slides):
                stype = analysis.slides[check_idx].slide_type
                if stype not in ("closing", "sources", "conclusion"):
                    last_content_idx = check_idx
                    break

        try:
            import copy
            from lxml import etree
            for dup_i in range(needed):
                # Copy the slide layout and content
                src_slide = prs.slides[last_content_idx]
                slide_layout = src_slide.slide_layout
                new_slide = prs.slides.add_slide(slide_layout)
                # Copy all elements from source
                for elem in src_slide.shapes:
                    el = copy.deepcopy(elem._element)
                    new_slide.shapes._spTree.append(el)
                # Overwrite text with extra content (tags already replaced, so write directly)
                extra_idx = num_template_slides + dup_i
                if extra_idx < num_content:
                    extra_content = slides_content[extra_idx]
                    if isinstance(extra_content, dict):
                        new_title = extra_content.get("title", "")
                        new_body = extra_content.get("content", extra_content.get("body", ""))
                        if not new_body:
                            pts = extra_content.get("points", [])
                            if pts:
                                new_body = "\n".join(f"• {p}" for p in pts if p)
                        title_set = False
                        body_set = False
                        for shape in _iter_all_shapes(new_slide):
                            if not shape.has_text_frame:
                                continue
                            tf = shape.text_frame
                            # Detect title vs body by placeholder type or font size
                            is_title_shape = False
                            if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                                ph_idx = shape.placeholder_format.idx
                                if ph_idx in (0, 13):  # title placeholders
                                    is_title_shape = True
                            elif tf.paragraphs and tf.paragraphs[0].runs:
                                fsize = tf.paragraphs[0].runs[0].font.size
                                if fsize and fsize.pt >= 20:
                                    is_title_shape = True
                            if is_title_shape and new_title and not title_set:
                                for para in tf.paragraphs:
                                    if para.runs:
                                        para.runs[0].text = new_title
                                        for r in para.runs[1:]:
                                            r.text = ""
                                title_set = True
                            elif not is_title_shape and new_body and not body_set:
                                for para in tf.paragraphs:
                                    if para.runs:
                                        para.runs[0].text = new_body
                                        for r in para.runs[1:]:
                                            r.text = ""
                                        new_body = ""  # only first body shape
                                body_set = True
                logger.debug(f"Duplicated slide {last_content_idx} → new slide with content idx {extra_idx}")
        except Exception as e:
            logger.warning(f"Slide duplication failed: {e}")

    # ── Remove excess slides if target_slides is set ──────────────────────
    if target_slides and target_slides < len(prs.slides):
        excess = len(prs.slides) - target_slides
        logger.info(f"Removing {excess} excess slides (template={len(prs.slides)}, target={target_slides})")
        # Remove slides from the end, but keep the very last slide if it's
        # a closing/thank-you slide — remove from just before it instead
        for _ in range(excess):
            # Remove the second-to-last slide (keep closing slide at end)
            remove_idx = len(prs.slides) - 2 if len(prs.slides) > 2 else len(prs.slides) - 1
            try:
                rId = prs.slides._sldIdLst[remove_idx].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[remove_idx]
            except Exception as e:
                logger.warning(f"Failed to remove slide {remove_idx}: {e}")
                break

    # Save
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
