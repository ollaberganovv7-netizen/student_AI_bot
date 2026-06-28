from __future__ import annotations
"""
Lecture PPTX Service
--------------------
Generates university-standard lecture presentations (1:1 copy of lecture_standard.pptx structure).

Structure:
  1. Titulniy (title)
  2. Reja (plan — 3 o'quv savol)
  3. Kirish (introduction)
  4. Maqsad va vazifalar (goals)
  5. O'quv savollari (questions overview)
  6. Adabiyotlar (references)
  7-13. I O'quv savoli (section header + content slides + mini conclusion)
  14-21. II O'quv savoli (same)
  22-29. III O'quv savoli (same)
  30. Iqtibos (quote)
  31. Xulosa (conclusion)
  32. Yakuniy (final slide)

Uses lecture_standard.pptx as the base template — clones its design exactly.
"""

import io
import os
import copy
import logging
from typing import List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from lxml import etree

logger = logging.getLogger(__name__)

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
LECTURE_TEMPLATE = os.path.join(BASE_DIR, "doc_examples", "lecture_standard.pptx")

# Slide indices in the template (0-based) — reference slides to clone from
REF_TITLE_SLIDE = 0          # Titulniy
REF_PLAN_SLIDE = 1            # Reja
REF_SECTION_HEADER_SLIDE = 6  # "BIRINCHI O'QUV SAVOLI" header
REF_CONTENT_SLIDE = 7         # Content slide with text
REF_CONTENT_2COL_SLIDE = 10   # Content with image/2 columns
REF_MINI_CONCLUSION = 25      # "II-O'QUV SAVOLGA XULOSA"  
REF_QUOTE_SLIDE = 31          # Quote slide
REF_CONCLUSION_SLIDE = 32     # XULOSA
REF_FINAL_SLIDE = 33          # Yakuniy


def _deep_copy_slide(prs, template_slide):
    """Deep copy a slide from the same presentation, preserving all formatting."""
    # Get the slide layout of the template slide
    slide_layout = template_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)
    
    # Copy background
    if template_slide.background.fill.type is not None:
        try:
            bg_elem = template_slide.background._element
            new_bg = copy.deepcopy(bg_elem)
            new_slide.background._element.getparent().replace(
                new_slide.background._element, new_bg
            )
        except:
            pass

    # Remove default placeholders from new slide
    for ph in list(new_slide.placeholders):
        try:
            sp = ph._element
            sp.getparent().remove(sp)
        except:
            pass

    # Copy all shapes from template
    for shape in template_slide.shapes:
        try:
            el = copy.deepcopy(shape._element)
            new_slide.shapes._spTree.append(el)
        except Exception as e:
            logger.warning(f"Could not copy shape {shape.name}: {e}")

    return new_slide


def _replace_text_in_slide(slide, replacements: dict):
    """
    Replace text in ALL text frames of a slide.
    replacements: {old_text_substring: new_text}
    Preserves formatting of the first run.
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            full_text = ''.join(run.text for run in para.runs)
            replaced = False
            new_text = full_text
            for old, new in replacements.items():
                if old in new_text:
                    new_text = new_text.replace(old, new)
                    replaced = True
            
            if replaced and para.runs:
                # Capture first run font
                first_run = para.runs[0]
                font_props = _capture_font(first_run)
                # Clear all runs
                for run in para.runs[1:]:
                    run.text = ""
                first_run.text = new_text
                _apply_font(first_run, font_props)


def _set_slide_text(slide, texts: dict):
    """
    Set text in slide shapes by matching shape name patterns.
    texts: {"title": "...", "body": "...", "section_label": "..."}
    """
    shapes_by_area = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        w = shape.width / 914400
        h = shape.height / 914400
        area = w * h
        text = shape.text_frame.text.strip()
        shapes_by_area.append((area, shape, text))
    
    # Sort by area — largest = body, medium = title, small = labels
    shapes_by_area.sort(key=lambda x: x[0], reverse=True)
    
    filled_title = False
    filled_body = False
    
    for area, shape, original_text in shapes_by_area:
        if not shape.has_text_frame:
            continue
        
        # Skip slide number shapes
        name_lower = shape.name.lower()
        if "номер" in name_lower or "number" in name_lower or "слайд" in name_lower:
            continue
        
        # Skip very small decorative shapes
        if area < 0.5:
            continue
            
        tf = shape.text_frame
        
        # Try to identify by content/size
        if not filled_body and area > 8 and "body" in texts:
            _fill_text_frame(tf, texts["body"])
            filled_body = True
        elif not filled_title and "title" in texts and area > 2:
            _fill_text_frame(tf, texts["title"])
            filled_title = True
        elif "section_label" in texts and area < 3:
            _fill_text_frame(tf, texts["section_label"])


def _fill_text_frame(tf, text: str):
    """Fill a text frame preserving first paragraph formatting."""
    if not tf.paragraphs:
        return
    
    # Capture formatting
    p0 = tf.paragraphs[0]
    alignment = p0.alignment
    font_props = {}
    if p0.runs:
        font_props = _capture_font(p0.runs[0])
    
    # Handle list/multiline
    if isinstance(text, list):
        lines = text
    elif "\n" in text:
        lines = [l.strip() for l in text.split("\n") if l.strip()]
    else:
        lines = [text]
    
    # Clear and refill
    tf.clear()
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = str(line)
        if alignment is not None:
            para.alignment = alignment
        for run in para.runs:
            _apply_font(run, font_props)
    
    try:
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except:
        pass


def _capture_font(run) -> dict:
    props = {}
    try:
        if run.font.size is not None: props['size'] = run.font.size
    except: pass
    try:
        if run.font.bold is not None: props['bold'] = run.font.bold
    except: pass
    try:
        if run.font.italic is not None: props['italic'] = run.font.italic
    except: pass
    try:
        if run.font.name is not None: props['name'] = run.font.name
    except: pass
    try:
        if run.font.color and run.font.color.rgb is not None:
            props['color_rgb'] = run.font.color.rgb
    except: pass
    return props


def _apply_font(run, props: dict):
    try:
        if 'size' in props: run.font.size = props['size']
        if 'bold' in props: run.font.bold = props['bold']
        if 'italic' in props: run.font.italic = props['italic']
        if 'name' in props: run.font.name = props['name']
        if 'color_rgb' in props: run.font.color.rgb = props['color_rgb']
    except: pass


def generate_lecture_pptx(
    content: dict,
    template_path: str = None,
) -> bytes:
    """
    Generate a university-standard lecture PPTX.

    Args:
        content: dict with ALL lecture content:
        {
            "university": "O'ZBEKISTON RESPUBLIKASI ...",
            "kafedra": "Harbiy pedagogika kafedrasi",
            "fan": "Harbiy pedagogika",
            "topic": "MA'RUZA MASHG'ULOTLARINI ...",
            "author": "dots. Aliyev A.A.",
            "year": "2026",
            "city": "Toshkent",

            "plan": ["1-savol nomi", "2-savol nomi", "3-savol nomi"],

            "kirish": "Mavzuning dolzarbligi ...",
            "maqsad": {
                "talimiy": "...",
                "tarbiyaviy": "...",
                "rivojlantiruvchi": "..."
            },
            "adabiyotlar": ["Kitob 1", "Kitob 2", ...],

            "savol_1": {
                "title": "Ma'ruza mashg'ulotlarning mohiyati va maqsadi",
                "slides": [
                    {"title": "...", "body": "bullet1\nbullet2\nbullet3"},
                    {"title": "...", "body": "..."},
                ],
                "xulosa": "I-savol bo'yicha xulosa matni"
            },
            "savol_2": { ... same structure ... },
            "savol_3": { ... same structure ... },

            "iqtibos": {"text": "Tsitata matni...", "author": "Muallif ismi"},
            "xulosa": "Umumiy xulosa matni...",
        }
        template_path: optional custom template, defaults to lecture_standard.pptx

    Returns:
        bytes of generated .pptx
    """
    tpl = template_path or LECTURE_TEMPLATE
    if not os.path.exists(tpl):
        raise FileNotFoundError(f"Lecture template not found: {tpl}")

    prs = Presentation(tpl)
    
    # We'll work with the existing slides — replace their text content
    slides = list(prs.slides)
    total_existing = len(slides)
    
    logger.info(f"Lecture generation: template has {total_existing} slides")

    # ── Strategy: Replace text in existing slides where possible,
    #    clone reference slides for additional content slides ──
    
    # SLIDE 1: Title
    if total_existing > 0:
        _replace_text_in_slide(slides[0], {
            "HARBIY TA'LIM MUASSALARIDA MA'RUZA": content.get("topic", "MAVZU"),
            "MASHG'ULOTLARINI TASHKILLASHTIRISH VA O'TKAZISH METODIKASI": "",
            "Toshkent 2026": f"{content.get('city', 'Toshkent')} {content.get('year', '2026')}",
        })
    
    # For the rest, we build a new presentation using the template's design
    # The approach: keep slides 0 (title) and last (final), rebuild middle content
    
    # Instead of complex slide manipulation, let's use a simpler approach:
    # Fill ALL text shapes in each slide with the appropriate content
    
    # Build ordered content list matching slide positions
    ordered_content = _build_slide_content_list(content)
    
    # Now fill each slide
    for idx, slide in enumerate(slides):
        if idx >= len(ordered_content):
            break
        sc = ordered_content[idx]
        if sc is None:
            continue  # skip — keep original
        _replace_all_text_shapes(slide, sc)
    
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _build_slide_content_list(content: dict) -> list:
    """Build list of content dicts, one per slide in the template."""
    result = []
    
    # Slide 1: Title (index 0) — handle separately above
    result.append(None)  # skip, already handled
    
    # Slide 2: Reja (index 1)
    plan = content.get("plan", [])
    plan_text = "\n".join(f"{i+1}. {p}" for i, p in enumerate(plan))
    result.append({"main_text": plan_text})
    
    # Slide 3: Kirish (index 2)
    result.append({"main_text": content.get("kirish", "")})
    
    # Slide 4: Maqsad (index 3)
    maqsad = content.get("maqsad", {})
    if isinstance(maqsad, dict):
        maqsad_text = (
            f"Ta'limiy maqsad: {maqsad.get('talimiy', '')}\n"
            f"Tarbiyaviy maqsad: {maqsad.get('tarbiyaviy', '')}\n"
            f"Rivojlantiruvchi maqsad: {maqsad.get('rivojlantiruvchi', '')}"
        )
    else:
        maqsad_text = str(maqsad)
    result.append({"main_text": maqsad_text})
    
    # Slide 5: O'quv savollari (index 4)
    savol_overview = "\n".join(
        f"{i+1}-savol: {p}" for i, p in enumerate(plan)
    )
    result.append({"main_text": savol_overview})
    
    # Slide 6: Adabiyotlar (index 5)
    adab = content.get("adabiyotlar", [])
    adab_text = "\n".join(f"{i+1}. {a}" for i, a in enumerate(adab))
    result.append({"main_text": adab_text})
    
    # I O'quv savoli slides
    result.extend(_build_savol_slides(content.get("savol_1", {}), "I"))
    
    # II O'quv savoli slides
    result.extend(_build_savol_slides(content.get("savol_2", {}), "II"))
    
    # III O'quv savoli slides
    result.extend(_build_savol_slides(content.get("savol_3", {}), "III"))
    
    # Quote slide
    iqtibos = content.get("iqtibos", {})
    if isinstance(iqtibos, dict):
        q_text = f"{iqtibos.get('text', '')}\n\n— {iqtibos.get('author', '')}"
    else:
        q_text = str(iqtibos)
    result.append({"main_text": q_text})
    
    # Xulosa
    result.append({"main_text": content.get("xulosa", "")})
    
    # Final slide
    result.append(None)  # keep original design
    
    return result


def _build_savol_slides(savol: dict, roman: str) -> list:
    """Build content list for one o'quv savol section."""
    result = []
    
    title = savol.get("title", f"{roman}-O'QUV SAVOLI")
    slides_content = savol.get("slides", [])
    xulosa = savol.get("xulosa", "")
    
    # Section header
    result.append({"main_text": f"{roman}-O'QUV SAVOLI\n{title}"})
    
    # Content slides
    for sc in slides_content:
        if isinstance(sc, dict):
            body = sc.get("body", sc.get("points", ""))
            if isinstance(body, list):
                body = "\n".join(body)
            text = f"{sc.get('title', '')}\n\n{body}"
        else:
            text = str(sc)
        result.append({"main_text": text})
    
    # Mini conclusion
    result.append({"main_text": f"{roman}-O'QUV SAVOLGA XULOSA\n\n{xulosa}"})
    
    return result


def _replace_all_text_shapes(slide, content: dict):
    """Replace text in the largest text shapes of a slide."""
    main_text = content.get("main_text", "")
    if not main_text:
        return
    
    # Find all text shapes, sorted by area (largest first)
    text_shapes = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        name = shape.name.lower()
        # Skip slide numbers, footers, tiny labels
        if any(x in name for x in ["номер", "number", "колонтитул", "footer", "дата", "date"]):
            continue
        w = shape.width / 914400
        h = shape.height / 914400
        area = w * h
        if area < 0.5:
            continue
        text_shapes.append((area, shape))
    
    text_shapes.sort(key=lambda x: x[0], reverse=True)
    
    if not text_shapes:
        return
    
    # Fill the largest text shape with main content
    biggest = text_shapes[0][1]
    _fill_text_frame(biggest.text_frame, main_text)
