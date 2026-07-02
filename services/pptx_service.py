from __future__ import annotations
import io
import os
import copy
from typing import List, Optional
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor


def analyze_template(template_path: str) -> dict:
    """
    Analyze a PPTX template and return detailed info about each slide's structure.
    Reports per-block word estimates so AI can generate properly sized content.
    
    Returns:
        {"total_slides": int, "slides": [{"index": int, "text_areas": int, 
         "has_title": bool, "has_body": bool, "body_blocks": int,
         "blocks": [{"words": int, "type": "body"|"textbox"}, ...],
         "estimated_words": int}, ...]}
    """
    if template_path and not os.path.isabs(template_path):
        base = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        full = os.path.join(base, template_path)
        if os.path.exists(full):
            template_path = full

    if not template_path or not os.path.exists(template_path):
        return {"total_slides": 0, "slides": []}

    prs = Presentation(template_path)
    result = {"total_slides": len(prs.slides), "slides": []}
    
    for i, slide in enumerate(prs.slides):
        info = {"index": i, "text_areas": 0, "has_title": False, 
                "has_body": False, "body_blocks": 0, "blocks": [],
                "estimated_words": 0}
        
        # Sort shapes geographically (left-to-right, then top-to-bottom)
        # We round X to nearest ~0.5 inch (457200 EMU) to group columns together
        sorted_shapes = sorted(
            [s for s in slide.shapes if s.has_text_frame],
            key=lambda s: (round(getattr(s, 'left', 0) / 457200), getattr(s, 'top', 0))
        )
        
        for shape in sorted_shapes:
            info["text_areas"] += 1
            if shape.is_placeholder:
                pf = shape.placeholder_format
                if pf.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                    info["has_title"] = True
                elif pf.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
                    info["has_body"] = True
                    info["body_blocks"] += 1
                    # Estimate capacity from placeholder size
                    w_inches = shape.width / 914400  # EMU to inches
                    h_inches = shape.height / 914400
                    area = w_inches * h_inches
                    block_words = max(15, int(area * 10))
                    info["blocks"].append({"words": block_words, "type": "body"})
                    info["estimated_words"] += block_words
            else:
                # Regular text box
                w_inches = shape.width / 914400
                h_inches = shape.height / 914400
                area = w_inches * h_inches
                block_words = max(10, int(area * 8))
                info["blocks"].append({"words": block_words, "type": "textbox"})
                info["estimated_words"] += block_words
        
        result["slides"].append(info)
    
    return result

# ── Slide shape analysis ─────────────────────────────────────────────────────

TITLE_PH = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}
BODY_PH  = {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT}
SUB_PH   = {PP_PLACEHOLDER.SUBTITLE}

def _classify_shapes(slide):
    """Return dict of shape lists: title, body, subtitle, textboxes."""
    title, body, subtitle, textboxes = [], [], [], []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.is_placeholder:
            pf = shape.placeholder_format
            if pf.type in TITLE_PH:
                title.append(shape)
            elif pf.type in BODY_PH:
                body.append(shape)
            elif pf.type in SUB_PH:
                subtitle.append(shape)
            else:
                textboxes.append(shape)
        else:
            textboxes.append(shape)
    return {"title": title, "body": body, "subtitle": subtitle, "other": textboxes}


# ── Smart Template Engine (auto-detect title/body in ANY template) ────────────

def _get_font_size_from_shape(shape):
    """Get the largest font size used in a shape (in Pt)."""
    max_size = 0
    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    size_pt = run.font.size.pt
                    if size_pt > max_size:
                        max_size = size_pt
    except Exception:
        pass
    return max_size


def _smart_classify_textboxes(textboxes, slide):
    """
    Intelligently classify regular text boxes as title or body
    based on position, font size, and dimensions.
    
    Rules:
    1. The shape closest to the TOP with the LARGEST font = Title
    2. The shape with the LARGEST area (excluding title) = Body
    3. Everything else = decorative (skip)
    """
    if not textboxes:
        return None, None
    
    # Filter out tiny decorative elements (less than 1.5 inches wide)
    candidates = []
    for s in textboxes:
        if s.width >= Inches(1.5) and s.height >= Inches(0.3):
            candidates.append(s)
    
    if not candidates:
        return None, None
    
    # Score each candidate for "title-ness":
    # Higher score = more likely to be the title
    # Factors: top position (higher = more title-like), big font, short text
    best_title = None
    best_title_score = -1
    
    for s in candidates:
        score = 0
        font_size = _get_font_size_from_shape(s)
        text_len = len(s.text_frame.text.strip())
        
        # Big font = likely title (font size contributes heavily)
        score += font_size * 3
        
        # Top of slide = likely title (invert top position so higher = better)
        # Slide height is typically ~6858000 EMU (7.5 inches)
        top_normalized = max(0, 6858000 - s.top) / 6858000  # 1.0 at top, 0.0 at bottom
        score += top_normalized * 100
        
        # Short text = likely title (penalize long text)
        if text_len < 80:
            score += 50
        elif text_len > 200:
            score -= 50
        
        if score > best_title_score:
            best_title_score = score
            best_title = s
    
    # Body = the LARGEST remaining text box (by area)
    best_body = None
    best_body_area = 0
    
    for s in candidates:
        if s is best_title:
            continue
        area = s.width * s.height
        if area > best_body_area:
            best_body_area = area
            best_body = s
    
    # If we only found one candidate, decide: is it a title or body?
    if best_body is None and best_title is not None:
        font_size = _get_font_size_from_shape(best_title)
        if font_size < 20:  # Small font = probably body, not title
            best_body = best_title
            best_title = None
    
    return best_title, best_body


# ── Format-preserving text fill ───────────────────────────────────────────────

def _capture_font(run):
    """Capture font properties from a Run."""
    f = run.font
    color_val = None
    color_type = None
    try:
        if getattr(f, 'color', None) and getattr(f.color, 'type', None):
            # python-pptx raises an error if you access .rgb on a scheme color
            if str(f.color.type).endswith('RGB') or f.color.type == 1: # MSO_COLOR_TYPE.RGB
                color_val = f.color.rgb
                color_type = "rgb"
            elif str(f.color.type).endswith('SCHEME') or f.color.type == 2: # MSO_COLOR_TYPE.SCHEME
                color_val = getattr(f.color, 'theme_color', None)
                color_type = "theme"
    except Exception:
        pass

    return {
        "size": getattr(f, 'size', None),
        "bold": getattr(f, 'bold', None),
        "italic": getattr(f, 'italic', None),
        "name": getattr(f, 'name', None),
        "color_val": color_val,
        "color_type": color_type,
    }


def _apply_font(run, props):
    """Apply captured font properties and force Arial."""
    try:
        run.font.name = "Arial"
        if props.get("size"):    run.font.size = props["size"]
        if props.get("bold") is not None: run.font.bold = props["bold"]
        if props.get("italic") is not None: run.font.italic = props["italic"]
        
        c_val = props.get("color_val")
        c_type = props.get("color_type")
        if c_val is not None:
            if c_type == "rgb":
                run.font.color.rgb = c_val
            elif c_type == "theme":
                run.font.color.theme_color = c_val
    except Exception:
        pass


def _fill_textframe(tf, content, is_title=False, compact=False):
    """Replace text in a TextFrame, preserving first-run formatting.
    compact=True: used on picture-layout slides, reduces font/word limits."""
    # capture formatting from first existing run
    font_props = {}
    alignment = None
    if tf.paragraphs:
        p0 = tf.paragraphs[0]
        alignment = p0.alignment
        if p0.runs:
            font_props = _capture_font(p0.runs[0])

    tf.clear()
    tf.word_wrap = True

    items = content if isinstance(content, list) else [content]
    
    # Limit body text to prevent overflow
    if not is_title:
        max_points = 2 if compact else 5
        if len(items) > max_points:
            items = items[:max_points]

    # Calculate total text length for dynamic font sizing
    total_text_len = sum(len(str(item)) for item in items)
    
    if is_title:
        if total_text_len < 40: default_size = Pt(40)
        elif total_text_len < 80: default_size = Pt(36)
        elif total_text_len < 120: default_size = Pt(32)
        else: default_size = Pt(28)
        max_size = Pt(44)
        max_words = 30
    else:
        if total_text_len < 100: default_size = Pt(28)
        elif total_text_len < 250: default_size = Pt(24)
        elif total_text_len < 400: default_size = Pt(20)
        else: default_size = Pt(18)
        max_size = Pt(32)
        max_words = 120

    for i, item in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        text = str(item).strip()

        if not is_title:
            # Trim very long points
            words = text.split()
            if len(words) > max_words:
                text = ' '.join(words[:max_words]) + '...'

            # Add bullet if not already present
            if text and not text[0] in ('•', '–', '-', '✓', '✔', '►', '▸', '●'):
                # Don't add bullet to numbered items like "1." "2."
                if not (len(text) > 1 and text[0].isdigit() and text[1] in '.):'):
                    text = '•  ' + text

            # Paragraph spacing for body items
            para.space_after = Pt(8)
            para.space_before = Pt(2)

        para.level = 0
        if alignment is not None:
            para.alignment = alignment

        import re
        parts = re.split(r'(\*\*.*?\*\*)', text)
        for part in parts:
            if not part: continue
            run = para.add_run()
            is_bold = part.startswith('**') and part.endswith('**')
            run.text = part[2:-2] if is_bold else part
            
            _apply_font(run, font_props)
            # Override font size based on dynamic calculation, ignoring original template size 
            # if it's too big or small, but if original is reasonable, we can keep it.
            # Actually, to guarantee readability as requested: force the dynamic default_size!
            run.font.size = default_size

            if is_bold:
                run.font.bold = True
                try:
                    from pptx.enum.dml import MSO_THEME_COLOR
                    run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
                except: pass

    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass


# ── Deep slide duplication (XML level) ────────────────────────────────────────

def _duplicate_slide(prs, source_slide):
    """
    Duplicate a slide preserving visual design elements.
    
    Uses add_slide(layout) for writable placeholders, then copies
    non-placeholder shapes (decorative elements, images, backgrounds)
    from the source slide.
    """
    new_slide = prs.slides.add_slide(source_slide.slide_layout)

    # Copy background from source if it exists
    src_bg = source_slide._element.find(
        '{http://schemas.openxmlformats.org/presentationml/2006/main}bg'
    )
    if src_bg is not None:
        new_bg = new_slide._element.find(
            '{http://schemas.openxmlformats.org/presentationml/2006/main}bg'
        )
        if new_bg is not None:
            new_slide._element.remove(new_bg)
        # Insert bg before spTree
        spTree = new_slide._element.find(
            './/{http://schemas.openxmlformats.org/presentationml/2006/main}cSld/'
            '{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing}spTree'
        )
        cSld = new_slide._element.find(
            '{http://schemas.openxmlformats.org/presentationml/2006/main}cSld'
        )
        if cSld is not None:
            cSld.insert(0, copy.deepcopy(src_bg))

    # Copy non-placeholder shapes (decorative elements, images, etc.)
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    
    src_spTree = source_slide._element.find(f'{{{ns_p}}}cSld/{{{ns_p}}}spTree')
    new_spTree = new_slide._element.find(f'{{{ns_p}}}cSld/{{{ns_p}}}spTree')
    
    if src_spTree is not None and new_spTree is not None:
        for child in src_spTree:
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            # Copy pics, connectors, group shapes (not sp placeholders)
            if tag in ('pic', 'cxnSp', 'grpSp', 'graphicFrame'):
                new_spTree.append(copy.deepcopy(child))
            elif tag == 'sp':
                # Only copy non-placeholder shapes (decorative text boxes, etc.)
                nvSpPr = child.find(f'.//{{{ns_p}}}nvSpPr/{{{ns_p}}}nvPr')
                if nvSpPr is None:
                    nvSpPr = child.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}nvPr')
                has_ph = False
                if nvSpPr is not None:
                    ph_elem = nvSpPr.find('{http://schemas.openxmlformats.org/presentationml/2006/main}ph')
                    has_ph = ph_elem is not None
                if not has_ph:
                    new_spTree.append(copy.deepcopy(child))

    # Copy image relationships from source
    for rel in source_slide.part.rels.values():
        if rel.is_external:
            continue
        try:
            if 'image' in rel.reltype.lower():
                new_slide.part.relate_to(rel.target_part, rel.reltype)
        except Exception:
            pass

    return new_slide


# ── Image helper ──────────────────────────────────────────────────────────────

def _find_picture_placeholder(slide):
    """Find picture placeholder in slide (the frame/рамка for images).
    Only detects actual PICTURE placeholders (type=18), not OBJECT/BODY."""
    for shape in slide.shapes:
        if shape.is_placeholder:
            pf = shape.placeholder_format
            # PP_PLACEHOLDER.PICTURE = 18 — this is the real picture frame
            if pf.type == PP_PLACEHOLDER.PICTURE:
                return shape
            # Also BITMAP, MEDIA_CLIP
            if pf.type in (PP_PLACEHOLDER.BITMAP, PP_PLACEHOLDER.MEDIA_CLIP):
                return shape
    return None


def _add_image(slide, img_bytes, prs, position="right"):
    """Insert image into picture placeholder if available, else overlay at position."""
    try:
        # ── PHASE 1: Try to find a picture placeholder (рамка) ──
        pic_ph = _find_picture_placeholder(slide)
        if pic_ph is not None:
            try:
                stream = io.BytesIO(img_bytes)
                pic_ph.insert_picture(stream)
                return  # Successfully inserted into placeholder frame
            except Exception:
                pass  # Placeholder doesn't support insert_picture, fall through

        # ── PHASE 2: Fallback — overlay image at top-right, compact size ──
        stream = io.BytesIO(img_bytes)
        w, h = prs.slide_width, prs.slide_height
        if position == "fullscreen":
            slide.shapes.add_picture(stream, 0, 0, w, h)
        elif position == "center":
            iw, ih = int(w * 0.30), int(h * 0.30)
            slide.shapes.add_picture(stream, int((w - iw) / 2), int(h * 0.55), iw, ih)
        else:
            iw, ih = int(w * 0.28), int(h * 0.40)
            slide.shapes.add_picture(stream, int(w - iw - Inches(0.2)),
                                     int(Inches(0.3)), iw, ih)
    except Exception as e:
        print(f"Image insert error: {e}")



# ── Fill a single slide with AI data ──────────────────────────────────────────

def _fill_slide(slide, slide_info):
    """Fill title + body on any slide."""
    title_text = slide_info.get("title", "")
    points = slide_info.get("content", slide_info.get("points", []))
    shapes = _classify_shapes(slide)
    
    # Check if this slide has a picture placeholder → compact mode
    has_picture = _find_picture_placeholder(slide) is not None
    is_plan_slide = title_text.lower().strip() in ("reja", "plan", "план", "mundarija")
    if has_picture and len(points) > 2 and not is_plan_slide:
        # Trim to 2 shorter points for picture-layout slides (but not plan slides)
        points = points[:2]
    
    has_placeholders = bool(shapes["title"] or shapes["body"] or shapes["subtitle"])

    # ── Fill title placeholder(s) ──
    title_filled = False
    for s in shapes["title"]:
        _fill_textframe(s.text_frame, title_text, is_title=True)
        title_filled = True

    # ── Fill body placeholder(s) — ONLY real placeholders ──
    body_filled = False
    use_compact = has_picture and not is_plan_slide
    if shapes["body"]:
        # Distribute points across all available body placeholders
        num_bodies = len(shapes["body"])
        points_per_body = max(1, len(points) // num_bodies)
        
        for i, s in enumerate(shapes["body"]):
            start_idx = i * points_per_body
            # If it's the last body, give it all the remaining points
            end_idx = (i + 1) * points_per_body if i < num_bodies - 1 else len(points)
            
            chunk = points[start_idx:end_idx]
            if chunk:
                _fill_textframe(s.text_frame, chunk, compact=use_compact)
                body_filled = True
            else:
                # If we ran out of points, just clear the placeholder so it doesn't show "Текст слайда"
                s.text_frame.clear()

    # ── Fallback: subtitle placeholder ──
    if not body_filled and points:
        for s in shapes["subtitle"]:
            _fill_textframe(s.text_frame, points)
            body_filled = True
            break

    # ── SMART FALLBACK: Auto-detect from regular text boxes ──
    if not has_placeholders and shapes["other"]:
        smart_title, smart_body = _smart_classify_textboxes(shapes["other"], slide)
        
        if smart_title and not title_filled:
            _fill_textframe(smart_title.text_frame, title_text, is_title=True)
            title_filled = True
        
        if smart_body and not body_filled and points:
            _fill_textframe(smart_body.text_frame, points)
            body_filled = True
    
    # ── Last resort: single largest "other" textbox ──
    if not body_filled and points and shapes["other"]:
        candidates = []
        for s in shapes["other"]:
            if s.width < Inches(1.5) or s.height < Inches(0.5):
                continue
            candidates.append(s)
        
        if candidates:
            biggest = max(candidates, key=lambda s: s.width * s.height)
            _fill_textframe(biggest.text_frame, points)
            body_filled = True


# ── Layout fallback (for templates with only 1 slide) ─────────────────────────

def _best_content_layout(prs):
    preferred = ["Заголовок и объект", "Title and Content", "Sarlavha va ob'ekt"]
    for name in preferred:
        for lay in prs.slide_layouts:
            if name.lower() in lay.name.lower():
                return lay
    for lay in prs.slide_layouts:
        types = {ph.placeholder_format.type for ph in lay.placeholders}
        if types & TITLE_PH and types & BODY_PH:
            return lay
    return prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]


# ── Remove extra slides ──────────────────────────────────────────────────────

def _remove_slide(prs, slide_index):
    """Remove slide at given index."""
    sldIdLst = prs.slides._sldIdLst
    items = list(sldIdLst)
    if slide_index < len(items):
        el = items[slide_index]
        try:
            prs.part.drop_rel(el.rId)
        except Exception:
            pass
        sldIdLst.remove(el)


# ── Academic title slide builder ──────────────────────────────────────────────

# University logo cache directory
_LOGOS_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                          "webapp", "images", "logos")

def _find_university_logo(university: str) -> Optional[str]:
    """Find university logo locally or download it."""
    if not university:
        return None
    os.makedirs(_LOGOS_DIR, exist_ok=True)
    
    # Normalize name for matching
    uni_lower = university.lower().replace("'", "").replace("ʻ", "").replace("'", "")
    
    # 1. Check local cache
    for fname in os.listdir(_LOGOS_DIR):
        name_lower = fname.lower().replace("'", "").replace("ʻ", "")
        name_no_ext = os.path.splitext(name_lower)[0]
        # Match if filename is part of uni name or vice versa
        if name_no_ext in uni_lower or uni_lower in name_no_ext:
            return os.path.join(_LOGOS_DIR, fname)
        # Also match abbreviations
        for part in name_no_ext.split("_"):
            if len(part) >= 3 and part in uni_lower:
                return os.path.join(_LOGOS_DIR, fname)
    
    # 2. Try to download via web search
    return _download_university_logo(university)


# Known Uzbek university → Wikipedia article name mapping
_UNI_WIKI_MAP = {
    "tdiu": "Tashkent_State_University_of_Economics",
    "iqtisodiyot": "Tashkent_State_University_of_Economics",
    "tsue": "Tashkent_State_University_of_Economics",
    "tatu": "Tashkent_University_of_Information_Technologies",
    "axborot": "Tashkent_University_of_Information_Technologies",
    "milliy": "National_University_of_Uzbekistan",
    "ozmu": "National_University_of_Uzbekistan",
    "texnika": "Tashkent_State_Technical_University",
    "toshtdtu": "Tashkent_State_Technical_University",
    "politexnika": "Turin_Polytechnic_University_in_Tashkent",
    "turin": "Turin_Polytechnic_University_in_Tashkent",
    "yuridik": "Tashkent_State_University_of_Law",
    "tdyu": "Tashkent_State_University_of_Law",
    "pedagogika": "Tashkent_State_Pedagogical_University",
    "tdpu": "Tashkent_State_Pedagogical_University",
    "samarqand": "Samarkand_State_University",
    "samdu": "Samarkand_State_University",
    "buxoro": "Bukhara_State_University",
    "buxdu": "Bukhara_State_University",
    "andijon": "Andijan_State_University",
    "fargona": "Fergana_State_University",
    "namangan": "Namangan_State_University",
    "navoiy": "Navoi_State_Mining_and_Technologies_University",
    "qoraqalpog": "Karakalpak_State_University",
    "nukus": "Karakalpak_State_University",
    "tashkent": "National_University_of_Uzbekistan",
    "transport": "Tashkent_State_Transport_University",
    "moliya": "Tashkent_Institute_of_Finance",
    "tibbiyot": "Tashkent_Medical_Academy",
    "tma": "Tashkent_Medical_Academy",
    "agrar": "Tashkent_State_Agrarian_University",
    "sharq": "University_of_World_Economy_and_Diplomacy",
    "jdpu": "Jizzakh_State_Pedagogical_University",
    "jizzax": "Jizzakh_State_Pedagogical_University",
    "webster": "Webster_University",
    "westminster": "Westminster_International_University_in_Tashkent",
    "inha": "Inha_University_in_Tashkent",
}


def _download_university_logo(university: str) -> Optional[str]:
    """Download university logo/emblem from Wikipedia and cache locally."""
    import urllib.request
    import json as _json
    import re
    
    safe_name = re.sub(r'[^\w\s-]', '', university.lower()).strip().replace(' ', '_')[:50]
    cache_path = os.path.join(_LOGOS_DIR, f"{safe_name}.png")
    if os.path.exists(cache_path):
        return cache_path
    
    uni_lower = university.lower().replace("'", "").replace("ʻ", "").replace("'", "")
    wiki_titles = []
    for keyword, wiki_title in _UNI_WIKI_MAP.items():
        if keyword in uni_lower:
            if wiki_title not in wiki_titles:
                wiki_titles.append(wiki_title)
    wiki_titles.append(university.replace(" ", "_"))
    
    ua = "StudentBot/1.0 (https://t.me/SLIDEMASTERAI1_BOT; education) Python/3"
    LOGO_WORDS = {"logo", "seal", "emblem", "coat", "gerb", "arms", "crest"}
    
    try:
        for wiki_q in wiki_titles:
            try:
                # Step 1: Get list of images on the Wikipedia page
                api = (
                    f"https://en.wikipedia.org/w/api.php?action=query"
                    f"&titles={wiki_q}&prop=images&format=json"
                )
                req = urllib.request.Request(api, headers={"User-Agent": ua})
                resp = urllib.request.urlopen(req, timeout=10)
                data = _json.loads(resp.read())
                
                logo_files = []
                emblem_files = []
                for page in data.get("query", {}).get("pages", {}).values():
                    for img in page.get("images", []):
                        title = img.get("title", "")
                        tl = title.lower()
                        # Skip generic icons (commons, symbol, icon, flag)
                        if any(skip in tl for skip in ("commons", "symbol", "icon", "flag", "oojs")):
                            continue
                        if "logo" in tl:
                            logo_files.append(title)  # High priority
                        elif any(w in tl for w in LOGO_WORDS):
                            # Skip city/country emblems
                            if not any(skip in tl for skip in ("tashkent", "uzbekistan", "samarkand")):
                                emblem_files.append(title)
                
                # Prefer 'logo' files, fallback to emblem
                all_candidates = logo_files + emblem_files
                if not all_candidates:
                    continue
                
                # Step 2: Get actual URL for the logo file
                for logo_title in all_candidates:
                    try:
                        safe_title = logo_title.replace(" ", "_")
                        info_url = (
                            f"https://en.wikipedia.org/w/api.php?action=query"
                            f"&titles={safe_title}&prop=imageinfo&iiprop=url"
                            f"&format=json"
                        )
                        req2 = urllib.request.Request(info_url, headers={"User-Agent": ua})
                        resp2 = urllib.request.urlopen(req2, timeout=10)
                        data2 = _json.loads(resp2.read())
                        
                        for page2 in data2.get("query", {}).get("pages", {}).values():
                            for ii in page2.get("imageinfo", []):
                                img_url = ii.get("url", "")
                                if not img_url:
                                    continue
                                # Download the image
                                req3 = urllib.request.Request(img_url, headers={"User-Agent": ua})
                                img_resp = urllib.request.urlopen(req3, timeout=15)
                                img_bytes = img_resp.read()
                                if len(img_bytes) > 500:
                                    with open(cache_path, 'wb') as f:
                                        f.write(img_bytes)
                                    print(f"✅ University logo downloaded: {university} → {logo_title}")
                                    return cache_path
                    except Exception:
                        continue
            except Exception:
                continue
    except Exception as e:
        print(f"Logo download error: {e}")
    
    # Fallback: generate emblem with Gemini AI
    return _generate_emblem_gemini(university, cache_path)


def _generate_emblem_gemini(university: str, cache_path: str) -> Optional[str]:
    """Generate a university emblem using Gemini AI as fallback."""
    import urllib.request
    import json as _json
    import base64
    
    gemini_key = os.environ.get("GEMINI_API_KEY", "")
    if not gemini_key:
        print(f"⚠️ No GEMINI_API_KEY, cannot generate emblem for: {university}")
        return None
    
    prompt = (
        f"Create a formal university emblem/seal for '{university}' (Uzbekistan). "
        "Style: classic academic heraldic emblem with a shield, laurel wreath, "
        "book symbol, and the university name written on a ribbon. "
        "Colors: dark blue and gold. Clean vector-like style on white background. "
        "NO photograph, NO realistic image — only a formal emblem/coat of arms design. "
        "NO extra text outside the emblem."
    )
    
    payload = _json.dumps({
        "contents": [{"parts": [{"text": prompt}]}],
        "generationConfig": {"responseModalities": ["IMAGE", "TEXT"]}
    }).encode()
    
    url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash-image:generateContent?key={gemini_key}"
    
    try:
        req = urllib.request.Request(
            url, data=payload,
            headers={"Content-Type": "application/json"}
        )
        resp = urllib.request.urlopen(req, timeout=60)
        data = _json.loads(resp.read())
        parts = data.get("candidates", [{}])[0].get("content", {}).get("parts", [])
        for p in parts:
            if "inlineData" in p:
                img_bytes = base64.b64decode(p["inlineData"]["data"])
                if len(img_bytes) > 500:
                    with open(cache_path, 'wb') as f:
                        f.write(img_bytes)
                    print(f"✅ University emblem generated (Gemini): {university}")
                    return cache_path
    except Exception as e:
        print(f"Gemini emblem generation error: {e}")
    
    print(f"⚠️ University logo/emblem not found for: {university}")
    return None


def _add_textbox(slide, left, top, width, height, text, font_size=14,
                 bold=False, alignment=PP_ALIGN.CENTER, color=None):
    """Helper: add a text box to a slide with given formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    for run in p.runs:
        run.font.name = "Arial"
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
    return txBox


def build_academic_title_slide(slide, topic: str, author: str = "",
                                subject: str = "", university: str = "",
                                doc_type_label: str = "", reviewer: str = "", prs=None):
    """
    Overlay academic title info on top of an existing title slide.
    Layout matches Uzbek university standard:
    
    O'ZBEKISTON RESPUBLIKASI OLIY TA'LIM...
    [University name]
    [University logo]
    Fan: ... (doc type)
    Mavzu: ...
    Bajardi: ... (right-aligned)
    """
    slide_w = prs.slide_width if prs else Inches(13.333)
    slide_h = prs.slide_height if prs else Inches(7.5)
    
    # Remove all text shapes from the slide (we'll add our own)
    for shape in list(slide.shapes):
        if shape.has_text_frame:
            try:
                sp = shape._element
                sp.getparent().remove(sp)
            except Exception:
                pass

    margin = Inches(0.5)
    content_w = slide_w - margin * 2
    y_pos = Inches(0.3)

    # 1. Ministry header
    header_text = (
        "O'ZBEKISTON RESPUBLIKASI OLIY TA'LIM, FAN VA\n"
        "INNOVATSIYALAR VAZIRLIGI"
    )
    _add_textbox(slide, margin, y_pos, content_w, Inches(0.7),
                 header_text, font_size=14, bold=True,
                 color=RGBColor(0, 0, 0))
    y_pos += Inches(0.7)

    # 2. University name
    if university:
        uni_upper = university.upper()
    else:
        uni_upper = "UNIVERSITET NOMI"
    _add_textbox(slide, margin, y_pos, content_w, Inches(0.5),
                 uni_upper, font_size=16, bold=True,
                 color=RGBColor(0, 0, 0))
    y_pos += Inches(0.6)

    # 3. University logo (if available)
    logo_path = _find_university_logo(university)
    if logo_path and prs:
        try:
            logo_size = Inches(1.8)
            logo_left = (slide_w - logo_size) // 2
            slide.shapes.add_picture(
                logo_path, logo_left, y_pos, logo_size, logo_size
            )
            y_pos += Inches(2.0)
        except Exception:
            y_pos += Inches(0.3)
    else:
        y_pos += Inches(0.3)

    # 4. Subject + doc type
    if subject or doc_type_label:
        fan_text = ""
        if subject:
            fan_text = f"Fan: {subject}"
        if doc_type_label:
            if fan_text:
                fan_text += f"\n({doc_type_label})"
            else:
                fan_text = doc_type_label
        _add_textbox(slide, margin, y_pos, content_w, Inches(0.7),
                     fan_text, font_size=16, bold=True,
                     color=RGBColor(0, 0, 0))
        y_pos += Inches(0.8)

    # 5. Topic (Mavzu)
    mavzu_text = f"MAVZU:\n{topic.upper()}"
    _add_textbox(slide, margin, y_pos, content_w, Inches(1.5),
                 mavzu_text, font_size=32, bold=True,
                 color=RGBColor(0, 0, 0))
    y_pos += Inches(1.5)

    # 6. Author and Reviewer (left-aligned, bottom area)
    if author or reviewer:
        author_text = ""
        if author: author_text += f"Bajardi: {author}\n"
        if reviewer: author_text += f"Qabul qildi: {reviewer}"
        author_left = Inches(1.0)
        author_top = slide_h - Inches(1.8)
        _add_textbox(slide, author_left, author_top, Inches(4.0), Inches(1.0),
                     author_text.strip(), font_size=16, bold=False,
                     alignment=PP_ALIGN.LEFT, color=RGBColor(0, 0, 0))


# ══════════════════════════════════════════════════════════════════════════════
# MAIN ENTRY POINT
# ══════════════════════════════════════════════════════════════════════════════

def generate_pptx(
    slides_data: List[dict],
    template_path: str,
    topic: str,
    author: str = "Foydalanuvchi",
    slide_images: dict = None,
    maket_mode: bool = False,
    subject: str = "",
    university: str = "",
    doc_type_label: str = "",
    reviewer: str = "",
) -> bytes:
    """
    Smart PPTX generator.

    • If the template has multiple pre-designed slides → fills them in-place,
      duplicating the last content slide when more are needed.
    • If the template has only a title slide → falls back to adding slides
      from the best matching layout (old behaviour).
    """

    # resolve path
    if template_path and not os.path.isabs(template_path):
        base = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        full = os.path.join(base, template_path)
        if os.path.exists(full):
            template_path = full

    prs = (Presentation(template_path)
           if template_path and os.path.exists(template_path)
           else Presentation())
    slide_images = slide_images or {}

    tpl_count = len(prs.slides)          # total slides in template
    tpl_content = max(0, tpl_count - 1)  # content slides (minus title)
    
    # In maket_mode, we only fill as many slides as exist in the template
    if maket_mode:
        need_content = min(len(slides_data), tpl_content)
    else:
        need_content = len(slides_data)      # how many AI content slides

    # ── 1. Title slide ────────────────────────────────────────────────────
    title_slide = prs.slides[0] if prs.slides else None
    if title_slide:
        # Use academic title slide if subject or university is provided
        if subject or university:
            build_academic_title_slide(
                title_slide, topic=topic, author=author,
                subject=subject, university=university,
                doc_type_label=doc_type_label, reviewer=reviewer, prs=prs
            )
        else:
            # Legacy: simple title + author
            shapes = _classify_shapes(title_slide)
            title_filled = False
            
            for s in shapes["title"]:
                _fill_textframe(s.text_frame, topic, is_title=True)
                title_filled = True
            for s in shapes["subtitle"]:
                _fill_textframe(s.text_frame, f"Bajardi: {author}")
            
            if not title_filled and shapes["other"]:
                smart_title, smart_body = _smart_classify_textboxes(shapes["other"], title_slide)
                if smart_title:
                    _fill_textframe(smart_title.text_frame, topic, is_title=True)
                if smart_body:
                    _fill_textframe(smart_body.text_frame, f"Bajardi: {author}")
        
        if 0 in slide_images:
            _add_image(title_slide, slide_images[0], prs, "fullscreen")

    # ── 2. Content slides ─────────────────────────────────────────────────
    if tpl_content == 0:
        # --- Template has only title → create from layout (old way) ---
        layout = _best_content_layout(prs)
        for i, info in enumerate(slides_data):
            ns = prs.slides.add_slide(layout)
            _fill_slide(ns, info)
            if (i + 1) in slide_images:
                _add_image(ns, slide_images[i + 1], prs)
    else:
        # --- Template has pre-designed content slides → use them ---

        # Classify template content slides: which have picture placeholders?
        tpl_slides = list(prs.slides)
        text_only_indices = []
        picture_indices = []

        for si in range(1, tpl_count):
            slide = tpl_slides[si]
            if _find_picture_placeholder(slide) is not None:
                picture_indices.append(si)
            else:
                text_only_indices.append(si)

        # Step A: Duplicate all needed content slides FIRST (sources still alive)
        if not maket_mode:
            text_cycle = 0
            for i in range(need_content):
                slide_idx = i + 1  # slide index in final presentation
                has_image = slide_idx in slide_images
                
                if has_image and picture_indices:
                    # Use the first picture layout (or could cycle if multiple)
                    _duplicate_slide(prs, tpl_slides[picture_indices[0]])
                else:
                    if text_only_indices:
                        # Alternate between available text layouts for variety
                        idx_to_use = text_only_indices[text_cycle % len(text_only_indices)]
                        _duplicate_slide(prs, tpl_slides[idx_to_use])
                        text_cycle += 1
                    elif picture_indices:
                        _duplicate_slide(prs, tpl_slides[picture_indices[0]])
                    else:
                        _duplicate_slide(prs, tpl_slides[tpl_count - 1])

        # Step B: Remove the ORIGINAL template content slides (they come before duplicates)
        # Original content slides are at indices 1..tpl_count-1
        # Remove them from last to first to avoid index shifting
        if not maket_mode:
            for _ in range(tpl_content):
                _remove_slide(prs, 1)  # always remove index 1 (first content slide)

        # Step C: Fill all content slides
        fresh = list(prs.slides)
        for i in range(need_content):
            idx = i + 1  # skip title
            if idx < len(fresh):
                _fill_slide(fresh[idx], slides_data[i])
                if (i + 1) in slide_images:
                    _add_image(fresh[idx], slide_images[i + 1], prs)

    # ── 3. Unique Design, Animations & Save ──────────────────────────────────────
    try:
        from services.premium_design import apply_premium_design, apply_premium_transitions
        apply_premium_design(prs, topic)
        apply_premium_transitions(prs)
    except Exception as e:
        import traceback
        print("Premium Design Error:", e)
        traceback.print_exc()
        pass  # non-critical
    # Clean up empty picture placeholders (to prevent dashed boxes)
    for slide in prs.slides:
        for shape in list(slide.shapes):
            if shape.is_placeholder:
                try:
                    if shape.placeholder_format.type == 18 and not shape.has_text_frame:
                        if not hasattr(shape, "image"):
                            sp = shape._element
                            sp.getparent().remove(sp)
                except Exception:
                    pass

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()

def apply_modern_animations(prs):
    from pptx.oxml import parse_xml
    import random
    transitions = [
        '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:fade/></p:transition>',
        '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:push dir="u"/></p:transition>',
        '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:cover dir="l"/></p:transition>'
    ]
    for slide in prs.slides:
        try:
            t_xml = random.choice(transitions)
            slide._element.insert(1, parse_xml(t_xml))
        except:
            pass

def apply_unique_ai_design(prs):
    """
    Applies a dynamically generated unique design to the presentation,
    moving away from static templates to create a truly unique and modern look.
    """
    import random
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    
    # Generate a random professional color palette
    palettes = [
        {"bg": RGBColor(20, 25, 35), "accent": RGBColor(255, 94, 58), "text": RGBColor(245, 245, 250)}, # Dark Modern
        {"bg": RGBColor(245, 247, 250), "accent": RGBColor(0, 102, 204), "text": RGBColor(30, 35, 40)},  # Light Corporate
        {"bg": RGBColor(25, 40, 55), "accent": RGBColor(46, 204, 113), "text": RGBColor(255, 255, 255)}, # Deep Forest
        {"bg": RGBColor(250, 245, 240), "accent": RGBColor(211, 84, 0), "text": RGBColor(40, 30, 20)},   # Warm Minimalist
    ]
    palette = random.choice(palettes)
    is_dark = (palette["bg"].r + palette["bg"].g + palette["bg"].b) < 300
    
    for idx, slide in enumerate(prs.slides):
        # 1. Apply background
        try:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = palette["bg"]
        except:
            pass
            
        # 2. Add decorative geometric shapes (Unique touch)
        # Left accent bar
        try:
            accent_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.2), prs.slide_height
            )
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = palette["accent"]
            accent_bar.line.fill.background()
            
            # Slide index specific decorations
            if idx == 0:
                # Title slide decoration
                circle = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL, prs.slide_width - Inches(2), -Inches(1), Inches(3), Inches(3)
                )
                circle.fill.solid()
                circle.fill.fore_color.rgb = palette["accent"]
                circle.line.fill.background()
            elif idx == len(prs.slides) - 1:
                # Final slide decoration (bottom thick bar)
                bottom_bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.5), prs.slide_width, Inches(0.5)
                )
                bottom_bar.fill.solid()
                bottom_bar.fill.fore_color.rgb = palette["accent"]
                bottom_bar.line.fill.background()
        except:
            pass

        # 3. Recolor text for contrast
        for shape in slide.shapes:
            if not shape.has_text_frame: continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.color.type == 0: # If no explicit color
                        run.font.color.rgb = palette["text"]
                    elif is_dark and hasattr(run.font.color, 'rgb') and run.font.color.rgb == RGBColor(0,0,0):
                        # Fix black text on dark background
                        run.font.color.rgb = palette["text"]