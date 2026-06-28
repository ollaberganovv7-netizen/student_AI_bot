from __future__ import annotations
"""
Akademik PPTX Service
---------------------
Specialized generator for 'akademik' category templates.
Unlike the universal generator, this service does NOT delete or clone slides.
It preserves the entire template design and only replaces [TAG] placeholders
with AI-generated content.

Supported tags:
  [TOPIC], [SUBJECT_NAME], [COMPLETED_BY], [TRAINING_SESSION],
  [EDUCATIONAL_GOALS], [STUDY_QUESTIONS],
  [QUESTION_1], [QUESTION_1_CONTENT], [QUESTION_1_CONCLUSION],
  [QUESTION_2], [QUESTION_2_CONTENT], [QUESTION_2_CONCLUSION],
  [QUESTION_3], [QUESTION_3_CONTENT], [QUESTION_3_CONTENT_IMAGE],
  [QUESTION_1_CONTENT_IMAGE], [QUESTION_2_CONTENT_IMAGE], [QUESTION_3_CONTENT_IMAGE],
  [GENERAL_CONCLUSION], [REFERENCES_LIST]
"""

import copy
import io
import re
from pptx import Presentation
from pptx.util import Emu


def generate_akademik_pptx(
    template_path: str,
    tag_data: dict,
    image_data: dict | None = None,
) -> bytes:
    """
    Open an akademik template and replace all [TAG] placeholders with content.

    Args:
        template_path: path to the .pptx template file
        tag_data: dict mapping tag names (without brackets) to replacement text
                  e.g. {"TOPIC": "Fizika asoslari", "QUESTION_1_CONTENT": "..."}
        image_data: dict mapping IMAGE tag names to image bytes
                    e.g. {"QUESTION_1_CONTENT_IMAGE": b"...png bytes..."}
    Returns:
        bytes of the generated .pptx file
    """
    prs = Presentation(template_path)
    image_data = image_data or {}

    for slide in prs.slides:
        _process_slide(slide, tag_data, image_data, prs)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _process_slide(slide, tag_data: dict, image_data: dict, prs):
    """Process a single slide: replace text tags and image tags."""
    # Collect shapes to remove after iteration (image placeholders)
    shapes_to_remove = []

    for shape in list(slide.shapes):
        if not shape.has_text_frame:
            continue

        full_text = _get_shape_full_text(shape)

        # Check if this shape contains an IMAGE tag
        image_tag_match = re.search(r'\[([A-Z0-9_]*IMAGE[A-Z0-9_]*)\]', full_text)
        if image_tag_match:
            tag_name = image_tag_match.group(1)
            if tag_name in image_data and image_data[tag_name]:
                # Replace this shape with an image at the same position/size
                _replace_shape_with_image(slide, shape, image_data[tag_name])
                shapes_to_remove.append(shape)
            else:
                # No image data for this tag — clear the tag text but keep shape
                _replace_tags_in_shape(shape, tag_data)
            continue

        # Regular text tag replacement
        _replace_tags_in_shape(shape, tag_data)

    # Remove shapes that were replaced by images
    for shape in shapes_to_remove:
        try:
            sp = shape._element
            sp.getparent().remove(sp)
        except Exception as e:
            print(f"Warning: could not remove placeholder shape: {e}")


def _get_shape_full_text(shape) -> str:
    """Get the full concatenated text of all paragraphs and runs in a shape."""
    text_parts = []
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            text_parts.append(run.text)
    return ''.join(text_parts)


def _replace_tags_in_shape(shape, tag_data: dict):
    """
    Replace [TAG] placeholders in a shape's text with values from tag_data.
    Preserves original formatting (font, color, size) of the first run.
    """
    for para in shape.text_frame.paragraphs:
        # Build full paragraph text from runs
        full_para_text = ''.join(run.text for run in para.runs)

        # Find all tags in this paragraph
        tags_in_para = re.findall(r'\[([A-Z0-9_]+)\]', full_para_text)

        if not tags_in_para:
            continue

        # Replace tags with actual content
        new_text = full_para_text
        for tag in tags_in_para:
            replacement = tag_data.get(tag, "")
            new_text = new_text.replace(f'[{tag}]', replacement)

        # Clean up whitespace
        new_text = new_text.strip()

        # Preserve first run's formatting
        if para.runs:
            first_run = para.runs[0]
            font_props = _capture_font(first_run)

            # Clear all runs except first
            for run in para.runs[1:]:
                run.text = ""
            
            # Set the text on the first run
            first_run.text = new_text

            # Restore original formatting
            _apply_font(first_run, font_props)
        else:
            # No runs — just set text directly
            para.text = new_text


def _capture_font(run) -> dict:
    """Capture font properties from a run."""
    font = run.font
    props = {}
    try:
        if font.size is not None:
            props['size'] = font.size
    except:
        pass
    try:
        if font.bold is not None:
            props['bold'] = font.bold
    except:
        pass
    try:
        if font.italic is not None:
            props['italic'] = font.italic
    except:
        pass
    try:
        if font.color and font.color.rgb is not None:
            props['color_rgb'] = font.color.rgb
    except:
        pass
    try:
        if font.name is not None:
            props['name'] = font.name
    except:
        pass
    return props


def _apply_font(run, props: dict):
    """Apply captured font properties to a run."""
    font = run.font
    if 'size' in props:
        font.size = props['size']
    if 'bold' in props:
        font.bold = props['bold']
    if 'italic' in props:
        font.italic = props['italic']
    if 'color_rgb' in props:
        font.color.rgb = props['color_rgb']
    if 'name' in props:
        font.name = props['name']


def _replace_shape_with_image(slide, shape, image_bytes: bytes):
    """
    Insert an image at the exact position and size of the given shape.
    The original shape should be removed by the caller after this.
    """
    left = shape.left
    top = shape.top
    width = shape.width
    height = shape.height

    image_stream = io.BytesIO(image_bytes)

    try:
        slide.shapes.add_picture(image_stream, left, top, width, height)
    except Exception as e:
        print(f"Error inserting image: {e}")
