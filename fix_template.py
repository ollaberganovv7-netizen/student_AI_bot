from __future__ import annotations
from pptx import Presentation
import re

prs = Presentation('templates/modern/test_example.pptx')
TAG = re.compile(r'\{\{(\w+)\}\}')

CYRILLIC_TO_LATIN = {
    '\u0410': 'A', '\u0412': 'B', '\u0421': 'C', '\u0415': 'E',
    '\u041D': 'H', '\u041A': 'K', '\u041C': 'M', '\u041E': 'O',
    '\u0420': 'P', '\u0422': 'T', '\u0425': 'X',
    '\u0430': 'a', '\u0432': 'b', '\u0441': 'c', '\u0435': 'e',
    '\u043D': 'h', '\u043A': 'k', '\u043C': 'm', '\u043E': 'o',
    '\u0440': 'p', '\u0442': 't', '\u0445': 'x',
}

def fix_text(text):
    """Replace Cyrillic lookalikes with Latin and fix typos"""
    fixed = ''
    for ch in text:
        fixed += CYRILLIC_TO_LATIN.get(ch, ch)
    fixed = fixed.replace('Subtotle', 'Subtitle')
    return fixed

fixes = 0
for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                # Get full paragraph text
                full_text = para.text
                if '{{' not in full_text:
                    continue
                fixed_full = fix_text(full_text)
                if fixed_full != full_text:
                    print(f'Slide {i+1}: {full_text} -> {fixed_full}')
                    fixes += 1
                    # Put all text into first run, clear others
                    if para.runs:
                        para.runs[0].text = fixed_full
                        for run in para.runs[1:]:
                            run.text = ''

prs.save('templates/modern/test_example.pptx')
print(f'Fixed {fixes} tags total')
