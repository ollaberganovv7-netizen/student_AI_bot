from pptx import Presentation
from pptx.oxml import parse_xml
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
_NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
xml_str = f'<p:transition xmlns:p="{_NS_P}" spd="slow"><p:fade/></p:transition>'
slide._element.insert(1, parse_xml(xml_str))
prs.save('test_anim.pptx')
print('Success')
