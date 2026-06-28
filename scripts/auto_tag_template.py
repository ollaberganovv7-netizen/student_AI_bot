from __future__ import annotations
"""
Auto-tag a PPTX template using GPT-4 Vision.
Converts each slide to an image, sends to GPT-4o,
and inserts {{title}}, {{body}}, {{subtitle}}, {{author}}, {{image}} tags.

Usage:
    python scripts/auto_tag_template.py templates/Business/template.pptx
"""

import sys
import os
import json
import subprocess
import tempfile
import base64
import asyncio
import logging
from pathlib import Path

# Add project root to path
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from openai import AsyncOpenAI
from config import OPENAI_API_KEY
from pptx import Presentation
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)
client = AsyncOpenAI(api_key=OPENAI_API_KEY)


def pptx_to_images(pptx_path: str, output_dir: str) -> list[str]:
    """Convert PPTX slides to PNG images using LibreOffice."""
    # First try LibreOffice
    try:
        subprocess.run([
            "soffice", "--headless", "--convert-to", "png",
            "--outdir", output_dir, pptx_path
        ], check=True, capture_output=True, timeout=60)
    except (subprocess.CalledProcessError, FileNotFoundError):
        # Try macOS LibreOffice path
        lo_path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
        if os.path.exists(lo_path):
            subprocess.run([
                lo_path, "--headless", "--convert-to", "png",
                "--outdir", output_dir, pptx_path
            ], check=True, capture_output=True, timeout=60)
        else:
            raise RuntimeError(
                "LibreOffice not found. Install it: brew install --cask libreoffice"
            )

    # LibreOffice creates one image per slide (or one combined)
    # Check for individual slide images
    images = sorted(Path(output_dir).glob("*.png"))
    if not images:
        raise RuntimeError("No images generated from PPTX")
    return [str(img) for img in images]


def pptx_to_pdf_to_images(pptx_path: str, output_dir: str) -> list[str]:
    """Convert PPTX → PDF → individual PNGs per slide."""
    lo_paths = ["soffice", "/Applications/LibreOffice.app/Contents/MacOS/soffice"]
    lo = None
    for p in lo_paths:
        if p == "soffice" or os.path.exists(p):
            lo = p
            break
    if not lo:
        raise RuntimeError("LibreOffice not found. Install: brew install --cask libreoffice")

    # Step 1: PPTX → PDF
    subprocess.run([
        lo, "--headless", "--convert-to", "pdf",
        "--outdir", output_dir, pptx_path
    ], check=True, capture_output=True, timeout=120)

    pdf_name = Path(pptx_path).stem + ".pdf"
    pdf_path = os.path.join(output_dir, pdf_name)

    if not os.path.exists(pdf_path):
        raise RuntimeError(f"PDF not created: {pdf_path}")

    # Step 2: PDF → PNGs using sips or convert
    images = []
    try:
        # Use Python's pdf2image if available
        from pdf2image import convert_from_path
        pil_images = convert_from_path(pdf_path, dpi=150)
        for i, img in enumerate(pil_images):
            img_path = os.path.join(output_dir, f"slide_{i+1:02d}.png")
            img.save(img_path, "PNG")
            images.append(img_path)
    except ImportError:
        # Fallback: use macOS Preview via automator or just use the PDF
        # We'll send PDF pages... but for simplicity, try pdftoppm
        try:
            subprocess.run([
                "pdftoppm", "-png", "-r", "150",
                pdf_path, os.path.join(output_dir, "slide")
            ], check=True, capture_output=True, timeout=60)
            images = sorted([str(p) for p in Path(output_dir).glob("slide-*.png")])
        except FileNotFoundError:
            raise RuntimeError("Install pdf2image: pip install pdf2image, or poppler: brew install poppler")

    return images


def get_shape_info(pptx_path: str) -> list[list[dict]]:
    """Extract shape info from each slide."""
    prs = Presentation(pptx_path)
    all_slides = []
    for slide_idx, slide in enumerate(prs.slides):
        shapes = []
        for shape in slide.shapes:
            info = {
                "name": shape.name,
                "left": round(shape.left / 914400, 2),
                "top": round(shape.top / 914400, 2),
                "width": round(shape.width / 914400, 2),
                "height": round(shape.height / 914400, 2),
                "has_text": shape.has_text_frame,
                "text": "",
                "has_image": False,
            }
            if shape.has_text_frame:
                info["text"] = shape.text_frame.text[:100]
            try:
                _ = shape.image
                info["has_image"] = True
            except:
                pass
            # Check shape type
            try:
                info["shape_type"] = str(shape.shape_type)
            except:
                info["shape_type"] = "unknown"

            shapes.append(info)
        all_slides.append(shapes)
    return all_slides


async def analyze_slide_with_vision(
    image_path: str,
    slide_idx: int,
    shapes_info: list[dict],
    total_slides: int,
) -> dict:
    """Send slide image + shape info to GPT-4o Vision for role classification."""

    with open(image_path, "rb") as f:
        img_b64 = base64.b64encode(f.read()).decode()

    shapes_desc = json.dumps(shapes_info, ensure_ascii=False, indent=2)

    prompt = f"""Анализируй слайд {slide_idx + 1} из {total_slides} презентации.

Вот список фигур (shapes) на этом слайде:
{shapes_desc}

Для каждой текстовой фигуры (has_text=true) определи её роль:
- "title" — заголовок слайда (крупный текст сверху)
- "subtitle" — подзаголовок (только на титульном слайде)
- "body" — основной текст/контент
- "author" — имя автора
- "skip" — декоративный элемент, номер, логотип, или текст который не нужно заполнять

Для фигур с изображениями (has_image=true) определи:
- "image" — если это место для пользовательской картинки
- "skip" — если это декоративный элемент/лого/иконка (маленькие < 1 inch)

Ответь ТОЛЬКО в формате JSON:
{{
  "slide_type": "title|content|plan|conclusion|references",
  "shapes": {{
    "shape_name": "role",
    ...
  }}
}}

Важно:
- Слайд 1 обычно title slide (title + subtitle + author)
- Последний слайд обычно "Спасибо"/"Thank you" (conclusion)
- Предпоследний может быть "Adabiyotlar"/"References"
- Маленькие фигуры с числами (01, 02) — skip
- Фоновые картинки на весь слайд — skip
- Маленькие иконки/логотипы (< 1 inch) — skip
"""

    resp = await client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:image/png;base64,{img_b64}",
                            "detail": "high",
                        },
                    },
                ],
            }
        ],
        max_tokens=1500,
        temperature=0.1,
        response_format={"type": "json_object"},
    )

    result = json.loads(resp.choices[0].message.content)
    return result


def apply_tags(pptx_path: str, output_path: str, all_roles: list[dict]):
    """Apply tags to shapes based on Vision analysis."""
    prs = Presentation(pptx_path)

    tag_map = {
        "title": "{{title}}",
        "subtitle": "{{subtitle}}",
        "body": "{{body}}",
        "author": "{{author}}",
        "image": "{{image}}",
    }

    tagged_count = 0
    for slide_idx, slide in enumerate(prs.slides):
        if slide_idx >= len(all_roles):
            break

        roles = all_roles[slide_idx].get("shapes", {})

        for shape in slide.shapes:
            role = roles.get(shape.name, "skip")

            if role == "skip" or role not in tag_map:
                continue

            if role == "image":
                # For images, we need to add {{image}} tag to a text frame
                if shape.has_text_frame:
                    shape.text_frame.paragraphs[0].text = "{{image}}"
                    tagged_count += 1
                continue

            if shape.has_text_frame:
                tag = tag_map[role]
                # Replace text with tag
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = ""
                shape.text_frame.paragraphs[0].text = tag
                tagged_count += 1

    prs.save(output_path)
    return tagged_count


async def auto_tag(pptx_path: str):
    """Main function: analyze template with Vision and apply tags."""
    pptx_path = os.path.abspath(pptx_path)
    basename = Path(pptx_path).stem

    print(f"📄 Analyzing: {pptx_path}")

    # Get shape info
    shapes_info = get_shape_info(pptx_path)
    total_slides = len(shapes_info)
    print(f"   {total_slides} slides found")

    # Convert to images
    with tempfile.TemporaryDirectory() as tmpdir:
        print("🖼️  Converting slides to images...")
        try:
            images = pptx_to_pdf_to_images(pptx_path, tmpdir)
        except Exception as e:
            print(f"   ⚠️  Image conversion failed: {e}")
            print("   Falling back to text-only analysis...")
            images = []

        if len(images) == 1 and total_slides > 1:
            # LibreOffice sometimes creates single image, we need per-slide
            print(f"   ⚠️  Only 1 image for {total_slides} slides — using text-only analysis")
            images = []

        # Analyze each slide
        all_roles = []
        for i in range(total_slides):
            print(f"   🔍 Slide {i+1}/{total_slides}...", end=" ", flush=True)

            try:
                if i < len(images):
                    result = await analyze_slide_with_vision(
                        images[i], i, shapes_info[i], total_slides
                    )
                else:
                    result = await analyze_slide_text_only(
                        i, shapes_info[i], total_slides
                    )

                slide_type = result.get("slide_type", "content")
                roles = result.get("shapes", {})
                tagged = [k for k, v in roles.items() if v != "skip"]
                print(f"type={slide_type}, tagged={len(tagged)} shapes")
                all_roles.append(result)
                await asyncio.sleep(1)  # rate limit protection
            except Exception as e:
                print(f"ERROR: {e}")
                # Fill remaining with empty
                all_roles.append({"slide_type": "content", "shapes": {}})
                if "quota" in str(e).lower() or "rate" in str(e).lower():
                    print("   ⚠️  API quota/rate limit — saving what we have so far...")
                    break

    # Apply tags
    output_path = pptx_path.replace(".pptx", "_tagged.pptx")
    tagged_count = apply_tags(pptx_path, output_path, all_roles)

    print(f"\n✅ Done! Tagged {tagged_count} shapes")
    print(f"📁 Saved: {output_path}")
    return output_path


async def analyze_slide_text_only(
    slide_idx: int,
    shapes_info: list[dict],
    total_slides: int,
) -> dict:
    """Fallback: analyze slide without image, using only shape metadata."""
    shapes_desc = json.dumps(shapes_info, ensure_ascii=False, indent=2)

    prompt = f"""Анализируй слайд {slide_idx + 1} из {total_slides} презентации.
    
Вот список фигур на слайде (позиция в дюймах, текст, тип):
{shapes_desc}

Определи роль каждой текстовой фигуры:
- "title" — заголовок (обычно крупный, сверху)
- "subtitle" — подзаголовок (только на титульном слайде 1)
- "body" — основной текст
- "author" — имя автора
- "image" — место для картинки (has_image=true и размер > 2 inch)
- "skip" — декор, номера, логотипы, маленькие элементы

JSON формат:
{{"slide_type": "title|content|plan|conclusion|references", "shapes": {{"shape_name": "role"}}}}
"""

    resp = await client.chat.completions.create(
        model="gpt-4o",
        messages=[{"role": "user", "content": prompt}],
        max_tokens=1500,
        temperature=0.1,
        response_format={"type": "json_object"},
    )
    return json.loads(resp.choices[0].message.content)


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python scripts/auto_tag_template.py <template.pptx>")
        sys.exit(1)

    asyncio.run(auto_tag(sys.argv[1]))
