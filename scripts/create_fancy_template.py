from __future__ import annotations
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def create_fancy_template(path):
    prs = Presentation()
    
    # ── Slide 1: Title Layout ──────────────────────────────────────────────
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42) # Dark Navy
    
    # Title Styling
    title = slide.shapes.title
    title.text = "Mavzu nomi"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(56, 189, 248) # Light Blue
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    
    # Subtitle
    subtitle = slide.placeholders[1]
    subtitle.text = "Bajardi: Ism Familiya"
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(148, 163, 184) # Gray
    
    # ── Slide 2: Content Layout ───────────────────────────────────────────
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42)
    
    # Title
    title = slide.shapes.title
    title.text = "Slayd sarlavhasi"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(56, 189, 248)
    
    # Body
    body = slide.placeholders[1]
    body.text_frame.text = "Asosiy ma'lumotlar bu yerda bo'ladi"
    body.text_frame.paragraphs[0].font.color.rgb = RGBColor(248, 250, 252) # White-ish
    
    # Add a decorative line
    line = slide.shapes.add_shape(
        1, # Rectangle
        Inches(0.5), Inches(1.1), Inches(9), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(56, 189, 248)
    line.line.fill.background()
    
    os.makedirs(os.path.dirname(path), exist_ok=True)
    prs.save(path)
    print(f"Fancy template created at: {path}")

if __name__ == "__main__":
    create_fancy_template("templates/Biology/Bio_Classic.pptx")
    create_fancy_template("templates/IT/kiberxafsizlik.pptx")
