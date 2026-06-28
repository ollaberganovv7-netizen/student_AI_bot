from __future__ import annotations
"""
Generate professional PPTX templates for each empty category.
Each template has: title slide + 9 content slides (5 varied layouts) + thank-you slide.
Every content slide includes an image placeholder.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

CATEGORIES = {
    "Biology": {
        "accent": RGBColor(0x2E, 0x7D, 0x32),
        "accent2": RGBColor(0x66, 0xBB, 0x6A),
        "bg": RGBColor(0xF1, 0xF8, 0xE9),
        "title_color": RGBColor(0x1B, 0x5E, 0x20),
        "body_color": RGBColor(0x33, 0x33, 0x33),
        "templates": ["nature_green", "cell_light"],
    },
    "Medicine": {
        "accent": RGBColor(0x15, 0x65, 0xC0),
        "accent2": RGBColor(0x42, 0xA5, 0xF5),
        "bg": RGBColor(0xE3, 0xF2, 0xFD),
        "title_color": RGBColor(0x0D, 0x47, 0xA1),
        "body_color": RGBColor(0x33, 0x33, 0x33),
        "templates": ["medical_blue", "health_clean"],
    },
    "Business": {
        "accent": RGBColor(0x37, 0x47, 0x4F),
        "accent2": RGBColor(0x78, 0x90, 0x9C),
        "bg": RGBColor(0xEC, 0xEF, 0xF1),
        "title_color": RGBColor(0x26, 0x32, 0x38),
        "body_color": RGBColor(0x45, 0x45, 0x45),
        "templates": ["corporate_dark", "business_minimal"],
    },
    "IT": {
        "accent": RGBColor(0x6A, 0x1B, 0x9A),
        "accent2": RGBColor(0xAB, 0x47, 0xBC),
        "bg": RGBColor(0xF3, 0xE5, 0xF5),
        "title_color": RGBColor(0x4A, 0x14, 0x8C),
        "body_color": RGBColor(0x33, 0x33, 0x33),
        "templates": ["tech_purple", "code_dark"],
    },
    "History": {
        "accent": RGBColor(0xBF, 0x36, 0x0C),
        "accent2": RGBColor(0xE6, 0x6A, 0x3A),
        "bg": RGBColor(0xFB, 0xE9, 0xE7),
        "title_color": RGBColor(0x8B, 0x1A, 0x00),
        "body_color": RGBColor(0x3E, 0x27, 0x23),
        "templates": ["vintage_warm", "history_classic"],
    },
    "Buxgalteriya": {
        "accent": RGBColor(0x00, 0x69, 0x5C),
        "accent2": RGBColor(0x26, 0xA6, 0x9A),
        "bg": RGBColor(0xE0, 0xF2, 0xF1),
        "title_color": RGBColor(0x00, 0x4D, 0x40),
        "body_color": RGBColor(0x33, 0x33, 0x33),
        "templates": ["finance_teal", "accounting_clean"],
    },
}

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def add_shape(slide, left, top, width, height, color, shape_type=MSO_SHAPE.RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Remove text frame so filler doesn't treat decorative shapes as text containers
    for child in list(shape._element):
        if child.tag.endswith('}txBody') or child.tag == 'txBody':
            shape._element.remove(child)
            break
    return shape


def add_textbox(slide, left, top, width, height, text, font_size, font_color,
                bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = font_name
    return txBox


def add_image_placeholder(slide, left, top, width, height, accent):
    """Add a styled image placeholder rectangle with {{image}} tag."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "{{image}}"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    tf.paragraphs[0].space_before = Pt(0)
    return shape


# ─── 5 Layout builders ──────────────────────────────────────────────────────

def layout_A_text_left_image_right(slide, accent, accent2, bg, tc, bc, dark):
    """Layout A: Title top, body left (60%), image right (35%)"""
    add_shape(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, bg)
    add_shape(slide, 0, 0, SLIDE_WIDTH, Inches(0.06), accent)
    add_shape(slide, Inches(12.8), 0, Inches(0.533), SLIDE_HEIGHT, accent)
    add_textbox(slide, Inches(0.7), Inches(0.3), Inches(7.5), Inches(1.0),
                "{{title}}", 34, tc, bold=True)
    add_shape(slide, Inches(0.7), Inches(1.2), Inches(2.0), Inches(0.05), accent)
    add_textbox(slide, Inches(0.7), Inches(1.5), Inches(7.0), Inches(5.2),
                "{{body}}", 16, bc)
    add_image_placeholder(slide, Inches(8.2), Inches(0.5), Inches(4.4), Inches(6.2), accent)


def layout_B_image_left_text_right(slide, accent, accent2, bg, tc, bc, dark):
    """Layout B: Image left (40%), title+body right (55%)"""
    add_shape(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, bg)
    add_shape(slide, 0, 0, Inches(0.12), SLIDE_HEIGHT, accent)
    add_image_placeholder(slide, Inches(0.5), Inches(0.5), Inches(5.0), Inches(6.3), accent)
    add_textbox(slide, Inches(6.0), Inches(0.5), Inches(6.8), Inches(1.0),
                "{{title}}", 32, tc, bold=True)
    add_shape(slide, Inches(6.0), Inches(1.4), Inches(3.0), Inches(0.05), accent2)
    add_textbox(slide, Inches(6.0), Inches(1.8), Inches(6.8), Inches(5.0),
                "{{body}}", 16, bc)


def layout_C_full_width_image_bottom(slide, accent, accent2, bg, tc, bc, dark):
    """Layout C: Title + body top half, wide image bottom half"""
    add_shape(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, bg)
    add_shape(slide, 0, 0, SLIDE_WIDTH, Inches(0.06), accent)
    add_shape(slide, 0, Inches(7.35), SLIDE_WIDTH, Inches(0.15), accent)
    add_textbox(slide, Inches(0.7), Inches(0.3), Inches(12), Inches(0.9),
                "{{title}}", 34, tc, bold=True)
    add_shape(slide, Inches(0.7), Inches(1.15), Inches(2.0), Inches(0.05), accent)
    add_textbox(slide, Inches(0.7), Inches(1.4), Inches(12), Inches(2.3),
                "{{body}}", 16, bc)
    add_image_placeholder(slide, Inches(0.7), Inches(3.9), Inches(11.9), Inches(3.3), accent)


def layout_D_two_column_with_image(slide, accent, accent2, bg, tc, bc, dark):
    """Layout D: Title across top, two text columns + image in right column"""
    add_shape(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, bg)
    # Colored header bar
    add_shape(slide, 0, 0, SLIDE_WIDTH, Inches(1.2), accent)
    add_textbox(slide, Inches(0.7), Inches(0.15), Inches(12), Inches(0.9),
                "{{title}}", 32, RGBColor(0xFF, 0xFF, 0xFF), bold=True)
    add_textbox(slide, Inches(0.7), Inches(1.5), Inches(5.5), Inches(5.2),
                "{{body}}", 16, bc)
    # Vertical divider
    add_shape(slide, Inches(6.5), Inches(1.5), Inches(0.04), Inches(5.2), accent2)
    add_image_placeholder(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.2), accent)


def layout_E_image_bg_text_overlay(slide, accent, accent2, bg, tc, bc, dark):
    """Layout E: Large image top, text overlay at bottom with accent bg strip"""
    add_shape(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, bg)
    add_image_placeholder(slide, Inches(0.5), Inches(0.4), Inches(12.3), Inches(3.8), accent)
    # Accent strip behind text
    strip_color = accent if not dark else accent2
    add_shape(slide, 0, Inches(4.4), SLIDE_WIDTH, Inches(0.08), strip_color)
    add_textbox(slide, Inches(0.7), Inches(4.7), Inches(12), Inches(0.9),
                "{{title}}", 30, tc, bold=True)
    add_textbox(slide, Inches(0.7), Inches(5.5), Inches(12), Inches(1.6),
                "{{body}}", 16, bc)


# ─── Text-only layouts (no image) ────────────────────────────────────────────

def layout_T1_full_text(slide, accent, accent2, bg, tc, bc, dark):
    """Text-only: Full width title + body with left accent bar."""
    add_shape(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, bg)
    add_shape(slide, 0, 0, Inches(0.15), SLIDE_HEIGHT, accent)
    add_shape(slide, 0, 0, SLIDE_WIDTH, Inches(0.06), accent)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(1.0),
                "{{title}}", 36, tc, bold=True)
    add_shape(slide, Inches(0.8), Inches(1.3), Inches(2.5), Inches(0.05), accent)
    add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.2),
                "{{body}}", 17, bc)


def layout_T2_accent_header(slide, accent, accent2, bg, tc, bc, dark):
    """Text-only: Colored accent header band + full body below."""
    add_shape(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, bg)
    add_shape(slide, 0, 0, SLIDE_WIDTH, Inches(1.4), accent)
    add_textbox(slide, Inches(0.8), Inches(0.2), Inches(11.5), Inches(1.0),
                "{{title}}", 34, RGBColor(0xFF, 0xFF, 0xFF), bold=True)
    add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5.0),
                "{{body}}", 17, bc)
    add_shape(slide, 0, Inches(7.3), SLIDE_WIDTH, Inches(0.2), accent2)


def layout_T3_two_column_text(slide, accent, accent2, bg, tc, bc, dark):
    """Text-only: Title top, two text columns below with divider."""
    add_shape(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, bg)
    add_shape(slide, 0, 0, SLIDE_WIDTH, Inches(0.06), accent)
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(1.0),
                "{{title}}", 34, tc, bold=True)
    add_shape(slide, Inches(0.8), Inches(1.2), Inches(2.5), Inches(0.05), accent)
    # Left column
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.3),
                "{{body}}", 16, bc)
    # Vertical divider
    add_shape(slide, Inches(6.6), Inches(1.5), Inches(0.04), Inches(5.3), accent2)
    # Right column — body_2 placeholder (will show continuation or points)
    add_textbox(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.3),
                "{{body}}", 16, bc)
    # Corner accent
    add_shape(slide, Inches(12.8), Inches(6.8), Inches(0.533), Inches(0.7), accent)


def layout_T4_quote_style(slide, accent, accent2, bg, tc, bc, dark):
    """Text-only: Large quote/highlight style with accent side bar."""
    add_shape(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, bg)
    # Big accent bar on left
    add_shape(slide, Inches(0.5), Inches(0.5), Inches(0.12), Inches(6.5), accent)
    add_textbox(slide, Inches(1.2), Inches(0.5), Inches(11), Inches(1.2),
                "{{title}}", 36, tc, bold=True)
    add_shape(slide, Inches(1.2), Inches(1.7), Inches(3.0), Inches(0.05), accent2)
    add_textbox(slide, Inches(1.2), Inches(2.0), Inches(11), Inches(4.8),
                "{{body}}", 17, bc)
    # Bottom right decorative shape
    add_shape(slide, Inches(11.5), Inches(6.0), Inches(1.5), Inches(1.2),
              accent2, MSO_SHAPE.ROUNDED_RECTANGLE)


def layout_T5_centered_text(slide, accent, accent2, bg, tc, bc, dark):
    """Text-only: Centered title with body, horizontal accent lines."""
    add_shape(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, bg)
    add_shape(slide, 0, 0, SLIDE_WIDTH, Inches(0.06), accent)
    add_shape(slide, 0, Inches(7.44), SLIDE_WIDTH, Inches(0.06), accent)
    add_textbox(slide, Inches(1.5), Inches(0.4), Inches(10.3), Inches(1.0),
                "{{title}}", 34, tc, bold=True, alignment=PP_ALIGN.CENTER)
    add_shape(slide, Inches(5.5), Inches(1.3), Inches(2.3), Inches(0.05), accent)
    add_textbox(slide, Inches(1.0), Inches(1.7), Inches(11.3), Inches(5.2),
                "{{body}}", 17, bc)


def layout_S1_process_steps(slide, accent, accent2, bg, tc, bc, dark):
    """SmartArt-like: 4 numbered process step boxes with arrows."""
    add_shape(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, bg)
    add_shape(slide, 0, 0, SLIDE_WIDTH, Inches(0.06), accent)
    add_textbox(slide, Inches(0.7), Inches(0.3), Inches(12), Inches(0.9),
                "{{title}}", 34, tc, bold=True)
    add_shape(slide, Inches(0.7), Inches(1.15), Inches(2.5), Inches(0.05), accent)

    # 4 step boxes
    box_w = Inches(2.6)
    box_h = Inches(4.0)
    gap = Inches(0.35)
    start_x = Inches(0.7)
    top_y = Inches(1.8)

    for step in range(4):
        x = start_x + step * (box_w + gap)
        # Card background
        card = add_shape(slide, x, top_y, box_w, box_h, accent2 if not dark else accent, MSO_SHAPE.ROUNDED_RECTANGLE)
        card.fill.solid()
        card.fill.fore_color.rgb = accent2 if not dark else accent
        # Number circle (needs text frame, so create manually)
        num_size = Inches(0.7)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.9), top_y + Inches(0.3), num_size, num_size)
        circle.fill.solid()
        circle.fill.fore_color.rgb = accent
        circle.line.fill.background()
        tf = circle.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = str(step + 1)
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        # Step text area
        add_textbox(slide, x + Inches(0.2), top_y + Inches(1.3), box_w - Inches(0.4), box_h - Inches(1.6),
                    "{{body}}", 14, RGBColor(0xFF, 0xFF, 0xFF) if dark else RGBColor(0x33, 0x33, 0x33))

        # Arrow between boxes (except after last)
        if step < 3:
            arrow_x = x + box_w + Inches(0.03)
            add_shape(slide, arrow_x, top_y + Inches(1.8), Inches(0.3), Inches(0.4),
                      accent, MSO_SHAPE.RIGHT_ARROW)

    # Bottom accent
    add_shape(slide, 0, Inches(7.3), SLIDE_WIDTH, Inches(0.2), accent)


def layout_S2_card_grid(slide, accent, accent2, bg, tc, bc, dark):
    """SmartArt-like: 2x2 card grid with colored headers."""
    add_shape(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, bg)
    # Colored header bar
    add_shape(slide, 0, 0, SLIDE_WIDTH, Inches(1.3), accent)
    add_textbox(slide, Inches(0.7), Inches(0.15), Inches(12), Inches(0.9),
                "{{title}}", 32, RGBColor(0xFF, 0xFF, 0xFF), bold=True)

    # 2x2 grid of cards
    card_w = Inches(5.7)
    card_h = Inches(2.6)
    gap_x = Inches(0.5)
    gap_y = Inches(0.4)
    start_x = Inches(0.7)
    start_y = Inches(1.7)

    card_colors = [accent, accent2, accent2, accent]
    for row in range(2):
        for col in range(2):
            idx = row * 2 + col
            x = start_x + col * (card_w + gap_x)
            y = start_y + row * (card_h + gap_y)
            # Card
            card = add_shape(slide, x, y, card_w, card_h, card_colors[idx], MSO_SHAPE.ROUNDED_RECTANGLE)
            # Card text
            add_textbox(slide, x + Inches(0.3), y + Inches(0.3), card_w - Inches(0.6), card_h - Inches(0.6),
                        "{{body}}", 15, RGBColor(0xFF, 0xFF, 0xFF), bold=False)


# ─── All available content layouts ───────────────────────────────────────────
IMAGE_LAYOUTS = [layout_A_text_left_image_right, layout_B_image_left_text_right, layout_C_full_width_image_bottom]
TEXT_LAYOUTS = [layout_T1_full_text, layout_T2_accent_header, layout_T3_two_column_text, layout_T4_quote_style, layout_T5_centered_text]
SMART_LAYOUTS = [layout_S1_process_steps, layout_S2_card_grid]

# 12 unique layout sequences — each template gets a different one
LAYOUT_VARIANTS = [
    # Variant 0: text → IMG → smart → text → IMG → smart → text → IMG
    [TEXT_LAYOUTS[0], IMAGE_LAYOUTS[0], SMART_LAYOUTS[0], TEXT_LAYOUTS[1], IMAGE_LAYOUTS[1], SMART_LAYOUTS[1], TEXT_LAYOUTS[3], IMAGE_LAYOUTS[2]],
    # Variant 1: smart → text → IMG → text → smart → IMG → text → IMG
    [SMART_LAYOUTS[1], TEXT_LAYOUTS[2], IMAGE_LAYOUTS[1], TEXT_LAYOUTS[0], SMART_LAYOUTS[0], IMAGE_LAYOUTS[0], TEXT_LAYOUTS[4], IMAGE_LAYOUTS[2]],
    # Variant 2: text → text → IMG → smart → text → IMG → smart → IMG
    [TEXT_LAYOUTS[1], TEXT_LAYOUTS[3], IMAGE_LAYOUTS[2], SMART_LAYOUTS[0], TEXT_LAYOUTS[0], IMAGE_LAYOUTS[1], SMART_LAYOUTS[1], IMAGE_LAYOUTS[0]],
    # Variant 3: IMG → smart → text → text → IMG → text → smart → IMG
    [IMAGE_LAYOUTS[0], SMART_LAYOUTS[0], TEXT_LAYOUTS[4], TEXT_LAYOUTS[1], IMAGE_LAYOUTS[2], TEXT_LAYOUTS[3], SMART_LAYOUTS[1], IMAGE_LAYOUTS[1]],
    # Variant 4: text → smart → IMG → smart → text → IMG → text → IMG
    [TEXT_LAYOUTS[2], SMART_LAYOUTS[1], IMAGE_LAYOUTS[0], SMART_LAYOUTS[0], TEXT_LAYOUTS[1], IMAGE_LAYOUTS[2], TEXT_LAYOUTS[0], IMAGE_LAYOUTS[1]],
    # Variant 5: smart → IMG → text → text → smart → IMG → text → IMG
    [SMART_LAYOUTS[0], IMAGE_LAYOUTS[1], TEXT_LAYOUTS[0], TEXT_LAYOUTS[4], SMART_LAYOUTS[1], IMAGE_LAYOUTS[0], TEXT_LAYOUTS[3], IMAGE_LAYOUTS[2]],
    # Variant 6: text → IMG → text → smart → IMG → text → smart → IMG
    [TEXT_LAYOUTS[4], IMAGE_LAYOUTS[2], TEXT_LAYOUTS[1], SMART_LAYOUTS[1], IMAGE_LAYOUTS[0], TEXT_LAYOUTS[2], SMART_LAYOUTS[0], IMAGE_LAYOUTS[1]],
    # Variant 7: IMG → text → smart → text → IMG → smart → text → IMG
    [IMAGE_LAYOUTS[1], TEXT_LAYOUTS[3], SMART_LAYOUTS[1], TEXT_LAYOUTS[0], IMAGE_LAYOUTS[2], SMART_LAYOUTS[0], TEXT_LAYOUTS[4], IMAGE_LAYOUTS[0]],
    # Variant 8: smart → text → IMG → text → text → IMG → smart → IMG
    [SMART_LAYOUTS[1], TEXT_LAYOUTS[1], IMAGE_LAYOUTS[0], TEXT_LAYOUTS[3], TEXT_LAYOUTS[2], IMAGE_LAYOUTS[1], SMART_LAYOUTS[0], IMAGE_LAYOUTS[2]],
    # Variant 9: text → smart → IMG → text → smart → IMG → text → IMG
    [TEXT_LAYOUTS[3], SMART_LAYOUTS[0], IMAGE_LAYOUTS[2], TEXT_LAYOUTS[4], SMART_LAYOUTS[1], IMAGE_LAYOUTS[1], TEXT_LAYOUTS[0], IMAGE_LAYOUTS[0]],
    # Variant 10: IMG → text → text → smart → IMG → text → smart → IMG
    [IMAGE_LAYOUTS[2], TEXT_LAYOUTS[0], TEXT_LAYOUTS[1], SMART_LAYOUTS[0], IMAGE_LAYOUTS[1], TEXT_LAYOUTS[4], SMART_LAYOUTS[1], IMAGE_LAYOUTS[0]],
    # Variant 11: text → IMG → smart → text → IMG → text → text → IMG
    [TEXT_LAYOUTS[2], IMAGE_LAYOUTS[1], SMART_LAYOUTS[0], TEXT_LAYOUTS[3], IMAGE_LAYOUTS[0], TEXT_LAYOUTS[1], TEXT_LAYOUTS[4], IMAGE_LAYOUTS[2]],
]


# ─── Title slide variants ────────────────────────────────────────────────────

def title_slide_v1(slide, accent, accent2, bg, tc, bc):
    """Left accent bar + decorative circle"""
    add_shape(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, bg)
    add_shape(slide, 0, 0, Inches(0.2), SLIDE_HEIGHT, accent)
    add_shape(slide, 0, Inches(7.2), SLIDE_WIDTH, Inches(0.3), accent)
    add_shape(slide, Inches(10.5), Inches(1.0), Inches(2.0), Inches(2.0), accent2, MSO_SHAPE.OVAL)
    add_textbox(slide, Inches(1.0), Inches(2.0), Inches(9), Inches(1.5),
                "{{title}}", 44, tc, bold=True)
    add_shape(slide, Inches(1.0), Inches(3.5), Inches(3.0), Inches(0.08), accent)
    add_textbox(slide, Inches(1.0), Inches(3.9), Inches(9), Inches(1.0), "{{subtitle}}", 22, bc)
    add_textbox(slide, Inches(1.0), Inches(5.5), Inches(9), Inches(0.7), "{{author}}", 18, bc)


def title_slide_v2(slide, accent, accent2, bg, tc, bc):
    """Centered with top/bottom accent bands"""
    add_shape(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, bg)
    add_shape(slide, 0, 0, SLIDE_WIDTH, Inches(0.8), accent)
    add_shape(slide, 0, Inches(6.7), SLIDE_WIDTH, Inches(0.8), accent)
    add_textbox(slide, Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.5),
                "{{title}}", 48, tc, bold=True, alignment=PP_ALIGN.CENTER)
    add_shape(slide, Inches(5.0), Inches(3.7), Inches(3.3), Inches(0.07), accent2)
    add_textbox(slide, Inches(1.5), Inches(4.0), Inches(10.3), Inches(1.0),
                "{{subtitle}}", 22, bc, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.5), Inches(5.3), Inches(10.3), Inches(0.7),
                "{{author}}", 18, bc, alignment=PP_ALIGN.CENTER)


def title_slide_v3(slide, accent, accent2, bg, tc, bc):
    """Right accent block + left-aligned title"""
    add_shape(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, bg)
    add_shape(slide, Inches(9.0), 0, Inches(4.333), SLIDE_HEIGHT, accent)
    add_shape(slide, Inches(8.7), Inches(2.0), Inches(0.12), Inches(3.5), accent2)
    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(7.5), Inches(1.5),
                "{{title}}", 44, tc, bold=True)
    add_shape(slide, Inches(0.8), Inches(3.5), Inches(2.5), Inches(0.07), accent)
    add_textbox(slide, Inches(0.8), Inches(3.8), Inches(7.5), Inches(1.0), "{{subtitle}}", 22, bc)
    add_textbox(slide, Inches(0.8), Inches(5.5), Inches(7.5), Inches(0.7), "{{author}}", 18, bc)


def title_slide_v4(slide, accent, accent2, bg, tc, bc):
    """Full accent background with white text"""
    add_shape(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, accent)
    add_shape(slide, Inches(0.5), Inches(0.5), Inches(1.2), Inches(1.2), accent2, MSO_SHAPE.OVAL)
    add_shape(slide, Inches(11.5), Inches(5.8), Inches(1.5), Inches(1.5), accent2, MSO_SHAPE.OVAL)
    white = RGBColor(0xFF, 0xFF, 0xFF)
    light = RGBColor(0xDD, 0xDD, 0xDD)
    add_textbox(slide, Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.5),
                "{{title}}", 48, white, bold=True, alignment=PP_ALIGN.CENTER)
    add_shape(slide, Inches(5.0), Inches(3.7), Inches(3.3), Inches(0.07), accent2)
    add_textbox(slide, Inches(1.5), Inches(4.0), Inches(10.3), Inches(1.0),
                "{{subtitle}}", 22, light, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.5), Inches(5.3), Inches(10.3), Inches(0.7),
                "{{author}}", 18, light, alignment=PP_ALIGN.CENTER)


TITLE_VARIANTS = [title_slide_v1, title_slide_v2, title_slide_v3, title_slide_v4]


# ─── Thank-you slide variants ───────────────────────────────────────────────

def thanks_slide_v1(slide, accent, accent2):
    add_shape(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, accent)
    add_shape(slide, Inches(0.5), Inches(0.5), Inches(1.5), Inches(1.5), accent2, MSO_SHAPE.OVAL)
    add_shape(slide, Inches(11.0), Inches(5.5), Inches(1.8), Inches(1.8), accent2, MSO_SHAPE.OVAL)
    add_textbox(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(2.0),
                "{{title}}", 48, RGBColor(0xFF, 0xFF, 0xFF), bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(1.0),
                "{{author}}", 22, RGBColor(0xDD, 0xDD, 0xDD), alignment=PP_ALIGN.CENTER)


def thanks_slide_v2(slide, accent, accent2):
    add_shape(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, accent)
    add_shape(slide, Inches(2.0), Inches(1.5), Inches(9.3), Inches(4.5), accent2, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, Inches(2.5), Inches(2.3), Inches(8.3), Inches(2.0),
                "{{title}}", 48, RGBColor(0xFF, 0xFF, 0xFF), bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2.5), Inches(4.3), Inches(8.3), Inches(1.0),
                "{{author}}", 22, RGBColor(0xDD, 0xDD, 0xDD), alignment=PP_ALIGN.CENTER)


THANKS_VARIANTS = [thanks_slide_v1, thanks_slide_v2]


def create_template(style_name, accent, accent2, bg, title_color, body_color, dark_mode=False, variant_idx=0):
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    if dark_mode:
        bg = RGBColor(0x1E, 0x1E, 0x2E)
        body_color = RGBColor(0xCC, 0xCC, 0xCC)
        title_color = RGBColor(0xFF, 0xFF, 0xFF)

    layout = prs.slide_layouts[6]  # Blank layout

    # ── Slide 1: Title (varied) ──
    slide = prs.slides.add_slide(layout)
    title_fn = TITLE_VARIANTS[variant_idx % len(TITLE_VARIANTS)]
    title_fn(slide, accent, accent2, bg, title_color, body_color)

    # ── Slides 2-9: Content (unique sequence per variant) ──
    content_layouts = LAYOUT_VARIANTS[variant_idx % len(LAYOUT_VARIANTS)]
    for i in range(8):
        slide = prs.slides.add_slide(layout)
        layout_fn = content_layouts[i % len(content_layouts)]
        layout_fn(slide, accent, accent2, bg, title_color, body_color, dark_mode)

    # ── Slide 10: Xulosa (Conclusion) ──
    slide = prs.slides.add_slide(layout)
    add_shape(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, bg)
    add_shape(slide, 0, 0, SLIDE_WIDTH, Inches(1.3), accent)
    add_textbox(slide, Inches(0.8), Inches(0.15), Inches(12), Inches(0.9),
                "{{title}}", 34, RGBColor(0xFF, 0xFF, 0xFF), bold=True)
    add_textbox(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(5.0),
                "{{body}}", 17, body_color)
    add_shape(slide, 0, Inches(7.3), SLIDE_WIDTH, Inches(0.2), accent2)

    # ── Slide 11: Foydalanilgan adabiyotlar (References) ──
    slide = prs.slides.add_slide(layout)
    add_shape(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, bg)
    add_shape(slide, 0, 0, Inches(0.15), SLIDE_HEIGHT, accent)
    add_shape(slide, 0, 0, SLIDE_WIDTH, Inches(0.06), accent)
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(1.0),
                "{{title}}", 34, title_color, bold=True)
    add_shape(slide, Inches(0.8), Inches(1.2), Inches(2.5), Inches(0.05), accent)
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5),
                "{{body}}", 16, body_color)

    # ── Slide 12: Thank you (varied) ──
    slide = prs.slides.add_slide(layout)
    thanks_fn = THANKS_VARIANTS[variant_idx % len(THANKS_VARIANTS)]
    thanks_fn(slide, accent, accent2)

    return prs


def main():
    base = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    templates_dir = os.path.join(base, "templates", "general")
    os.makedirs(templates_dir, exist_ok=True)

    total = 0
    variant_counter = 0
    for cat, cfg in CATEGORIES.items():
        for idx, style in enumerate(cfg["templates"]):
            dark_mode = idx == 1
            prs = create_template(
                style_name=style,
                accent=cfg["accent"],
                accent2=cfg["accent2"],
                bg=cfg["bg"],
                title_color=cfg["title_color"],
                body_color=cfg["body_color"],
                dark_mode=dark_mode,
                variant_idx=variant_counter,
            )

            path = os.path.join(templates_dir, f"{style}.pptx")
            prs.save(path)
            total += 1
            variant_counter += 1
            print(f"  \u2705 general/{style}.pptx (variant {variant_counter})")

    print(f"\nDone! Created {total} templates in general/ — all unique layouts.")


if __name__ == "__main__":
    main()
