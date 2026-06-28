from __future__ import annotations
from services.template_filler_service import analyze_template_deep, fill_template
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

analysis = analyze_template_deep('templates/modern/vena_shablon.pptx')

slides_content = [
    {'title': 'Vena tizimi', 'subtitle': 'Bajardi: Arslan'},
    {'title': 'Reja', 'point_1': 'Kirish', 'point_2': 'Struktura', 'point_3': 'Funksiya', 'point_4': 'Kasalliklar', 'point_5': 'Xulosa'},
    {'title': 'Kirish', 'body': 'Vena tizimi qon aylanish tizimining muhim qismi.'},
    {'point_1': 'Tuzilish', 'subtitle_1': 'Katta venalar', 'body_1': 'Yirik venalar tuzilishi', 'subtitle_2': 'Kapillyarlar', 'body_2': 'Kichik qon tomirlari', 'subtitle_3': 'Valflar', 'body_3': 'Qon oqimini boshqaradi'},
    {'title': 'Qismlar', 'body_1': 'Katta venalar', 'body_2': 'Kichik venalar', 'body_3': 'Kapillyarlar'},
    {'title': 'Valflar', 'body': 'Valflar qon harakatini nazorat qiladi.'},
    {'point_2': 'Funksiya', 'subtitle_1': 'Qon qaytarish', 'subtitle_2': 'Chiqindi', 'subtitle_3': 'Harorat'},
    {'title': 'Funksiyalar', 'body_1': 'A', 'body_2': 'B', 'body_3': 'C'},
    {'title': 'Kasalliklar', 'body': 'Varikoz, tromboz.'},
    {'point_3': 'Kasalliklar', 'subtitle_1': 'Varikoz', 'subtitle_2': 'Tromboz'},
    {'title': 'Davolash', 'body_1': 'Varikoz davolash', 'body_2': 'Tromboz davolash'},
    {'title': 'Saqlash', 'body': 'Mashqlar va parhez.'},
    {'title': 'Xulosa', 'body': 'Vena tizimi juda muhim.'},
    {'title': 'Rahmat'},
]

result_bytes = fill_template(
    template_path='templates/modern/vena_shablon.pptx',
    slides_content=slides_content,
    topic='Vena tizimi',
    author='Arslan',
    analysis=analysis,
)

with open('/tmp/test_fill_result.pptx', 'wb') as f:
    f.write(result_bytes)

prs = Presentation('/tmp/test_fill_result.pptx')
for si, slide in enumerate(prs.slides):
    texts = []
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub in sh.shapes:
                if sub.has_text_frame and sub.text_frame.text.strip():
                    texts.append(sub.text_frame.text.strip()[:50])
        elif sh.has_text_frame and sh.text_frame.text.strip():
            texts.append(sh.text_frame.text.strip()[:50])
    if texts:
        print(f'Slide {si+1}: {texts}')
    else:
        print(f'Slide {si+1}: [EMPTY]')
