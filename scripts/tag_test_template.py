from __future__ import annotations
from pptx import Presentation
import shutil

src = "templates/power point/test.pptx"
shutil.copy(src, "templates/power point/test_backup.pptx")

prs = Presentation(src)
slides = list(prs.slides)

tags = {
    0: {"Заголовок 1": "{{title}}", "Подзаголовок 2": "{{subtitle}}"},
    1: {"Заголовок 1": "Reja"},
    2: {"Заголовок 5": "{{title}}", "Объект 6": "{{body}}"},
    3: {"Заголовок 7": "{{point_1}}"},
    4: {"Заголовок 3": "{{title}}", "Объект 4": "{{body_1}}", "Объект 5": "{{body_2}}", "Объект 6": "{{body_3}}"},
    5: {"Заголовок 7": "{{title}}", "Текст 9": "{{body}}"},
    6: {"Заголовок 7": "{{point_2}}"},
    7: {"Заголовок 3": "{{title}}", "Объект 4": "{{body_1}}", "Объект 5": "{{body_2}}", "Объект 6": "{{body_3}}"},
    8: {"Заголовок 4": "{{title}}", "Текст 6": "{{body}}"},
    9: {"Заголовок 1": "{{point_3}}"},
    10: {"Заголовок 3": "{{title}}", "Объект 4": "{{body_1}}", "Объект 5": "{{body_2}}"},
    11: {"Заголовок 4": "{{title}}", "Текст 6": "{{body}}"},
    12: {"Заголовок 1": "{{title}}", "Объект 2": "{{body}}"},
}

for si, slide in enumerate(slides):
    if si not in tags:
        continue
    for shape in slide.shapes:
        if shape.name in tags[si] and shape.has_text_frame:
            shape.text_frame.paragraphs[0].text = tags[si][shape.name]

prs.save(src)
print("Done! Tags added to test.pptx")

prs2 = Presentation(src)
for i, sl in enumerate(prs2.slides):
    for sh in sl.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            print(f"  Slide {i+1}: [{sh.name}] = {sh.text_frame.text.strip()[:40]}")
