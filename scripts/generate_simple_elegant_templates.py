from __future__ import annotations
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

def create_template(path, bg_color, title_color, body_color, accent_color, is_dark=False):
    prs = Presentation()
    
    # ── Slide 1: Title Layout ──────────────────────────────────────────────
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    
    # Title Styling
    title = slide.shapes.title
    title.text = "Taqdimot Mavzusi"
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    title.text_frame.paragraphs[0].font.size = Pt(48)
    title.text_frame.paragraphs[0].font.bold = True
    
    # Subtitle
    subtitle = slide.placeholders[1]
    subtitle.text = "Muallif: Ism Familiya"
    subtitle.text_frame.paragraphs[0].font.color.rgb = body_color
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    
    # Decorative Top Bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()

    # ── Slide 2: Content Layout ───────────────────────────────────────────
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    
    # Title
    title = slide.shapes.title
    title.text = "Slayd sarlavhasi"
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    
    # Body
    body = slide.placeholders[1]
    body.text_frame.text = "• Asosiy ma'lumotlar bu yerda bo'ladi"
    body.text_frame.paragraphs[0].font.color.rgb = body_color
    body.text_frame.paragraphs[0].font.size = Pt(20)
    
    # Decorative line under title
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(9), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = accent_color
    line.line.fill.background()
    
    os.makedirs(os.path.dirname(path), exist_ok=True)
    prs.save(path)
    print(f"Created: {path}")

if __name__ == "__main__":
    # 1. Corporate Dark Blue
    create_template(
        "templates/General/Corporate_Dark.pptx",
        bg_color=RGBColor(15, 23, 42),      # Dark navy
        title_color=RGBColor(56, 189, 248), # Light Blue
        body_color=RGBColor(241, 245, 249), # Off-white
        accent_color=RGBColor(56, 189, 248),
        is_dark=True
    )
    
    # 2. Minimalist Light
    create_template(
        "templates/General/Minimalist_Light.pptx",
        bg_color=RGBColor(250, 250, 250),   # Almost white
        title_color=RGBColor(30, 41, 59),   # Slate dark
        body_color=RGBColor(51, 65, 85),    # Slate medium
        accent_color=RGBColor(245, 158, 11), # Amber
        is_dark=False
    )
    
    # 3. Elegant Emerald
    create_template(
        "templates/General/Elegant_Emerald.pptx",
        bg_color=RGBColor(6, 78, 59),       # Dark emerald
        title_color=RGBColor(167, 243, 208), # Light emerald
        body_color=RGBColor(255, 255, 255),  # White
        accent_color=RGBColor(52, 211, 153), # Emerald accent
        is_dark=True
    )
    
    # 4. Premium Purple
    create_template(
        "templates/General/Premium_Purple.pptx",
        bg_color=RGBColor(46, 16, 101),      # Deep purple
        title_color=RGBColor(233, 213, 255), # Light purple
        body_color=RGBColor(250, 250, 250),  # White
        accent_color=RGBColor(192, 132, 252), # Purple accent
        is_dark=True
    )
    
    # 5. Medicine (Red/White)
    create_template(
        "templates/Medicine/Medical_Clean.pptx",
        bg_color=RGBColor(255, 255, 255),    # White
        title_color=RGBColor(220, 38, 38),   # Red
        body_color=RGBColor(55, 65, 81),     # Gray
        accent_color=RGBColor(239, 68, 68),  # Light Red
        is_dark=False
    )
    
    # 6. Engineering (Orange/Dark Gray)
    create_template(
        "templates/IT/Engineering_Pro.pptx",
        bg_color=RGBColor(31, 41, 55),       # Dark Gray
        title_color=RGBColor(249, 115, 22),  # Orange
        body_color=RGBColor(229, 231, 235),  # Light Gray
        accent_color=RGBColor(251, 146, 60), # Light Orange
        is_dark=True
    )

    # 7. History (Brown/Gold)
    create_template(
        "templates/General/History_Classic.pptx",
        bg_color=RGBColor(67, 20, 7),        # Dark Brown
        title_color=RGBColor(253, 230, 138), # Gold
        body_color=RGBColor(254, 243, 199),  # Light Yellow
        accent_color=RGBColor(245, 158, 11), # Amber
        is_dark=True
    )

    # 8. Law (Navy/Gold)
    create_template(
        "templates/General/Law_Justice.pptx",
        bg_color=RGBColor(30, 58, 138),      # Deep Navy
        title_color=RGBColor(253, 224, 71),  # Yellow Gold
        body_color=RGBColor(241, 245, 249),  # Off-white
        accent_color=RGBColor(234, 179, 8),  # Gold
        is_dark=True
    )

    # 9. Education (Blue/Yellow)
    create_template(
        "templates/General/Education_Bright.pptx",
        bg_color=RGBColor(240, 249, 255),    # Light Sky Blue
        title_color=RGBColor(3, 105, 161),   # Dark Sky Blue
        body_color=RGBColor(15, 23, 42),     # Navy Text
        accent_color=RGBColor(250, 204, 21), # Bright Yellow
        is_dark=False
    )
    
    print("Done generating 9 simple templates!")
