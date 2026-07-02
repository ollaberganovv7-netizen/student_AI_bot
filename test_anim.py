from pptx import Presentation
from pptx.oxml import parse_xml
from pptx.util import Inches

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
shape = slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(2), Inches(2))
# Try to add simple entrance animation
spid = shape.shape_id
anim_xml = f'''
<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>
                <p:par>
                  <p:cTn id="3" fill="hold">
                    <p:stCondLst>
                      <p:cond delay="indefinite"/>
                    </p:stCondLst>
                    <p:childTnLst>
                      <p:par>
                        <p:cTn id="4" fill="hold">
                          <p:stCondLst>
                            <p:cond delay="0"/>
                          </p:stCondLst>
                          <p:childTnLst>
                            <p:par>
                              <p:cTn id="5" presetID="1" presetClass="entr" presetSubtype="0" fill="hold" nodeType="clickEffect">
                                <p:stCondLst>
                                  <p:cond delay="0"/>
                                </p:stCondLst>
                                <p:childTnLst>
                                  <p:set>
                                    <p:cBhvr>
                                      <p:cTn id="6" dur="1" fill="hold">
                                        <p:stCondLst>
                                          <p:cond delay="0"/>
                                        </p:stCondLst>
                                      </p:cTn>
                                      <p:tgtEl>
                                        <p:spTgt spid="{spid}"/>
                                      </p:tgtEl>
                                      <p:attrNameLst>
                                        <p:attrName>style.visibility</p:attrName>
                                      </p:attrNameLst>
                                    </p:cBhvr>
                                    <p:to>
                                      <p:strVal val="visible"/>
                                    </p:to>
                                  </p:set>
                                </p:childTnLst>
                              </p:cTn>
                            </p:par>
                          </p:childTnLst>
                        </p:cTn>
                      </p:par>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
              </p:childTnLst>
            </p:cTn>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>
'''
try:
    slide._element.append(parse_xml(anim_xml))
    prs.save('test_element_anim.pptx')
    print('Success')
except Exception as e:
    print('Error:', e)
