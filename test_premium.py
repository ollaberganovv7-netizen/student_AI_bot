import os, sys, io
sys.path.append(r'C:\Users\Шухрат\Desktop\Новая папка\student_AI_bot')

# Test 1: Import check
print('Test 1: Importing premium_design...')
from services.premium_design import apply_premium_design, apply_premium_transitions
print('  OK: Import successful')

# Test 2: Palette selection
print('Test 2: Palette selection...')
from services.premium_design import _select_palette
p1 = _select_palette('Kompyuter texnologiyalari')
print(f'  Tech topic -> {p1["name"]}')
p2 = _select_palette("Tabiatni muhofaza qilish")
print(f'  Nature topic -> {p2["name"]}')
p3 = _select_palette("Iqtisodiyot nazariyasi")
print(f'  Economics topic -> {p3["name"]}')
p4 = _select_palette("Tibbiyot asoslari")
print(f'  Medicine topic -> {p4["name"]}')

# Test 3: Full generation with template
print('Test 3: Full PPTX generation...')
from services.pptx_service import generate_pptx

slides = [
    {'title': 'IQTIBOS', 'content': ['Test quote', '- Author']},
    {'title': 'Reja', 'content': ['1. Kirish', '2. Asosiy qism', '3. Xulosa']},
    {'title': 'Kirish', 'content': ['Bu test slayd uchun kirish matni.']},
    {'title': 'I O\'QUV SAVOLI\nAsosiy tushunchalar', 'content': ['Birinchi savol matni']},
    {'title': 'Tahlil', 'content': ['Tahlil matni bu yerda']},
    {'title': 'Xulosa', 'content': ['Xulosa matni']},
]

tpl = os.path.join(r'C:\Users\Шухрат\Desktop\Новая папка\student_AI_bot', 'templates', 'modern', 'beige-minimalism.pptx')
pptx_bytes = generate_pptx(
    slides_data=slides,
    template_path=tpl,
    topic='Kompyuter texnologiyalari',
    university='TDIU',
    author='Test User',
    reviewer='Prof. Test',
)

from pptx import Presentation
prs = Presentation(io.BytesIO(pptx_bytes))
print(f'  Total slides: {len(prs.slides)}')

for i, s in enumerate(prs.slides):
    has_bg = s._element.find('{http://schemas.openxmlformats.org/presentationml/2006/main}bg') is not None
    shape_count = len(s.shapes)
    texts = [sh.text[:30] for sh in s.shapes if sh.has_text_frame and sh.text.strip()]
    print(f'  Slide {i}: bg={has_bg}, shapes={shape_count}, texts={texts[:3]}...')

# Save for visual inspection
out_path = os.path.join(r'C:\Users\Шухрат\Desktop\Новая папка\student_AI_bot', 'test_premium.pptx')
with open(out_path, 'wb') as f:
    f.write(pptx_bytes)
print(f'  Saved to: {out_path}')

print('\nALL TESTS PASSED!')
