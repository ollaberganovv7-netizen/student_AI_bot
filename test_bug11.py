import os, sys, io
sys.path.append(r'C:\Users\Шухрат\Desktop\Новая папка\student_AI_bot')
from pptx import Presentation
from services.pptx_service import generate_pptx

slides = [{'title': 'IQTIBOS', 'content': ['Quote', 'Author']}, {'title': 'Kirish', 'content': 'Test2'}]
pptx_bytes = generate_pptx(
    slides_data=slides,
    template_path=os.path.join(r'C:\Users\Шухрат\Desktop\Новая папка\student_AI_bot', 'templates', 'modern', 'beige-minimalism.pptx'),
    topic='Test Topic',
    university='',
    author='Me',
)
prs = Presentation(io.BytesIO(pptx_bytes))
print('Total slides:', len(prs.slides))
print('Slide 0:', [shape.text.replace('\n', ' ') for shape in prs.slides[0].shapes if shape.has_text_frame])
