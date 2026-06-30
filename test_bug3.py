import os, sys, io
sys.path.append(r'C:\Users\Шухрат\Desktop\Новая папка\student_AI_bot')
from pptx import Presentation
from services.pptx_service import _remove_slide

path = os.path.join(r'C:\Users\Шухрат\Desktop\Новая папка\student_AI_bot', 'templates', 'classic.pptx')
prs = Presentation(path)
print('Before:', len(prs.slides))
_remove_slide(prs, 1)
print('After 1 remove:', len(prs.slides))
_remove_slide(prs, 1)
print('After 2 removes:', len(prs.slides))
for i, s in enumerate(prs.slides):
    texts = []
    for shape in s.shapes:
        if shape.has_text_frame:
            texts.append(shape.text.replace('\n', ' '))
    print(f'Slide {i}: {texts}')
