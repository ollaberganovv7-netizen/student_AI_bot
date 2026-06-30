import os, sys, io
sys.path.append(r'C:\Users\Шухрат\Desktop\Новая папка\student_AI_bot')
from pptx import Presentation
from services.pptx_service import generate_pptx

slides = [{'title': 'IQTIBOS', 'content': ['Quote', 'Author']}, {'title': 'Kirish', 'content': 'Test2'}]
pptx_bytes = generate_pptx(
    slides_data=slides,
    template_path=os.path.join(r'C:\Users\Шухрат\Desktop\Новая папка\student_AI_bot', 'templates', 'modern', 'beige-minimalism.pptx'),
    topic='Test Topic',
    university='Test Uni',
    author='Me',
)
prs = Presentation(io.BytesIO(pptx_bytes))
print('Total slides:', len(prs.slides))
for i, s in enumerate(prs.slides):
    texts = []
    for shape in s.shapes:
        if shape.has_text_frame:
            texts.append(shape.text.replace('\n', ' '))
    print(f'Slide {i}: {texts}')
